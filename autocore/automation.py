# automation.py header
# AutoCore - Automate Core Actions
# Author: Ash
# GitHub: https://github.com/AshAutomates/AutoCore
# Supports: Windows, Linux
# Install: pip install autocore
# Usage: from autocore import *

# Standard library imports
import atexit
import configparser
import csv
import datetime
import email
import io
import json
import logging
import os
import platform
import re
import sqlite3
import subprocess
import sys
import time
import tkinter as tk
import traceback
import xml.etree.ElementTree as et
from difflib import get_close_matches
from logging.handlers import RotatingFileHandler
from pathlib import Path
from typing import cast, Tuple, Optional

# Third-party imports - Image processing
import cv2  # pip install opencv-python
import numpy as np
from PIL import Image, ImageTk

# Third-party imports - Automation
import pyautogui
import pyperclip
import keyboard

# Third-party imports - OCR
import easyocr  # pip install "numpy<2" easyocr

# Third-party imports - Web scraping & parsing
from bs4 import BeautifulSoup

# Third-party imports - Selenium
from selenium import webdriver
from selenium.common.exceptions import *
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service as ChromeService
from selenium.webdriver.common.action_chains import ActionChains
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import Select
from webdriver_manager.chrome import ChromeDriverManager

# Third-party imports - Windows automation
if platform.system() == "Windows":
    import win32con
    import win32gui

# Third-party imports - Document processing
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import ebooklib
from ebooklib import epub
from odf import teletype
from odf.opendocument import load as odf_load
from odf.text import P
import extract_msg
from striprtf.striprtf import rtf_to_text

# Third-party imports - Utilities
import pyttsx3
import yaml

# to avoid 'RuntimeError: maximum recursion depth exceeded'
sys.setrecursionlimit(1500)

# Enable PyAutoGUI fail-safe: move mouse to any screen corner to abort script
pyautogui.FAILSAFE = True

# Global variables to track logging state
_log_file_handler: Optional[RotatingFileHandler] = None
_original_stdout = None
_original_stderr = None
_log_folder: Path = Path("logs")
_script_had_error = False  # Track if unhandled exception occurred


def _preprocess_for_ocr(image):
    """
    Preprocess image for better OCR accuracy.
    Accepts numpy array, returns cleaned binary image ready for EasyOCR.
    """
    gray_image = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
    gray_image = cv2.resize(gray_image, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
    denoised = cv2.fastNlMeansDenoising(gray_image, h=10)
    threshold_img = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 1))
    cleaned = cv2.morphologyEx(threshold_img, cv2.MORPH_CLOSE, kernel)
    return cleaned


def _click_word_by_ocr(word_to_search, occurrences_to_click, button='left'):
    """
    Click on text found via OCR on screen.

    Args:
        word_to_search: Text to search for on screen
        occurrences_to_click: 0 = click all occurrences, N = click Nth occurrence only
        button: 'left' for left-click (default), 'right' for right-click

    Returns:
        True if any click happened, False otherwise

    Note:
        This function is used internally by click() and click_right() functions.
        Cross-platform compatible (Windows, Linux).

        For this function to work on Linux, ensure scrot is installed:
        # Ubuntu/Debian
        sudo apt-get install scrot

        # RHEL/CentOS/Fedora
        sudo yum install scrot
    """
    click_count = 0

    try:
        # Get shared OCR reader
        reader = _get_ocr_reader()

        # Capture screen using pyautogui (cross-platform)
        screenshot_img = pyautogui.screenshot()

        # Convert PIL image to numpy array
        image = np.array(screenshot_img)

        # Convert RGB to BGR for OpenCV
        image = cv2.cvtColor(image, cv2.COLOR_RGB2BGR)
        cleaned = _preprocess_for_ocr(image)

        # Extract text with bounding boxes using EasyOCR
        results = reader.readtext(cleaned)

        # Filter matching text and collect positions
        matches = []
        for (bbox, text, confidence) in results:
            if word_to_search.lower() in text.lower():
                # Calculate center point (adjust for 2x upscaling)
                center_x = int((bbox[0][0] + bbox[2][0]) / 2 / 2)
                center_y = int((bbox[0][1] + bbox[2][1]) / 2 / 2)
                matches.append((center_x, center_y, text))

        # Sort matches by position (top-to-bottom, then left-to-right)
        matches.sort(key=lambda m: (m[1], m[0]))

        total_matches = len(matches)

        # Handle case when no matches found
        if total_matches == 0:
            print(f"'{word_to_search}' not found on screen")
            return False

        # Determine click action based on button parameter
        click_action = pyautogui.rightClick if button == 'right' else pyautogui.click
        action_name = "right-clicking" if button == 'right' else "clicking"

        # Click on matches based on occurrences_to_click
        if occurrences_to_click == 0:
            # Click all occurrences
            if total_matches == 1:
                print(f"1 occurrence of '{word_to_search}' found, {action_name} it")
            else:
                print(f"{total_matches} occurrences of '{word_to_search}' found, {action_name} all")
            for (x, y, text) in matches:
                click_action(x, y)
                click_count += 1
        elif 0 < occurrences_to_click <= total_matches:
            # Click specific occurrence (1-indexed)
            if total_matches == 1:
                print(f"1 occurrence of '{word_to_search}' found, {action_name} it")
            else:
                print(
                    f"{total_matches} occurrences of '{word_to_search}' found, {action_name} occurrence #{occurrences_to_click}")
            x, y, text = matches[occurrences_to_click - 1]
            click_action(x, y)
            click_count = 1
        else:
            # Requested occurrence doesn't exist
            if total_matches == 1:
                print(
                    f"1 occurrence of '{word_to_search}' found on screen, but you tried to {action_name[:-3]} occurrence #{occurrences_to_click} which doesn't exist")
            else:
                print(
                    f"{total_matches} occurrences of '{word_to_search}' found on screen, but you tried to {action_name[:-3]} occurrence #{occurrences_to_click} which doesn't exist")
            return False

        # Return True if any click happened
        return click_count > 0

    except Exception as e:
        print(f"Error in OCR click: {e}")
        return False


def _get_ocr_reader():
    """
    Lazy initialize OCR reader - only creates when first needed.
    Tries GPU first, falls back to CPU if GPU unavailable.
    Shared across all OCR functions.
    On first run, downloads OCR models (~100MB) and prints status messages.
    Subsequent runs load instantly from cache.
    """
    if not hasattr(_get_ocr_reader, 'reader'):
        try:
            print("Initializing OCR engine (first run downloads models, please wait)...")
            # Try GPU first
            _get_ocr_reader.reader = easyocr.Reader(['en'], gpu=True, verbose=False)
        except:
            # Fall back to CPU if GPU fails
            _get_ocr_reader.reader = easyocr.Reader(['en'], gpu=False, verbose=False)
        print("OCR engine ready.")
    return _get_ocr_reader.reader


def _get_web_element(driver_obj, selector_type, selector):
    """
    Internal helper function to locate a web element using Selenium.

    Args:
        driver_obj: Selenium WebDriver instance
        selector_type: Type of selector ('id', 'xpath', 'class', 'name', 'css', 'tag', 'text', 'partial')
        selector: Selector value/string

    Returns:
        WebElement if found, None if not found

    Examples:
        element = get_web_element(driver, 'id', 'submit-button')
        element = get_web_element(driver, 'xpath', '//button[@type="submit"]')
        element = get_web_element(driver, 'class', 'btn-primary')

    Note:
        This is an internal helper function used by other functions like click(), write(), etc.
        Users typically don't call this directly.
    """
    try:
        if selector_type == "id":
            return driver_obj.find_element(By.ID, selector)
        elif selector_type == "xpath":
            return driver_obj.find_element(By.XPATH, selector)
        elif selector_type == "class":
            return driver_obj.find_element(By.CLASS_NAME, selector)
        elif selector_type == "name":
            return driver_obj.find_element(By.NAME, selector)
        elif selector_type == "css":
            return driver_obj.find_element(By.CSS_SELECTOR, selector)
        elif selector_type == "tag":
            return driver_obj.find_element(By.TAG_NAME, selector)
        elif selector_type == "text":  # remember, text search is case-sensitive
            return driver_obj.find_element(By.LINK_TEXT, selector)
        elif selector_type == "partial":  # remember, partial text search is case-sensitive
            return driver_obj.find_element(By.PARTIAL_LINK_TEXT, selector)
        else:
            return None
    except NoSuchElementException:
        return None


class _CustomRotatingFileNameHandler(RotatingFileHandler):
    """Custom rotating handler that keeps .txt extension with part numbers and cleans up old logs"""

    def rotation_filename(self, default_name):
        # Convert log.txt.1 → log_part_1.txt format
        match = re.search(r'\.txt\.(\d+)$', default_name)
        if match:
            part_num = match.group(1)
            return default_name.replace(f'.txt.{part_num}', f'_part_{part_num}.txt')
        return default_name

    def doRollover(self):
        """Override rollover to cleanup old logs after each rotation"""
        # Do the normal rotation first
        super().doRollover()

        # Now cleanup old logs if folder exceeds 100MB
        self._cleanup_old_logs()

    def _cleanup_old_logs(self):
        """Delete oldest log files if total folder size exceeds 100MB"""
        try:
            log_folder = Path("logs")
            if not log_folder.exists():
                return

            log_files = []
            total_size = 0

            for log_file in log_folder.glob("log_*.txt*"):
                size = log_file.stat().st_size
                mtime = log_file.stat().st_mtime
                log_files.append((log_file, size, mtime))
                total_size += size

            # Convert to MB
            total_size_mb = total_size / (1024 * 1024)
            max_total_mb = 100

            # Keep deleting oldest files until under 100MB
            if total_size_mb > max_total_mb:
                # Sort by modification time (oldest first)
                log_files.sort(key=lambda x: x[2])

                for log_file, size, _ in log_files:
                    if total_size_mb <= max_total_mb:
                        break

                    try:
                        log_file.unlink() # deletes the file from the filesystem.
                        total_size_mb -= size / (1024 * 1024)
                        print(f"Deleted old log file: {log_file.name} (folder size: {total_size_mb:.1f}MB)")
                    except Exception as e:
                        print(f"Warning: Could not delete {log_file.name}: {e}")

        except Exception as e:
            print(f"Warning: Log cleanup failed: {e}")


class _LogCapture:
    """Captures print statements and writes to both terminal and log file"""

    def __init__(self, original_stream, logger, level):
        self.original_stream = original_stream
        self.logger = logger
        self.level = level

    def write(self, message):
        # Write to original stream (terminal)
        self.original_stream.write(message)
        self.original_stream.flush()

        # Write to log file (strip to avoid double newlines)
        if message.strip():
            self.logger.log(self.level, message.rstrip())

    def flush(self):
        self.original_stream.flush()


def browser(url, headless=False, implicit_wait=30, cookie_path=None):
    """
    Initialize and return a browser instance for web automation.

    Args:
        url: Target URL to navigate to
        headless: Run browser in headless mode (default: False)
        implicit_wait: Maximum seconds to wait for elements to appear (default: 30)
        cookie_path: Path to cookies JSON file (optional)
                    - Cookies MUST be in JSON format
                    - Export from Chrome using "Cookie-Editor" extension
                    - Cookie domain must match the target URL

    Returns:
        WebDriver: Browser instance, or None if initialization fails

    Examples:
        # Basic usage
        driver = browser('https://google.com')
        click(driver, 'id', 'search-button')

        # Slow-loading site
        driver = browser('https://slow-site.gov', implicit_wait=90)
        click(driver, 'id', 'submit-btn')  # Waits up to 90s

        # Fast site testing
        driver = browser('https://fast-site.com', implicit_wait=5)
        click(driver, 'id', 'login-btn')  # Fails fast in 5s

        # With cookies
        driver = browser('https://site.com', cookie_path='cookies.json')

        # Headless mode
        driver = browser('https://google.com', headless=True)

    Note:
    Requires Google Chrome to be installed.

    Windows:
        winget install Google.Chrome

    Linux (Ubuntu/Debian/Mint):
        wget https://dl.google.com/linux/direct/google-chrome-stable_current_amd64.deb
        sudo dpkg -i google-chrome-stable_current_amd64.deb
        sudo apt-get install -f -y

    Linux (RHEL/CentOS/Fedora):
        wget https://dl.google.com/linux/direct/google-chrome-stable_current_x86_64.rpm
        sudo rpm -i google-chrome-stable_current_x86_64.rpm
    """
    # Auto-add https:// if protocol is missing
    if not url.startswith(('http://', 'https://')):
        url = 'https://' + url
        print(f"Protocol not specified. Using: {url}")

    # Initialize options for Chrome
    options = Options()

    # confirming headless mode
    if headless:
        options.add_argument("--headless=new")
        options.add_argument("--window-size=1920,1080")

    # Set a user-agent string to make the automated browser look like a regular Chrome browser
    options.add_argument(
        "user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/143.0.7499.193 Safari/537.36")

    # Additional options to enhance realism and disable Selenium detection
    options.add_argument('start-maximized')  # Start browser maximized
    options.add_argument('--disable-blink-features=AutomationControlled')  # Disable automation flags
    options.add_experimental_option("excludeSwitches", ["enable-automation"])  # Exclude automation switches
    options.add_experimental_option('useAutomationExtension', False)  # Disable automation extension

    # Set preferences to avoid unnecessary pop-ups and block notifications
    prefs = {
        "profile.default_content_setting_values.notifications": 2,
        'credentials_enable_service': False,
        'profile': {'password_manager_enabled': False}
    }
    options.add_experimental_option("prefs", prefs)

    # Attempt to initialize the Chrome driver
    try:
        driver_instance = webdriver.Chrome(options=options)
    except Exception as e:
        print(f"Error initializing Chrome Driver: {e}")
        return None

    # Set an implicit wait for elements to be found
    driver_instance.implicitly_wait(implicit_wait)

    # Fallback URLs to ensure the driver is initialized properly
    fallback_urls = [
        "https://www.python.org/static/img/python-logo.png",  # Python logo (fallback 1)
        "https://upload.wikimedia.org/wikipedia/commons/f/f8/Python_logo_and_wordmark.svg",
        # Python logo SVG (fallback 2)
        "https://www.google.com"  # Google (fallback 3)
    ]

    # Attempt to load each fallback URL until successful
    for fallback_url in fallback_urls:
        try:
            driver_instance.get(fallback_url)
            break
        except Exception as e:
            print(f"Error loading URL {fallback_url}: {e}")

    # Load cookies from the specified path, if available.
    if cookie_path and os.path.exists(cookie_path):
        # Load target page to inject cookie
        driver_instance.get(url)

        try:
            with open(cookie_path, "r", encoding="utf-8") as f:
                cookies = json.load(f)

            for cookie in cookies:
                # Remove fields added by browser extensions that Selenium does not support
                cookie.pop("sameSite", None)
                cookie.pop("hostOnly", None)
                cookie.pop("session", None)
                cookie.pop("storeId", None)

                # Selenium rejects cookies with a leading dot in the domain, while browsers accept them, so the dot must be removed for compatibility.
                if cookie.get("domain", "").startswith("."):
                    cookie["domain"] = cookie["domain"][1:]

                driver_instance.add_cookie(cookie)

            # Reload page so cookies show effect
            driver_instance.get(url)

        except Exception as e:
            print(f"Cookie loading failed: {e}")
            return None

    else:
        # No cookies provided so single page load is enough
        driver_instance.get(url)

    # Return the driver instance - user can name it anything!
    return driver_instance


def click(*where):
    """
    Performs left-click based on different input types.

    Modes:
        1. Image matching: Click on visual element
        2. OCR text matching: Click on text found on screen (with nth occurrence support)
        3. Coordinates: Click at specific x, y position
        4. Color matching: Click on specific color in region
        5. Selenium web element: Click on element in browser

    Args:
        *where: Variable arguments depending on click mode

    Examples:
        # Image matching
        click('button.png')

        # OCR text matching
        click('Submit')                    # Click first occurrence
        click('Submit', 2)                 # Click 2nd occurrence
        click('Login', 0)                  # Click all occurrences

        # Coordinates
        click(100, 200)

        # Color matching in region
        click(x1, y1, x2, y2, r, g, b)              # Find and click color
        click(x1, y1, x2, y2, r, g, b, tolerance)   # With tolerance

        # Selenium (pass driver object first)
        click(driver, 'id', 'submit-button')
        click(driver, 'xpath', '//button[@id="submit"]')
        click(driver, 'class', 'btn-primary')
        click(driver, 'name', 'username')
        click(driver, 'css', 'button.submit')
        click(driver, 'tag', 'button')
        click(driver, 'text', 'Click Here')
        click(driver, 'partial', 'Click')

    Returns:
        True if successful, False otherwise
    """

    # Check if first argument is a WebDriver object
    if len(where) > 0 and hasattr(where[0], 'find_element'):
        # SELENIUM MODE - driver object passed
        driver_obj = where[0]

        if len(where) < 3:
            print("Error: Selenium click requires 3 arguments: click(driver, selector_type, selector)")
            return False

        selector_type = where[1]
        selector = where[2]

        try:
            # Use get_web_element with driver object
            element = _get_web_element(driver_obj, selector_type, selector)
            if element:
                element.click()
                return True
            else:
                print(f"Element not found: {selector_type} - {selector}")
                return False
        except NoSuchElementException:
            print(f"Element of {selector_type} - {selector} not found.")
            return False
        except Exception as e:
            print(f"Error clicking element: {e}")
            return False

    # PYAUTOGUI MODES - no driver object
    elif len(where) == 1:
        # Click on an image file if the argument contains a file extension
        if '.' in where[0]:
            # Assuming 'where[0]' is the path to an image file
            pyautogui.click(where[0])
            return True
        else:
            # Click based on OCR requests - 'where[0]' is assumed to be text for OCR
            # Click first occurrence by default
            result = _click_word_by_ocr(where[0], 1, button='left')
            return result

    elif len(where) == 2:
        # Click on coordinates if both arguments are integers
        if isinstance(where[0], int) and isinstance(where[1], int):
            pyautogui.click(where[0], where[1])
            return True
        # Check if it's OCR with occurrence number (text, occurrence)
        elif isinstance(where[0], str) and isinstance(where[1], int):
            # OCR mode with nth occurrence
            result = _click_word_by_ocr(where[0], where[1], button='left')
            return result
        else:
            print("Error: Invalid arguments for click()")
            return False

    # Click by color in a region (7 or 8 arguments)
    elif len(where) in [7, 8]:
        x_from, y_from, x_to, y_to, r, g, b = where[:7]
        tolerance = where[7] if len(where) == 8 else 0

        try:
            # Take a screenshot of the specified area
            screenshot_img = pyautogui.screenshot(region=(x_from, y_from, x_to - x_from, y_to - y_from))
            screenshot_img = np.array(screenshot_img)
            screenshot_img = cv2.cvtColor(screenshot_img, cv2.COLOR_RGB2BGR)

            # Define the lower and upper bounds of the target color
            lower = np.array([b - tolerance, g - tolerance, r - tolerance])
            upper = np.array([b + tolerance, g + tolerance, r + tolerance])

            # Find the color
            mask = cv2.inRange(screenshot_img, lower, upper)
            points = cv2.findNonZero(mask)

            if points is not None:
                # Click the first matching pixel
                click_x, click_y = points[0][0]
                pyautogui.click(x_from + click_x, y_from + click_y)
                print(f'Pixel found and clicked at ({x_from + click_x}, {y_from + click_y}).')
                return True

            print('Pixel not found.')
            return False

        except Exception as e:
            print(f"Error during color search: {e}")
            return False

    else:
        print("Error: Invalid arguments for click()")
        return False


def click_right(*where):
    """
    Performs right-click (context menu) based on different input types.

    Modes:
        1. Image matching: Right-click on visual element
        2. OCR text matching: Right-click on text found on screen (with nth occurrence support)
        3. Coordinates: Right-click at specific x, y position
        4. Color matching: Right-click on specific color in region
        5. Selenium web element: Right-click on element in browser

    Args:
        *where: Variable arguments depending on click mode

    Examples:
        # Image matching
        click_right('button.png')

        # OCR text matching
        click_right('Submit')              # Right-click first occurrence
        click_right('Submit', 2)           # Right-click 2nd occurrence
        click_right('Login', 0)            # Right-click all occurrences

        # Coordinates
        click_right(100, 200)

        # Color matching in region
        click_right(x1, y1, x2, y2, r, g, b)              # Find and right-click color
        click_right(x1, y1, x2, y2, r, g, b, tolerance)   # With tolerance

        # Selenium (pass driver object first)
        click_right(driver, 'id', 'submit-button')
        click_right(driver, 'xpath', '//button[@id="submit"]')
        click_right(driver, 'class', 'btn-primary')
        click_right(driver, 'name', 'username')
        click_right(driver, 'css', 'button.submit')
        click_right(driver, 'tag', 'button')
        click_right(driver, 'text', 'Click Here')
        click_right(driver, 'partial', 'Click')

    Returns:
        True if successful, False otherwise
    """

    # Check if first argument is a WebDriver object
    if len(where) > 0 and hasattr(where[0], 'find_element'):
        # SELENIUM MODE - driver object passed
        driver_obj = where[0]

        if len(where) < 3:
            print("Error: Selenium click_right requires 3 arguments: click_right(driver, selector_type, selector)")
            return False

        selector_type = where[1]
        selector = where[2]

        try:
            # Use get_web_element with driver object
            element = _get_web_element(driver_obj, selector_type, selector)
            if element:
                # Perform right-click using ActionChains
                ActionChains(driver_obj).context_click(element).perform()
                return True
            else:
                print(f"Element not found: {selector_type} - {selector}")
                return False
        except NoSuchElementException:
            print(f"Element of {selector_type} - {selector} not found.")
            return False
        except Exception as e:
            print(f"Error right-clicking element: {e}")
            return False

    # PYAUTOGUI MODES - no driver object
    elif len(where) == 1:
        # Right-click on an image file if the argument contains a file extension
        if '.' in where[0]:
            # Assuming 'where[0]' is the path to an image file
            pyautogui.rightClick(where[0])
            return True
        else:
            # Right-click based on OCR requests - 'where[0]' is assumed to be text for OCR
            # Right-click first occurrence by default
            result = _click_word_by_ocr(where[0], 1, button='right')
            return result

    elif len(where) == 2:
        # Right-click on coordinates if both arguments are integers
        if isinstance(where[0], int) and isinstance(where[1], int):
            pyautogui.rightClick(where[0], where[1])
            return True
        # Check if it's OCR with occurrence number (text, occurrence)
        elif isinstance(where[0], str) and isinstance(where[1], int):
            # OCR mode with nth occurrence
            result = _click_word_by_ocr(where[0], where[1], button='right')
            return result
        else:
            print("Error: Invalid arguments for click_right()")
            return False

    # Right-click by color in a region (7 or 8 arguments)
    elif len(where) in [7, 8]:
        x_from, y_from, x_to, y_to, r, g, b = where[:7]
        tolerance = where[7] if len(where) == 8 else 0

        try:
            # Take a screenshot of the specified area
            screenshot_img = pyautogui.screenshot(region=(x_from, y_from, x_to - x_from, y_to - y_from))
            screenshot_img = np.array(screenshot_img)
            screenshot_img = cv2.cvtColor(screenshot_img, cv2.COLOR_RGB2BGR)

            # Define the lower and upper bounds of the target color
            lower = np.array([b - tolerance, g - tolerance, r - tolerance])
            upper = np.array([b + tolerance, g + tolerance, r + tolerance])

            # Find the color
            mask = cv2.inRange(screenshot_img, lower, upper)
            points = cv2.findNonZero(mask)

            if points is not None:
                # Right-click the first matching pixel
                click_x, click_y = points[0][0]
                pyautogui.rightClick(x_from + click_x, y_from + click_y)
                print(f'Pixel found and right-clicked at ({x_from + click_x}, {y_from + click_y}).')
                return True

            print('Pixel not found.')
            return False

        except Exception as e:
            print(f"Error during color search: {e}")
            return False

    else:
        print("Error: Invalid arguments for click_right()")
        return False


def copy(*where):
    """
    Copies text from various sources: screen, clipboard, Selenium elements, or web pages.

    Modes:
        1. Active window: Copy all content from current window
        2. Clipboard: Get current clipboard content
        3. Screen coordinates: Click at position and copy
        4. Selenium webpage: Copy entire page content
        5. Selenium element: Copy element text or attribute value

    Args:
        *where: Variable arguments depending on copy mode

    Examples:
        # Active window - Copy everything from current window
        copy()                                      # Ctrl+A, Ctrl+C from active window

        # Clipboard
        copy('clipboard')                           # Get current clipboard content

        # Screen coordinates
        copy(500, 300)                              # Click at (500, 300) and copy

        # Selenium webpage - Copy entire page
        copy(driver)                                # Copy all webpage content

        # Selenium element - Copy text
        copy(driver, 'id', 'username-display')
        copy(driver, 'xpath', '//div[@class="name"]')
        copy(driver, 'class', 'user-info')
        copy(driver, 'name', 'description')
        copy(driver, 'css', 'p.content')
        copy(driver, 'tag', 'h1')
        copy(driver, 'text', 'Welcome')
        copy(driver, 'partial', 'Hello')

        # Selenium element - Copy attribute
        copy(driver, 'id', 'download-link', 'href')         # Get link URL
        copy(driver, 'class', 'product-img', 'src')         # Get image source
        copy(driver, 'id', 'email-field', 'value')          # Get input value
        copy(driver, 'xpath', '//a[@id="link"]', 'title')   # Get title attribute

    Returns:
        str: Copied text, or None if error/not found
    """

    # in this variable the copied content will be stored and returned at end
    result = None

    # clear clipboard before we copy anything
    pyperclip.copy('')

    # NO ARGUMENTS - Copy everything from active window
    if len(where) == 0:
        pyautogui.hotkey('ctrl', 'a')  # Select all
        time.sleep(1)
        pyautogui.hotkey('ctrl', 'c')  # Copy
        time.sleep(1)
        result = pyperclip.paste().strip()

    # ONE ARGUMENT
    elif len(where) == 1:

        # Copy from clipboard
        if where[0] == 'clipboard':
            result = pyperclip.paste().strip()

        # Check if it's a WebDriver object - copy entire webpage
        elif hasattr(where[0], 'find_element'):
            driver_obj = where[0]
            try:
                action = ActionChains(driver_obj)
                action.key_down(Keys.CONTROL).send_keys('a').key_up(Keys.CONTROL).perform()  # Select all
                time.sleep(1)
                action.key_down(Keys.CONTROL).send_keys('c').key_up(Keys.CONTROL).perform()  # Copy
                time.sleep(1)
                result = pyperclip.paste().strip()
                pyperclip.copy('')  # clear clipboard
                driver_obj.execute_script("window.getSelection().removeAllRanges();")  # Deselect webpage text
            except Exception as e:
                print(f"Error copying webpage content: {e}")

        # If string but not 'clipboard', it's invalid
        else:
            print(f"Invalid argument: {where[0]}")

    # TWO ARGUMENTS
    elif len(where) == 2:

        # Case 1: Both integers - screen coordinates
        if isinstance(where[0], int) and isinstance(where[1], int):
            pyautogui.click(where[0], where[1])  # Click at specified coordinates
            time.sleep(1)
            pyautogui.hotkey('ctrl', 'a')  # Select all
            time.sleep(1)
            pyautogui.hotkey('ctrl', 'c')  # Copy
            time.sleep(1)
            result = pyperclip.paste().strip()

        # Case 2: First is WebDriver - invalid, selector type is missing
        elif hasattr(where[0], 'find_element'):
            print("Error: Selenium copy requires 3 arguments: copy(driver, selector_type, selector)")
            print("Example: copy(driver, 'id', 'username')")

        # Invalid combination
        else:
            print("Invalid arguments for copy()")

    # THREE OR FOUR ARGUMENTS - Selenium element text or attribute
    elif len(where) in [3, 4]:

        # Check if first argument is WebDriver
        if hasattr(where[0], 'find_element'):
            driver_obj = where[0]
            selector_type = where[1]
            selector = where[2]
            attribute = where[3] if len(where) == 4 else None

            try:
                element = _get_web_element(driver_obj, selector_type, selector)

                if element:
                    if attribute:
                        attr_value = element.get_attribute(attribute)
                        result = attr_value.strip() if attr_value else ''
                    else:
                        result = element.text.strip()
                else:
                    print(f"Element not found: {selector_type} - {selector}")

            except Exception as e:
                print(f"Error copying from element: {e}")

        # Invalid - no driver object
        else:
            print("Error: For Selenium mode, first argument must be driver object")

    # Invalid number of arguments
    else:
        print("Invalid arguments provided for copy()")

    # Final check before returning
    if result == '' or result is None:
        print("No content returned from copy()")
        return ''
    else:
        # Returning the non-blank result
        return result


def csv_to_xlsx(csv_file=None, delete_csv=True):
    """
    Converts CSV file(s) to XLSX format.

    Args:
        csv_file: Path to CSV file, or None to auto-detect single CSV in current directory
        delete_csv: If True, deletes original CSV after conversion (default: True)

    Returns:
        str: Path to created XLSX file, or None if error

    Examples:
        # Auto-detect single CSV in current directory (deletes CSV by default)
        csv_to_xlsx()                              # Finds, converts, and deletes CSV

        # Specific file (deletes CSV by default)
        csv_to_xlsx('data.csv')                    # Converts and deletes data.csv

        # Keep original CSV
        csv_to_xlsx('report.csv', delete_csv=False) # Keeps report.csv
    """
    # Auto-detect CSV if not provided
    if csv_file is None:
        csv_files = list(Path('.').glob('*.csv'))

        if len(csv_files) == 0:
            print("Error: No CSV files found in current directory")
            return None
        elif len(csv_files) == 1:
            csv_file = str(csv_files[0])
            print(f"Auto-detected CSV file: {csv_file}")
        else:
            print(f"Error: Multiple CSV files found ({len(csv_files)}). Please specify which one to convert:")
            for f in csv_files:
                print(f"  - {f.name}")
            return None

    new_wb = Workbook()
    new_ws = new_wb.active

    try:
        with open(csv_file, 'r', encoding='utf-8', newline='') as f:
            reader = csv.reader(f)
            for row in reader:
                new_ws.append(row)

        xlsx_file = Path(csv_file).with_suffix(".xlsx")
        new_wb.save(xlsx_file)
        print(f"Converted: {csv_file} → {xlsx_file.name}")

        if delete_csv:
            os.remove(csv_file)
            print(f"Deleted: {csv_file}")

        return str(xlsx_file)

    except Exception as e:
        print(f"Error converting CSV to XLSX: {e}")
        return None


def date():
    """
    Get current day of month.

    Returns:
        int: Current day (1-31)

    Example:
        date()  # 24
    """
    return time.localtime().tm_mday


def day():
    """
    Get current day of week.

    Returns:
        str: Day name in lowercase (monday, tuesday, wednesday, thursday, friday, saturday, sunday)

    Examples:
        # Weekday check
        if day() == 'monday':
            print("It is Monday today.")
    """
    days = ['monday', 'tuesday', 'wednesday', 'thursday', 'friday', 'saturday', 'sunday']
    return days[time.localtime().tm_wday]


def drag(*args):
    """
    Drag from source to target.

    Args:
        PyAutoGUI: (x1, y1, x2, y2)
        Selenium: (driver, src_type, src_selector, tgt_type, tgt_selector)

    Examples:
        # Screen drag (PyAutoGUI) - 2 second duration
        drag(100, 200, 500, 600)

        # Web element drag (Selenium)
        drag(driver, 'id', 'card-1', 'class', 'done-column')
        drag(driver, 'xpath', '//li[1]', 'xpath', '//li[5]')

        # Multiple drivers
        driver1 = selenium('https://trello.com')
        driver2 = selenium('https://jira.com')
        drag(driver1, 'id', 'task-1', 'id', 'done-column')
        drag(driver2, 'class', 'issue', 'class', 'backlog')

    Returns:
        True if successful, False otherwise
    """

    # PYAUTOGUI MODE - 4 args, all integers
    if len(args) == 4 and all(isinstance(arg, int) for arg in args):
        x1, y1, x2, y2 = args
        try:
            pyautogui.moveTo(x1, y1, 1, pyautogui.easeInOutQuad)
            pyautogui.dragTo(x2, y2, duration=2, button='left')  # 2 seconds default
            print(f"Dragged from ({x1}, {y1}) to ({x2}, {y2})")
            return True
        except Exception as e:
            print(f"Error during drag: {e}")
            return False

    # SELENIUM MODE - 5 args, first is WebDriver object
    elif len(args) == 5 and hasattr(args[0], 'find_element'):
        driver_obj = args[0]
        src_type = args[1]
        src_selector = args[2]
        tgt_type = args[3]
        tgt_selector = args[4]

        # Validate selector types
        valid_selectors = ['id', 'xpath', 'class', 'name', 'css', 'tag', 'text', 'partial']
        if src_type not in valid_selectors:
            raise ValueError(f"Invalid source selector type '{src_type}'. Valid: {', '.join(valid_selectors)}")
        if tgt_type not in valid_selectors:
            raise ValueError(f"Invalid target selector type '{tgt_type}'. Valid: {', '.join(valid_selectors)}")

        try:

            # Find elements using the passed driver object
            source = _get_web_element(driver_obj, src_type, src_selector)
            target = _get_web_element(driver_obj, tgt_type, tgt_selector)

            if not source:
                print(f"Source element not found: {src_type} = '{src_selector}'")
                return False
            if not target:
                print(f"Target element not found: {tgt_type} = '{tgt_selector}'")
                return False

            # Perform drag and drop using the passed driver object
            ActionChains(driver_obj).drag_and_drop(source, target).perform()
            print(f"Dragged element from {src_type}='{src_selector}' to {tgt_type}='{tgt_selector}'")
            return True

        except Exception as e:
            print(f"Error during Selenium drag: {e}")
            return False

    else:
        raise ValueError(
            "Invalid arguments. Use drag(x1, y1, x2, y2) or drag(driver, src_type, src_selector, tgt_type, tgt_selector)")


def dropdown_select(driver_obj, selector_type, selector, selection_criteria):
    """
    Selects an item from a dropdown menu based on the provided criteria.

    Args:
        driver_obj: Selenium WebDriver instance
        selector_type (str): The type of selector ('id', 'name', 'xpath', 'class', 'css', 'tag', 'text', 'partial')
        selector (str): The value of the selector
        selection_criteria (int or str): The index (int) or visible text (str) for selection

    Examples:
        # Select by index
        dropdown_select(driver, 'id', 'country-dropdown', 0)        # Select first option
        dropdown_select(driver, 'id', 'country-dropdown', 2)        # Select third option

        # Select by visible text
        dropdown_select(driver, 'id', 'country-dropdown', 'United States')
        dropdown_select(driver, 'name', 'language', 'English')
        dropdown_select(driver, 'xpath', '//select[@name="city"]', 'New York')

        # Different selector types
        dropdown_select(driver, 'class', 'form-select', 'Option 1')
        dropdown_select(driver, 'css', 'select.dropdown', 'Value')

    Returns:
        True if successful, False otherwise
    """
    try:
        # Using the get_web_element function to locate the dropdown element
        dropdown_element = _get_web_element(driver_obj, selector_type, selector)

        if dropdown_element is None:
            print(f"Dropdown element not found: {selector_type} - {selector}")
            return False

        select = Select(dropdown_element)

        # Selecting based on the type of 'selection_criteria'
        if isinstance(selection_criteria, int):
            select.select_by_index(selection_criteria)
            print(f"Selected dropdown option by index: {selection_criteria}")
        elif isinstance(selection_criteria, str):
            select.select_by_visible_text(selection_criteria)
            print(f"Selected dropdown option by text: '{selection_criteria}'")
        else:
            print("Invalid selection criteria. Must be an integer or string.")
            return False

        return True

    except NoSuchElementException:
        print(f"Option not found in the dropdown: {selection_criteria}")
        return False
    except Exception as e:
        print(f"Error during dropdown selection: {e}")
        return False


def erase(*args):
    """
    Erase/clear text from input fields.

    Args:
        *args: Variable arguments depending on mode

    Examples:
        # PyAutoGUI mode (erase active window)
        erase()                                  # Select all and delete (Ctrl+A, Delete)

        # Selenium mode (erase specific element)
        erase(driver, 'id', 'username')          # Clear username field
        erase(driver, 'xpath', '//input[@name="email"]')  # Clear email field
        erase(driver, 'class', 'search-box')     # Clear search box

    Returns:
        True if successful, False otherwise
    """

    try:
        # ============================================================
        # PYAUTOGUI MODE (no arguments)
        # ============================================================
        if len(args) == 0:
            # Erase active window (Ctrl+A, Delete)
            pyautogui.hotkey('ctrl', 'a', 'delete')
            print("Erased content in active window")
            return True

        # ============================================================
        # SELENIUM MODE (driver object + selector)
        # ============================================================
        elif len(args) == 3 and hasattr(args[0], 'find_element'):
            driver_obj = args[0]
            selector_type = args[1]
            selector = args[2]

            # Find and clear element
            element = _get_web_element(driver_obj, selector_type, selector)
            if element:
                element.clear()
                print(f"Cleared element: {selector_type} - {selector}")
                return True
            else:
                print(f"Element not found: {selector_type} - {selector}")
                return False

        # ============================================================
        # ERROR - Invalid arguments
        # ============================================================
        else:
            print("Error: Invalid arguments for erase()")
            print("Use: erase() or erase(driver, selector_type, selector)")
            return False

    except Exception as e:
        print(f"Error erasing content: {e}")
        return False


def find_browser(*args):
    """
    Find text in browser using Ctrl+F (find function).

    Args:
        *args: Variable arguments depending on mode

    Examples:
        # PyAutoGUI mode (any window)
        find_browser('Python')              # Find in active window
        find_browser('error message')       # Find phrase

        # Selenium mode (browser)
        find_browser(driver, 'Python')      # Find in Selenium browser
        find_browser(driver, 'contact us')  # Find phrase in browser

    Returns:
        True if successful, False otherwise

    Note:
        - PyAutoGUI: Opens browser find dialog (Ctrl+F), types search term, presses Enter, then Esc
        - Selenium: Uses JavaScript to find and highlight text on page
        - Default wait time: 1 second between actions
    """

    wait_time = 1  # Default wait time (1 second)

    try:
        # ============================================================
        # SELENIUM MODE (driver object passed)
        # ============================================================
        if len(args) >= 2 and hasattr(args[0], 'execute_script'):
            driver_obj = args[0]
            search_text = args[1]

            # Use JavaScript to find and highlight text
            # Using raw string (r"...") to avoid escape sequence warnings
            script = r"""
            var searchText = arguments[0];
            var body = document.body;
            var innerHTML = body.innerHTML;

            // Remove previous highlights
            innerHTML = innerHTML.replace(/<mark style="background-color: yellow;">([^<]*)<\/mark>/gi, "$1");

            // Add new highlights
            var regex = new RegExp(searchText, "gi");
            innerHTML = innerHTML.replace(regex, '<mark style="background-color: yellow;">$&</mark>');

            body.innerHTML = innerHTML;

            // Scroll to first match
            var firstMatch = document.querySelector('mark');
            if (firstMatch) {
                firstMatch.scrollIntoView({ behavior: 'smooth', block: 'center' });
                return true;
            }
            return false;
            """

            # Pass search_text as argument to avoid string interpolation issues
            result = driver_obj.execute_script(script, search_text)

            if result:
                print(f"[Selenium] Text '{search_text}' found and highlighted")
                return True
            else:
                print(f"[Selenium] Text '{search_text}' not found on page")
                return False

        # ============================================================
        # PYAUTOGUI MODE (no driver object)
        # ============================================================
        elif len(args) == 1:
            search_text = args[0]

            # Open browser find dialog
            pyautogui.hotkey('ctrl', 'f')
            time.sleep(wait_time)

            # Type search text
            pyautogui.typewrite(search_text)
            time.sleep(wait_time)

            # Press Enter to find
            pyautogui.press('enter')
            time.sleep(wait_time)

            # Close find dialog
            pyautogui.press('esc')
            time.sleep(wait_time)

            print(f"[PyAutoGUI] Searched for '{search_text}'")
            return True

        else:
            print("Error: Invalid arguments for find_browser()")
            print("Use: find_browser(text) or find_browser(driver, text)")
            return False

    except Exception as e:
        print(f"Error during browser find: {e}")
        return False


def find_key(data, key):
    """
    Recursively finds all values of a specified key in nested data structures.

    Supports dictionaries, lists, and tuples.

    Args:
        data: Data structure to search (dict, list, or tuple)
        key: Key name to find

    Returns:
        list: All values found for the key (empty list if not found)

    Examples:
        # Single occurrence
        data = {'name': 'John', 'age': 30}
        name = find_key(data, 'name')[0]          # 'John'

        # Multiple occurrences
        data = {
            'user': {'id': 1, 'name': 'Alice'},
            'admin': {'id': 2, 'name': 'Bob'}
        }
        ids = find_key(data, 'id')                # [1, 2]

        # Nested lists/tuples
        data = {'users': [{'age': 25}, {'age': 30}]}
        ages = find_key(data, 'age')              # [25, 30]
    """
    results = []

    def extract(obj):
        """Recursively search for key in nested structure"""
        if isinstance(obj, dict):
            for k, v in obj.items():
                if k == key:
                    results.append(v)
                if isinstance(v, (dict, list, tuple)):
                    extract(v)
        elif isinstance(obj, (list, tuple)):
            for item in obj:
                extract(item)

    extract(data)
    return results


def find_str(string, starts_after, ends_before, index=0):
    """
    Extracts substring between two markers.

    Args:
        string: Text to search in
        starts_after: Start extraction after this
        ends_before: End extraction before this
        index: Which match (0=first, -1=last, 1=second, etc.)

    Returns:
        str: Extracted string, or None if not found

    Examples:
        # Get last error message
        error = find_str(logs, 'ERROR: ', '\n', -1)
        if error:
            print(f"Error: {error}")
    """
    try:
        # Escape special regex characters
        starts_after = re.escape(starts_after)
        ends_before = re.escape(ends_before)

        # Build regex pattern
        pattern = f'{starts_after}(.*?){ends_before}'
        matches = re.findall(pattern, string)

        if not matches:
            return None

        try:
            return matches[index].strip()
        except IndexError:
            return None

    except re.error as e:
        print(f"Regex error in find_str: {e}")
        return None
    except Exception as e:
        print(f"Error in find_str: {e}")
        return None


def hour():
    """
    Get current hour.

    Returns:
        int: Current hour (0-23) in 24-hour format

    Example:
        hour()  # 14 (2 PM)
    """
    return time.localtime().tm_hour


def inspect():
    """
    Opens GUI to inspect pixel position and color with zoomed preview.
    Move mouse, press 'ESC' to capture. Copies 'x, y, r, g, b' to clipboard.

    Platform: Windows only
    Note: This function uses global keyboard detection which requires
          administrator privileges on Linux/macOS, so it's Windows-only.
    """

    # Check if running on Windows
    if platform.system() != "Windows":
        print(f"inspect() is only supported on Windows.")
        print(f"Current OS: {platform.system()}")
        return

    # Flag to keep track of the window state
    window_open = True

    # noinspection PyTypeChecker
    def update_color_and_position():
        nonlocal window_open
        if not window_open:
            return

        # Check if 'ESC' key is pressed to capture
        if keyboard.is_pressed('esc'):
            try:
                x, y = pyautogui.position()
                color = pyautogui.pixel(x, y)
                hex_color = f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}".upper()

                # Format output and copy to clipboard
                formatted_output = f"Pixel at ({x}, {y}): RGB {color} | HEX {hex_color}"
                clipboard_text = f"{x}, {y}, {color[0]}, {color[1]}, {color[2]}"
                print(formatted_output)
                print(f"Copied to clipboard: {clipboard_text}")
                pyperclip.copy(clipboard_text)

                # Close window
                window_open = False
                root.destroy()
                return
            except Exception as e:
                print(f"Error capturing pixel: {e}")
                # Don't close window if capture failed
                root.after(100, update_color_and_position)
                return

        # Update display with current mouse position and color
        try:
            x, y = pyautogui.position()

            # Get screen size to validate position
            screen_width, screen_height = pyautogui.size()

            # Check if position is valid
            if x < 0 or y < 0 or x >= screen_width or y >= screen_height:
                root.after(100, update_color_and_position)
                return

            color = pyautogui.pixel(x, y)
            hex_color = f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}".upper()

            # Update labels
            position_label.config(text=f"Position: ({x}, {y})")
            color_label.config(text=f"RGB: {color}")
            hex_label.config(text=f"HEX: {hex_color}")
            color_display.config(bg=hex_color.lower())

            # Update zoomed preview
            # Capture area around cursor (30x30 pixels for 8x zoom)
            zoom_size = 30
            left = max(0, x - zoom_size // 2)
            top = max(0, y - zoom_size // 2)

            # Ensure we don't go beyond screen boundaries
            left = min(left, screen_width - zoom_size)
            top = min(top, screen_height - zoom_size)

            # Take screenshot of area around cursor
            screenshot = pyautogui.screenshot(region=(left, top, zoom_size, zoom_size))

            # Resize to make it bigger (8x magnification)
            zoomed = screenshot.resize((240, 240), Image.Resampling.NEAREST)

            # Convert to PhotoImage for tkinter
            photo = ImageTk.PhotoImage(zoomed)

            # Update canvas
            zoom_canvas.delete("all")
            zoom_canvas.create_image(120, 120, image=photo)
            zoom_canvas.image = photo

            # Draw crosshair at center
            zoom_canvas.create_line(120, 110, 120, 130, fill='red', width=2)
            zoom_canvas.create_line(110, 120, 130, 120, fill='red', width=2)

        except Exception as e:
            print(f"Pixel preview update failed: {e}")
            return

        # Schedule next update
        root.after(100, update_color_and_position)

    # Create main window
    root = tk.Tk()
    root.title("Pixel Inspector")
    root.geometry("400x560")
    root.configure(bg='#4D4D4D')

    # Position label
    position_label = tk.Label(root, text="Position:", font=("Helvetica", 12), bg='#4D4D4D', fg='white')
    position_label.pack()

    # Color label (RGB)
    color_label = tk.Label(root, text="RGB:", font=("Helvetica", 12), bg='#4D4D4D', fg='white')
    color_label.pack()

    # HEX label
    hex_label = tk.Label(root, text="HEX:", font=("Helvetica", 12), bg='#4D4D4D', fg='white')
    hex_label.pack()

    # Color display frame (wider)
    color_display = tk.Frame(root, height=60, width=200)
    color_display.pack(pady=10)

    # Zoom preview label
    zoom_label = tk.Label(root, text="Zoomed Preview (8x):", font=("Helvetica", 10), bg='#4D4D4D', fg='white')
    zoom_label.pack()

    # Canvas for zoomed preview (bigger)
    zoom_canvas = tk.Canvas(root, width=240, height=240, bg='#2D2D2D', highlightthickness=2,
                            highlightbackground='white')
    zoom_canvas.pack(pady=5)

    # Instruction label
    instruction_label = tk.Label(root, text="Press 'ESC' to capture and exit", font=("Helvetica", 10),
                                 bg='#4D4D4D', fg='white')
    instruction_label.pack(pady=10)

    # Start the update loop
    update_color_and_position()

    # Run tkinter event loop
    root.mainloop()


def log_setup(title):
    """
    Sets up logging and terminal styling for the script.

    Combines terminal setup with comprehensive logging and automatic color-coded status indication.
    Creates a logs folder and saves all output with timestamps.
    Shows output in terminal while also saving to file.

    Terminal colors change automatically based on script outcome:
    - Blue background: Script is running
    - Green background: Script completed successfully
    - Red background: Script crashed (unhandled exception)

    Args:
        title: Name for both terminal title and log file

    Examples:
        log_setup("MyScript")
        print("This gets logged")
        # ... script runs ...
        # Terminal turns GREEN on success or RED on crash

    Log file format:
        logs/log_MyScript_2026-03-26_14-30-45_IST_session-1.txt          (active - newest logs)
        logs/log_MyScript_2026-03-26_14-30-45_IST_session-1_part_1.txt   (2nd newest - rotated)
        logs/log_MyScript_2026-03-26_14-30-45_IST_session-1_part_2.txt   (3rd newest)
        ...
        logs/log_MyScript_2026-03-26_14-30-45_IST_session-1_part_9.txt   (oldest backup)

        Session numbering:
        - session-1: First run of this script
        - session-2: Second run of this script
        - session-3: Third run, etc.
        - Automatically increments based on existing log files

    Features:
        - Sets terminal title and colors (blue bg, white text)
        - Automatic color changes: Blue→Green (success) or Blue→Red (crash)
        - Automatic session numbering (increments from previous runs)
        - Captures all print() statements
        - Captures all errors and exceptions
        - Adds timestamp to each entry
        - Shows output in terminal AND saves to file
        - Automatic file rotation (10MB per file, max 10 files = 100MB per session)
        - Automatic cleanup (keeps max 100MB total logs across all sessions)

    Platform Support:
         Windows (Command Prompt): Full support
         Windows (PowerShell): Full support
         Linux: Full support
         macOS: Logging works, terminal colors not supported
    """
    global _log_file_handler, _original_stdout, _original_stderr, _log_folder

    # ========== TERMINAL SETUP ==========
    system = platform.system()

    if system == "Windows":
        # Clear the screen
        os.system('cls')

        # Set the window title (works in both CMD and PowerShell)
        os.system(f'title {title}')

        # Try CMD color command first (Blue background, white text)
        result = os.system('color 1F')

        # If color command failed (PowerShell), use ANSI escape codes
        if result != 0:
            # Enable ANSI color support in PowerShell
            print("\033[97m\033[44m", end="", flush=True)

    elif system == "Linux":
        # Set the window title using ANSI escape codes
        print(f"\33]0;{title}\a", end="", flush=True)
        # Set terminal text color to bright white and background to blue
        # ANSI: 97 = bright white text, 44 = blue background
        print("\033[97m\033[44m", end="")
        # Clear the screen
        os.system('clear')

    elif system == "Darwin":
        print("Terminal colors not supported on macOS (logging will still work)")

    # ========== LOGGING SETUP ==========

    # Create logs folder if it doesn't exist
    _log_folder.mkdir(exist_ok=True)

    # ========== DETERMINE SESSION NUMBER ==========

    # Find all existing log files for this script title
    existing_logs = list(_log_folder.glob(f"log_{title}_*_session-*.txt"))

    # Extract session numbers from existing log files
    session_numbers = []
    for log_file in existing_logs:
        # Pattern: log_title_timestamp_timezone_session-N.txt
        filename = log_file.stem  # Get filename without extension

        # Extract session number using regex
        match = re.search(r'_session-(\d+)', filename)
        if match:
            session_numbers.append(int(match.group(1)))

    # Determine next session number
    if session_numbers:
        next_session = max(session_numbers) + 1
    else:
        next_session = 1

    print(f"Starting session {next_session} for '{title}'")

    # ========== GET TIMESTAMP WITH TIMEZONE ==========

    # Get current time with system's local timezone
    try:
        now = datetime.datetime.now().astimezone()  # auto-detects system timezone
        time_zone_name = now.tzname() or "LOCAL"  # e.g., IST, EST, PDT
    except Exception:
        now = datetime.datetime.now()
        time_zone_name = "LOCAL"

    # Create log filename with session number
    timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
    log_filename = f"log_{title}_{timestamp}_{time_zone_name}_session-{next_session}.txt"
    log_filepath = _log_folder / log_filename

    # ========== CONFIGURE LOGGER ==========

    # Configure logger
    logger = logging.getLogger(title)
    logger.setLevel(logging.DEBUG)

    # Remove existing handlers
    logger.handlers.clear()

    # Create file handler with rotation (10MB per file)
    _log_file_handler = _CustomRotatingFileNameHandler(
        log_filepath,
        maxBytes=10 * 1024 * 1024,  # 10MB
        backupCount=9  # 10 files total (1 active + 9 backups) = 100MB
    )

    # Create formatter with timestamp
    formatter = logging.Formatter(
        '%(asctime)s -> %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    _log_file_handler.setFormatter(formatter)
    logger.addHandler(_log_file_handler)

    # Redirect stdout and stderr
    _original_stdout = sys.stdout
    _original_stderr = sys.stderr

    sys.stdout = _LogCapture(_original_stdout, logger, logging.INFO)
    sys.stderr = _LogCapture(_original_stderr, logger, logging.ERROR)

    # Log start message
    logger.info(f"Logging started for: {title} (Session {next_session})")
    logger.info(f"Log file: {log_filepath}")

    # ========== SETUP SUCCESS/ERROR COLOR HANDLERS ==========

    global _script_had_error  # Declare at start of this section

    # Store original exception handler
    original_excepthook = sys.excepthook

    def set_success_color():
        """Change terminal to green on successful completion"""
        global _original_stdout, _original_stderr, _script_had_error, _log_file_handler

        # Don't show success if there was an error
        if _script_had_error:
            return

        # LOG SUCCESS MESSAGE directly to file (before restoring streams)
        if _log_file_handler:
            try:
                # Create a log record manually
                record = logging.LogRecord(
                    name="completion",
                    level=logging.INFO,
                    pathname="",
                    lineno=0,
                    msg="Script execution completed successfully",
                    args=(),
                    exc_info=None
                )
                _log_file_handler.emit(record)
                _log_file_handler.flush()  # Ensure it's written
            except Exception as e:
                print(f"Could not write script completion status to log file: {e}")

        # NOW restore original stdout/stderr
        if _original_stdout:
            sys.stdout = _original_stdout
            sys.stderr = _original_stderr

        # Change color (terminal only)
        if platform.system() == "Windows":
            os.system('color 2F')  # Green background, white text
        elif platform.system() == "Linux":
            print("\033[97m\033[42m", end="", flush=True)  # Green bg, white text

    def set_error_color(exc_type, exc_value, exc_traceback):
        """Change terminal to red on unhandled exception"""
        global _original_stdout, _original_stderr, _script_had_error

        # Mark that an error occurred
        _script_had_error = True

        # LOG THE EXCEPTION FIRST (while stderr is still redirected to log)
        error_msg = ''.join(traceback.format_exception(exc_type, exc_value, exc_traceback))
        print(error_msg, file=sys.stderr, end='')

        # NOW restore original stdout/stderr
        if _original_stdout:
            sys.stdout = _original_stdout
            sys.stderr = _original_stderr

        # Change color
        if platform.system() == "Windows":
            os.system('color 4F')  # Red background, white text
        elif platform.system() == "Linux":
            print("\033[97m\033[41m", end="", flush=True)  # Red bg, white text

        # DON'T call original_excepthook - we already printed the error!
        # (Calling it would print the error a second time)

    # Register success handler (runs on normal exit)
    atexit.register(set_success_color)

    # Register error handler (runs on unhandled exception)
    sys.excepthook = set_error_color


def minute():
    """
    Get current minute.

    Returns:
        int: Current minute (0-59)

    Example:
        minute()  # 23
    """
    return time.localtime().tm_min


def month():
    """
    Get current month.

    Returns:
        int: Current month (1-12)

    Example:
        month()  # 2 (February)
    """
    return time.localtime().tm_mon


def press(*keys):
    """
    Press keyboard keys with support for Selenium, PyAutoGUI, and key combinations.

    Args:
        *keys: Variable arguments for key presses

    Examples:
        # Selenium element + key (requires driver object now)
        press(driver, "xpath", "//input", "enter")
        press(driver, "id", "username", "tab")

        # Selenium driver keys (pass driver object)
        press(driver, "tab")
        press(driver, "tab", 5)           # Press tab 5 times
        press(driver, "tab", -5)          # Press shift+tab 5 times
        press(driver, "ctrl", "c")
        press(driver, "ctrl", "shift", "s")

        # PyAutoGUI keys (no driver needed)
        press("tab")                    # Single key
        press("tab", 5)                 # Press 5 times
        press("tab", -5)                # Press 5 times with shift held
        press("ctrl", "a")              # Two-key combo
        press("alt", "ctrl", "z")       # Three-key combo
        press("num5")                   # Numpad 5
        press("volumeup")               # Volume up
        press("mute")                   # Volume mute (short form)
        press("back")                   # Browser back (short form)
        press("forward")                # Browser forward (short form)
    """

    # Key mappings for special keys
    key_map = {
        'enter': Keys.RETURN,
        'tab': Keys.TAB,
        'up': Keys.ARROW_UP,
        'down': Keys.ARROW_DOWN,
        'left': Keys.ARROW_LEFT,
        'right': Keys.ARROW_RIGHT,
        'alt': Keys.ALT,
        'backspace': Keys.BACKSPACE,
        'esc': Keys.ESCAPE,
        'escape': Keys.ESCAPE,
        'home': Keys.HOME,
        'end': Keys.END,
        'insert': Keys.INSERT,
        'delete': Keys.DELETE,
        'pageup': Keys.PAGE_UP,
        'pagedown': Keys.PAGE_DOWN,
        'shift': Keys.SHIFT,
        'ctrl': Keys.CONTROL,
        'control': Keys.CONTROL,
        'space': Keys.SPACE,
        'pause': Keys.PAUSE,
        'f1': Keys.F1, 'f2': Keys.F2, 'f3': Keys.F3, 'f4': Keys.F4,
        'f5': Keys.F5, 'f6': Keys.F6, 'f7': Keys.F7, 'f8': Keys.F8,
        'f9': Keys.F9, 'f10': Keys.F10, 'f11': Keys.F11, 'f12': Keys.F12,
        'windows': Keys.COMMAND,
        'command': Keys.COMMAND
    }

    # PyAutoGUI-only keys (not supported in Selenium)
    pyautogui_only_keys = {
        'num0', 'num1', 'num2', 'num3', 'num4', 'num5', 'num6', 'num7', 'num8', 'num9',
        'numlock', 'volumeup', 'volumedown', 'volumemute', 'playpause',
        'browserback', 'browserforward', 'printscreen', 'capslock', 'scrolllock',
        # Short forms
        'back', 'forward', 'mute'
    }

    # Short form mappings for PyAutoGUI keys
    pyautogui_short_forms = {
        'back': 'browserback',
        'forward': 'browserforward',
        'mute': 'volumemute'
    }

    # Valid PyAutoGUI keys (comprehensive list) to check valid keys
    valid_pyautogui_keys = {
        '\t', '\n', '\r', ' ', '!', '"', '#', '$', '%', '&', "'", '(',
        ')', '*', '+', ',', '-', '.', '/', '0', '1', '2', '3', '4', '5', '6', '7',
        '8', '9', ':', ';', '<', '=', '>', '?', '@', '[', '\\', ']', '^', '_', '`',
        'a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n', 'o',
        'p', 'q', 'r', 's', 't', 'u', 'v', 'w', 'x', 'y', 'z', '{', '|', '}', '~',
        'accept', 'add', 'alt', 'altleft', 'altright', 'apps', 'backspace',
        'browserback', 'browserfavorites', 'browserforward', 'browserhome',
        'browserrefresh', 'browsersearch', 'browserstop', 'capslock', 'clear',
        'convert', 'ctrl', 'ctrlleft', 'ctrlright', 'decimal', 'del', 'delete',
        'divide', 'down', 'end', 'enter', 'esc', 'escape', 'execute', 'f1', 'f10',
        'f11', 'f12', 'f13', 'f14', 'f15', 'f16', 'f17', 'f18', 'f19', 'f2', 'f20',
        'f21', 'f22', 'f23', 'f24', 'f3', 'f4', 'f5', 'f6', 'f7', 'f8', 'f9',
        'final', 'fn', 'hanguel', 'hangul', 'hanja', 'help', 'home', 'insert', 'junja',
        'kana', 'kanji', 'launchapp1', 'launchapp2', 'launchmail',
        'launchmediaselect', 'left', 'modechange', 'multiply', 'nexttrack',
        'nonconvert', 'num0', 'num1', 'num2', 'num3', 'num4', 'num5', 'num6',
        'num7', 'num8', 'num9', 'numlock', 'pagedown', 'pageup', 'pause', 'pgdn',
        'pgup', 'playpause', 'prevtrack', 'print', 'printscreen', 'prntscrn',
        'prtsc', 'prtscr', 'return', 'right', 'scrolllock', 'select', 'separator',
        'shift', 'shiftleft', 'shiftright', 'sleep', 'space', 'stop', 'subtract', 'tab',
        'up', 'volumedown', 'volumemute', 'volumeup', 'win', 'winleft', 'winright', 'yen',
        'command', 'option', 'optionleft', 'optionright',
        # Short forms
        'back', 'forward', 'mute'
    }

    try:
        # ============================================================
        # SELENIUM MODE - Check if first arg is WebDriver object
        # ============================================================
        if len(keys) > 0 and hasattr(keys[0], 'find_element'):
            driver_obj = keys[0]

            # ----------------------------------------------------------
            # SELENIUM ELEMENT + KEY (driver, selector_type, selector, key)
            # ----------------------------------------------------------
            if len(keys) >= 4 and keys[1] in ('id', 'name', 'xpath', 'text', 'partial', 'tag', 'class', 'css'):
                selector_type = keys[1]
                selector = keys[2]
                key_name = keys[3].lower()

                # Check if it's a PyAutoGUI-only key
                if key_name in pyautogui_only_keys:
                    print(f"Warning: '{key_name}' key not supported in Selenium. Use PyAutoGUI instead.")
                    return False

                key = key_map.get(key_name, key_name)  # Use mapping or raw key

                # Find element and send key
                element = _get_web_element(driver_obj, selector_type, selector)
                if element:
                    element.send_keys(key)
                    return True
                else:
                    print(f"Element not found with {selector_type}: {selector}")
                    return False

            # ----------------------------------------------------------
            # SELENIUM DRIVER KEYS (driver + keys)
            # ----------------------------------------------------------
            else:
                action = ActionChains(driver_obj)

                # Single key on driver: press(driver, "tab")
                if len(keys) == 2:
                    key_name = keys[1].lower()

                    # Check if it's a PyAutoGUI-only key
                    if key_name in pyautogui_only_keys:
                        print(
                            f"Warning: '{key_name}' key not supported in Selenium. Use PyAutoGUI press('{key_name}') instead.")
                        return False

                    key = key_map.get(key_name, keys[1])  # Use mapping or raw string
                    action.send_keys(key).perform()
                    return True

                # Key + count: press(driver, "tab", 5)
                elif len(keys) == 3 and isinstance(keys[2], int):
                    count = keys[2]
                    key_name = keys[1].lower()

                    # Check if it's a PyAutoGUI-only key
                    if key_name in pyautogui_only_keys:
                        print(f"Warning: '{key_name}' key not supported in Selenium.")
                        return False

                    key = key_map.get(key_name, keys[1])

                    # Positive count: press normally
                    if count >= 0:
                        for _ in range(count):
                            action.send_keys(key).perform()
                            time.sleep(1)
                    # Negative count: press with shift held
                    else:
                        for _ in range(abs(count)):
                            action.key_down(Keys.SHIFT).send_keys(key).key_up(Keys.SHIFT).perform()
                            time.sleep(1)

                    return True

                # Multiple keys (combination): press(driver, "ctrl", "c")
                elif len(keys) >= 3:
                    # Check if any key is PyAutoGUI-only
                    for i in range(1, len(keys)):
                        if keys[i].lower() in pyautogui_only_keys:
                            print(f"Warning: '{keys[i]}' key not supported in Selenium.")
                            return False

                    # Press all keys down in order
                    for i in range(1, len(keys)):
                        key_name = keys[i].lower()
                        key = key_map.get(key_name, keys[i])

                        # If it's a modifier key, hold it down
                        if key in (Keys.CONTROL, Keys.SHIFT, Keys.ALT, Keys.COMMAND):
                            action.key_down(key)
                        else:
                            # Regular key - just send it
                            action.send_keys(key)

                    # Release all modifier keys in reverse order
                    for i in range(len(keys) - 1, 0, -1):
                        key_name = keys[i].lower()
                        key = key_map.get(key_name, keys[i])

                        if key in (Keys.CONTROL, Keys.SHIFT, Keys.ALT, Keys.COMMAND):
                            action.key_up(key)

                    action.perform()
                    return True

        # ============================================================
        # PYAUTOGUI MODE - No driver object
        # ============================================================
        else:
            # Single key press
            if len(keys) == 1:
                key_name = keys[0].lower()
                # Translate short form to full form
                actual_key = pyautogui_short_forms.get(key_name, keys[0])

                # Validate key
                if actual_key.lower() not in valid_pyautogui_keys and len(actual_key) > 1:
                    print(f"Error: Invalid key '{keys[0]}'")
                    return False

                pyautogui.press(actual_key)
                return True

            # Key + count (press n times)
            elif len(keys) == 2 and isinstance(keys[1], int):
                count = keys[1]
                key_name = keys[0].lower()
                # Translate short form to full form
                actual_key = pyautogui_short_forms.get(key_name, keys[0])

                # Validate key
                if actual_key.lower() not in valid_pyautogui_keys and len(actual_key) > 1:
                    print(f"Error: Invalid key '{keys[0]}'")
                    return False

                # Positive count: press normally
                if count >= 0:
                    for _ in range(count):
                        pyautogui.press(actual_key)
                        time.sleep(1)
                # Negative count: press with shift held
                else:
                    pyautogui.keyDown('shift')
                    for _ in range(abs(count)):
                        pyautogui.press(actual_key)
                        time.sleep(1)
                    pyautogui.keyUp('shift')
                return True

            # Multiple key combination (hotkey)
            else:
                # Translate short forms to full forms and validate
                translated_keys = []
                for key in keys:
                    key_name = key.lower() if isinstance(key, str) else key
                    translated_key = pyautogui_short_forms.get(key_name, key)

                    # Validate each key
                    if translated_key.lower() not in valid_pyautogui_keys and len(str(translated_key)) > 1:
                        print(f"Error: Invalid key '{key}' in combination")
                        return False

                    translated_keys.append(translated_key)

                pyautogui.hotkey(*translated_keys)
                return True

    except Exception as e:
        print(f"Error pressing keys {keys}: {e}")
        return False


def read(*args):
    """
    Extract text from screen (using OCR), files (by parsing file format), or a Selenium browser window.

    Modes:
        No arguments: OCR full screen
        2 integers: OCR from (x, y) to bottom-right corner
        4 integers: OCR specific region (x, y, width, height)
        1 string: Read file by parsing its format
        1 driver object: Take screenshot of Selenium browser and read using OCR

    Supported file formats:
        Documents: PDF, DOCX, PPTX, ODT, RTF
        Tabular: CSV, TSV, XLSX, SQLite
        Structured: JSON, YAML, XML, INI/CFG
        Text: TXT, LOG, MD
        Web: HTML
        Email: EML, MSG
        eBooks: EPUB
        Scripts: SH, BAT, PY

    Args:
        *args: Variable arguments depending on mode

    Returns:
        str: Extracted text (lowercase for OCR), or None if error

    Examples:
        # OCR - Read entire screen
        text = read()

        # OCR - Read from (100, 200) to bottom-right corner
        text = read(100, 200)

        # OCR - Read specific region: x=100, y=200, width=400, height=300
        text = read(100, 200, 400, 300)

        # Selenium - Read text from browser window using OCR
        d1 = browser('https://example.com')
        text = read(d1)

        # File - Read with extension
        text = read('report.pdf')
        text = read('data.csv')
        text = read('script.py')

        # File - Read without extension (auto-detects)
        text = read('report')      # Finds report.pdf automatically
        text = read('config')      # Finds config.yaml automatically

        # Check if text on screen
        if 'login' in read():
            print("Login visible!")

    Note:
        OCR first run downloads models (~100MB), subsequent runs are fast.
    """

    # ============================================================
    # DETERMINE MODE AND VALIDATE ARGUMENTS
    # ============================================================
    is_ocr = False
    region = None

    # MODE 1: OCR - No arguments (full screen)
    if len(args) == 0:
        is_ocr = True
        region = None  # Full screen

    # MODE 2: OCR - 2 integers (from x,y to bottom-right corner)
    elif len(args) == 2 and all(isinstance(arg, int) for arg in args):
        is_ocr = True
        x, y = args

        # Validate coordinates
        if x < 0 or y < 0:
            print(f"Error: Coordinates must be non-negative (x={x}, y={y})")
            return None

        # Get screen size
        screen_width, screen_height = pyautogui.size()

        # Calculate width and height to bottom-right corner
        width = screen_width - x
        height = screen_height - y

        if width <= 0 or height <= 0:
            print(f"Error: Starting point ({x}, {y}) is outside screen bounds")
            return None

        region = (x, y, width, height)

    # MODE 3: OCR - 4 integers (screen region: x, y, width, height)
    elif len(args) == 4 and all(isinstance(arg, int) for arg in args):
        is_ocr = True
        x, y, width, height = args

        # Validate region
        if width <= 0 or height <= 0:
            print(f"Error: Invalid region dimensions (width={width}, height={height})")
            return None

        if x < 0 or y < 0:
            print(f"Error: Coordinates must be non-negative (x={x}, y={y})")
            return None

        region = (x, y, width, height)

    # ============================================================
    # PERFORM OCR (if OCR mode detected)
    # ============================================================
    if is_ocr:
        try:
            # Get shared OCR reader
            reader = _get_ocr_reader()

            # Capture screenshot (region or full screen)
            if region:
                image = pyautogui.screenshot(region=region)
            else:
                image = pyautogui.screenshot()

            # Convert PIL image to numpy array
            image = np.array(image)
            cleaned = _preprocess_for_ocr(image)

            # Extract text using EasyOCR
            results = reader.readtext(cleaned, detail=0)

            # If no text was found
            if not results:
                print(f"OCR: No text found.")
                return ""  # Return blank string (not None) to avoid errors in blank text case

            # Join all detected text with spaces and convert to lowercase
            text = ' '.join(results).lower()

            return text.strip()

        except ImportError as e:
            print(f"Error: Required library not installed - {e}")
            print("Install with: pip install easyocr opencv-python")
            return None

        except Exception as e:
            print(f"Error during OCR: {e}")
            return None

    # ============================================================
    # MODE 4: Selenium driver - screenshot browser and OCR
    # ============================================================
    elif len(args) == 1 and hasattr(args[0], 'find_element'):
        driver_obj = args[0]

        try:
            reader = _get_ocr_reader()

            # Get screenshot as bytes directly without saving to disk
            png_bytes = driver_obj.get_screenshot_as_png()
            image = np.array(Image.open(io.BytesIO(png_bytes)))
            cleaned = _preprocess_for_ocr(image)

            results = reader.readtext(cleaned, detail=0)

            if not results:
                print("OCR: No text found in browser window.")
                return ""

            text = ' '.join(results).lower()
            return text.strip()

        except Exception as e:
            print(f"Could not read text from browser window: {e}")
            return None

    # ============================================================
    # MODE 5: File reading - 1 string (file path)
    # ============================================================
    elif len(args) == 1 and isinstance(args[0], str):
        file = args[0]

        # ============================================================
        # AUTO-DETECT FILE EXTENSION IF NOT PROVIDED
        # ============================================================
        if not os.path.splitext(file)[1]:  # No extension provided
            # Get directory and filename
            if os.path.dirname(file):
                directory = os.path.dirname(file)
                filename = os.path.basename(file)
            else:
                directory = '.'
                filename = file

            # Find files matching the name
            try:
                matching_files = [f for f in os.listdir(directory)
                                  if os.path.splitext(f)[0] == filename]

                if len(matching_files) == 1:
                    # Found exactly one file - use it
                    file = os.path.join(directory, matching_files[0])
                    print(f"Auto-detected file: {file}")
                elif len(matching_files) > 1:
                    print(f"Error: Multiple files found with name '{filename}':")
                    for f in matching_files:
                        print(f"  - {f}")
                    print("Please specify file extension.")
                    return None
                else:
                    print(f"Error: No file found with name '{filename}'")
                    return None

            except Exception as e:
                print(f"Error searching for file: {e}")
                return None

        # Get file extension
        _, ext = os.path.splitext(file)
        ext = ext.lower()

        try:
            # ============================================================
            # DOCUMENTS
            # ============================================================

            # PDF files
            if ext == '.pdf':
                with open(file, 'rb') as pdf_file:
                    pdf_reader = PdfReader(pdf_file)
                    return '\r\n'.join([pdf_reader.pages[i].extract_text() for i in range(len(pdf_reader.pages))])

            # Word documents
            elif ext in ['.docx', '.doc']:
                doc = Document(file)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])

            # PowerPoint presentations
            elif ext in ['.pptx', '.ppt']:
                prs = Presentation(file)
                text = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    text.append(f"=== Slide {slide_num} ===")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                    text.append('')
                return '\n'.join(text)

            # OpenDocument Text (LibreOffice)
            elif ext == '.odt':
                textdoc = odf_load(file)
                allparas = textdoc.getElementsByType(P)
                return '\n'.join([teletype.extractText(p) for p in allparas])

            # Rich Text Format
            elif ext == '.rtf':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return rtf_to_text(f.read())

            # ============================================================
            # TABULAR DATA (with >>>Row_X: format)
            # ============================================================

            # CSV files
            elif ext == '.csv':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = list(reader)

                    if not rows:
                        return ''

                    text = []
                    # First row as header
                    text.append('\t'.join(rows[0]))
                    text.append('-' * 60)

                    # Data rows with >>>Row_X: format
                    for row_num, row in enumerate(rows[1:], 1):
                        row_data = '\t'.join(row)
                        text.append(f">>>Row_{row_num}:\t{row_data}")

                    return '\n'.join(text)

            # TSV files (Tab-separated values)
            elif ext == '.tsv':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f, delimiter='\t')
                    rows = list(reader)

                    if not rows:
                        return ''

                    text = []
                    # First row as header
                    text.append('\t'.join(rows[0]))
                    text.append('-' * 60)

                    # Data rows with >>>Row_X: format
                    for row_num, row in enumerate(rows[1:], 1):
                        row_data = '\t'.join(row)
                        text.append(f">>>Row_{row_num}:\t{row_data}")

                    return '\n'.join(text)

            # Excel files
            elif ext in ['.xlsx', '.xlsm']:
                wb = load_workbook(file, data_only=True)
                text = []

                for sheet in wb.worksheets:
                    text.append(f"\n{'=' * 60}")
                    text.append(f"Sheet: {sheet.title}")
                    text.append('=' * 60)

                    rows = list(sheet.iter_rows(values_only=True))

                    if rows:
                        # First row as header
                        text.append('\t'.join([str(cell) if cell is not None else '' for cell in rows[0]]))
                        text.append('-' * 60)

                        # Data rows with >>>Row_X: format
                        for row_num, row in enumerate(rows[1:], 1):
                            row_data = '\t'.join([str(cell) if cell is not None else 'NULL' for cell in row])
                            text.append(f">>>Row_{row_num}:\t{row_data}")

                    text.append('')

                return '\n'.join(text)

            # SQLite databases
            elif ext in ['.db', '.sqlite', '.sqlite3']:
                conn = sqlite3.connect(file)
                cursor = conn.cursor()
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = cursor.fetchall()
                text = []

                for table in tables:
                    table_name = table[0]
                    text.append(f"\n{'=' * 60}")
                    text.append(f"Table: {table_name}")
                    text.append('=' * 60)

                    # Get column names
                    cursor.execute(f"PRAGMA table_info({table_name})")
                    columns = [col[1] for col in cursor.fetchall()]

                    # Add column headers
                    text.append('\t'.join(columns))
                    text.append('-' * 60)

                    # Get data (limit to 100 rows)
                    cursor.execute(f"SELECT * FROM {table_name} LIMIT 100")
                    rows = cursor.fetchall()

                    # Add rows with >>>Row_X: format
                    for row_num, row in enumerate(rows, 1):
                        row_data = '\t'.join([str(cell) if cell is not None else 'NULL' for cell in row])
                        text.append(f">>>Row_{row_num}:\t{row_data}")

                    text.append(f"\nTotal_rows: {len(rows)}\n")

                conn.close()
                return '\n'.join(text)

            # ============================================================
            # STRUCTURED DATA
            # ============================================================

            # JSON files
            elif ext == '.json':
                with open(file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data, indent=2)

            # YAML files
            elif ext in ['.yaml', '.yml']:
                with open(file, 'r', encoding='utf-8') as f:
                    data = yaml.safe_load(f)
                    return yaml.dump(data, default_flow_style=False, sort_keys=False)

            # XML files
            elif ext == '.xml':
                tree = et.parse(file)
                return et.tostring(tree.getroot(), encoding='unicode', method='xml')

            # INI/CFG files
            elif ext in ['.ini', '.cfg']:
                config = configparser.ConfigParser()
                config.read(file)
                text = []
                for section in config.sections():
                    text.append(f"[{section}]")
                    for key, value in config.items(section):
                        text.append(f"{key} = {value}")
                    text.append('')
                return '\n'.join(text)

            # ============================================================
            # TEXT FILES
            # ============================================================

            # Plain text, log files, markdown
            elif ext in ['.txt', '.log', '.md']:
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # ============================================================
            # WEB
            # ============================================================

            # HTML files
            elif ext in ['.html', '.htm']:
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    soup = BeautifulSoup(f.read(), 'html.parser')
                    return soup.get_text('\n').strip()

            # ============================================================
            # EMAIL
            # ============================================================

            # EML files (standard email)
            elif ext == '.eml':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    msg = email.message_from_file(f)
                    text = []
                    text.append(f"From: {msg.get('From')}")
                    text.append(f"To: {msg.get('To')}")
                    text.append(f"Subject: {msg.get('Subject')}")
                    text.append(f"Date: {msg.get('Date')}")
                    text.append("\n--- Body ---")

                    if msg.is_multipart():
                        for part in msg.walk():
                            if part.get_content_type() == "text/plain":
                                payload = part.get_payload(decode=True)
                                if payload:
                                    text.append(payload.decode('utf-8', errors='ignore'))
                    else:
                        payload = msg.get_payload(decode=True)
                        if payload:
                            text.append(payload.decode('utf-8', errors='ignore'))

                    return '\n'.join(text)

            # MSG files (Outlook email)
            elif ext == '.msg':
                msg = extract_msg.Message(file)
                text = []
                text.append(f">>>From: {msg.sender}")
                text.append(f">>>To: {msg.to}")
                text.append(f">>>Subject: {msg.subject}")
                text.append(f">>>Date: {msg.date}")
                text.append("\nBody ---")
                text.append(msg.body)
                msg.close()
                return '\n'.join(text)

            # ============================================================
            # EBOOKS
            # ============================================================

            # EPUB files
            elif ext == '.epub':
                book = epub.read_epub(file)
                text = []
                for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text.append(soup.get_text())
                return '\n\n'.join(text)

            # ============================================================
            # SCRIPT FILES
            # ============================================================

            # Shell scripts (Linux)
            elif ext == '.sh':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # Batch files (Windows)
            elif ext == '.bat':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # Python scripts
            elif ext == '.py':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # ============================================================
            # UNSUPPORTED FORMAT
            # ============================================================

            else:
                print(f"Unsupported file format: {ext}")
                print(f"Supported formats: PDF, DOCX, PPTX, ODT, RTF, CSV, TSV, XLSX, SQLite,")
                print(f"                   JSON, YAML, XML, INI/CFG, TXT, LOG, MD, HTML,")
                print(f"                   EML, MSG, EPUB, SH, BAT, PY")
                return None

        except Exception as e:
            print(f"Error reading file {file}: {e}")
            return None

    # ============================================================
    # INVALID ARGUMENTS
    # ============================================================
    else:
        print("Error: Invalid arguments for read()")
        print("Use: read() - OCR full screen")
        print("     read(x, y) - OCR from (x,y) to bottom-right")
        print("     read(x, y, width, height) - OCR specific region")
        print("     read(driver) - OCR browser window")
        print("     read('file.pdf') - Read file")
        return None


def run(item):
    """
    Runs a file or application on Windows and Linux.

    Platform Support:
        Windows: Uses 'start' command
        Linux: Uses 'xdg-open' for files, direct execution for apps
        macOS: Not supported

    Behavior:
        - If item is a file path: Opens with default application
        - If item is an application name: Launches the application
        - For applications, the command must be available in system PATH

    Args:
        item (str): Path to a file or name of an application to run

    Examples:
        # Open files with default application
        run("sample.txt")           # Opens in default text editor
        run("document.pdf")         # Opens in default PDF viewer
        run("C:\\Users\\file.xlsx") # Opens Excel file

        # Launch applications (Windows)
        run("calc")                 # Calculator
        run("notepad")              # Notepad
        run("mspaint")              # Paint

        # Launch applications (Linux)
        run("gedit")                # Text editor
        run("firefox")              # Browser
        run("gnome-calculator")     # Calculator

    Linux Dependencies:
        xdg-utils package (usually pre-installed)
        sudo apt-get install xdg-utils    # If needed

    Returns:
        None

    Raises:
        NotImplementedError: If called on macOS
    """

    # Check if macOS and reject
    if sys.platform.startswith('darwin') or platform.system() == "Darwin":
        raise NotImplementedError(
            "run() is not supported on macOS.\n"
            "Supported platforms: Windows, Linux\n"
            "Use subprocess.Popen() directly for macOS-specific needs."
        )

    try:
        # ============================================================
        # WINDOWS
        # ============================================================
        if sys.platform.startswith('win'):
            if os.path.isfile(item):
                # File exists - open with default application
                # Using os.startfile is more reliable than shell=True
                os.startfile(item)
            else:
                # Assume it's an application name
                # Use shell=True to allow PATH resolution
                subprocess.Popen(item, shell=True)

        # ============================================================
        # LINUX
        # ============================================================
        elif sys.platform.startswith('linux'):
            if os.path.isfile(item):
                # File exists - open with default application using xdg-open
                item_path = os.path.abspath(item)
                subprocess.Popen(['xdg-open', item_path])
            else:
                # Assume it's an application name - try to execute directly
                subprocess.Popen([item])

        else:
            # Unknown platform
            raise NotImplementedError(
                f"run() is not supported on platform: {sys.platform}\n"
                "Supported platforms: Windows (win32), Linux (linux)"
            )

    except FileNotFoundError:
        print(f"Error: Application or file not found: '{item}'")
        print("For applications, ensure the command is in your system PATH.")
    except PermissionError:
        print(f"Error: Permission denied to run: '{item}'")
        print("Check file permissions or try running with appropriate privileges.")
    except Exception as e:
        print(f"Error running '{item}': {e}")


def say(text, volume=1.0):
    """
    Speak text using offline Text-to-Speech.

    Args:
        text: Text to speak
        volume: Volume level 0.0 to 1.0 (default: 1.0)

    Examples:
        say("Download complete", volume=0.7)
        say("Error occurred")

    Note:
        Speech rate is fixed at 130 words/minute for optimal clarity.
        Automatically logs spoken text when log_setup() is active.
    """
    # Validate text input
    if not isinstance(text, str):
        raise TypeError("Text to speak must be a string")

    # Validate volume range
    if not (0.0 <= volume <= 1.0):
        raise ValueError("Volume must be between 0.0 and 1.0")

    # Initialize fresh engine (prevents silent failures in repeated usage)
    engine = pyttsx3.init()

    # Set properties
    engine.setProperty('rate', 130)  # Fixed rate for clarity across platforms
    engine.setProperty('volume', volume)

    # Speak text
    engine.say(text)
    engine.runAndWait()

    # Cleanup to prevent stuck engine
    engine.stop()

    # Log spoken text (integrates with log_setup for audit trail)
    print(f"Spoken: {text}")


def screenshot(*args):
    """
    Takes a screenshot and saves it.

    Arguments:
        *args: Variable arguments depending on usage:
            - () → Full screen, auto-named
            - (filename) → Full screen, custom filename
            - (x, y) → From (x,y) to screen edge, auto-named
            - (x, y, filename) → From (x,y) to screen edge, custom filename
            - (x, y, width, height) → Specific region, auto-named
            - (x, y, width, height, filename) → Specific region, custom filename
            - (driver, ...) → Same as above but uses Selenium driver

    Where:
        driver: Selenium WebDriver instance (for Selenium mode)
        x, y: Coordinates of top-left corner of the screenshot region
        width, height: Dimensions of the screenshot region
        filename: Name to save the screenshot (auto-generates if not provided)

    Default Naming:
        If no filename provided, auto-generates: screenshot_YYYY-MM-DD_HH-MM-SS_<unix>.png
        Example: screenshot_2025-02-18_14-30-45_1708268445.png

    Examples:
        # Full screen (PyAutoGUI)
        screenshot()                                    # Full screen, auto-named
        screenshot('desktop.png')                       # Full screen, custom name

        # Selenium full page
        screenshot(driver)                              # Selenium full page, auto-named
        screenshot(driver, 'webpage.png')               # Selenium full page, custom name

        # From top-left point to edge (PyAutoGUI)
        screenshot(100, 200)                            # From (100,200) to edge, auto-named
        screenshot(100, 200, 'crop.png')                # From (100,200) to edge, custom name

        # From top-left point to edge (Selenium)
        screenshot(driver, 100, 200)                    # Selenium from (100,200) to edge
        screenshot(driver, 100, 200, 'page.png')        # Selenium, custom name

        # Specific region (PyAutoGUI)
        screenshot(0, 0, 500, 300)                      # Region: top-left (0,0), 500x300, auto-named
        screenshot(0, 0, 500, 300, 'region.png')        # Region: top-left (0,0), 500x300, custom name

        # Specific region (Selenium)
        screenshot(driver, 0, 0, 800, 600)              # Selenium region, auto-named
        screenshot(driver, 0, 0, 800, 600, 'sel.png')   # Selenium region, custom name

    Returns:
        True if successful, False otherwise
    """

    try:
        # Determine if using Selenium or PyAutoGUI
        if len(args) > 0 and hasattr(args[0], 'save_screenshot'):
            # First argument is a WebDriver object
            use_driver = True
            driver_obj = args[0]
            remaining_args = args[1:]
        else:
            # PyAutoGUI mode
            use_driver = False
            driver_obj = None
            remaining_args = args

        # Default values
        x, y, width, height, filename = None, None, None, None, None

        # Parse remaining arguments
        if len(remaining_args) == 0:
            # screenshot() or screenshot(driver)
            # Full screen, auto-named
            x, y, width, height, filename = 0, 0, None, None, None

        elif len(remaining_args) == 1:
            # screenshot('file.png') or screenshot(driver, 'file.png')
            # Full screen, custom name
            x, y, width, height = 0, 0, None, None
            filename = remaining_args[0]

        elif len(remaining_args) == 2:
            # screenshot(100, 200) or screenshot(driver, 100, 200)
            # From (x,y) to screen edge, auto-named
            x, y = remaining_args[0], remaining_args[1]
            width, height, filename = None, None, None

        elif len(remaining_args) == 3:
            # screenshot(100, 200, 'file.png') or screenshot(driver, 100, 200, 'file.png')
            # From (x,y) to screen edge, custom name
            x, y, filename = remaining_args[0], remaining_args[1], remaining_args[2]
            width, height = None, None

        elif len(remaining_args) == 4:
            # screenshot(0, 0, 500, 300) or screenshot(driver, 0, 0, 500, 300)
            # Specific region, auto-named
            x, y, width, height = remaining_args[0], remaining_args[1], remaining_args[2], remaining_args[3]
            filename = None

        elif len(remaining_args) == 5:
            # screenshot(0, 0, 500, 300, 'file.png') or screenshot(driver, 0, 0, 500, 300, 'file.png')
            # Specific region, custom name
            x, y, width, height, filename = remaining_args

        else:
            raise ValueError("Too many arguments")

        # Convert coordinates to integers
        x = int(x) if x is not None else 0
        y = int(y) if y is not None else 0

        # Get screen dimensions
        screen_width, screen_height = pyautogui.size()

        # Calculate width and height if not provided
        if width is None:
            width = screen_width - x  # From x to right edge
        if height is None:
            height = screen_height - y  # From y to bottom edge

        # Convert to integers
        width, height = int(width), int(height)

        # Validate
        if width <= 0 or height <= 0:
            raise ValueError("Width and height must be positive")
        if x < 0 or y < 0:
            raise ValueError("Coordinates must be non-negative")

        # Generate filename if not provided
        if filename is None:
            now = datetime.datetime.now()
            timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
            unix_time = int(now.timestamp())
            filename = f"screenshot_{timestamp}_{unix_time}.png"
        else:
            # Add .png extension if missing
            if not os.path.splitext(filename)[1]:
                filename += '.png'

        # Get full path
        full_path = os.path.join(os.getcwd(), filename)

        # Take screenshot
        if use_driver:
            # Selenium: capture full page, then crop
            temp_path = 'temp_full_screenshot.png'
            driver_obj.save_screenshot(temp_path)

            image = Image.open(temp_path)
            cropped = image.crop((x, y, x + width, y + height))
            cropped.save(full_path)

            os.remove(temp_path)
        else:
            # PyAutoGUI: capture region directly
            img = pyautogui.screenshot(region=(x, y, width, height))
            img.save(full_path)

        print(f"Screenshot saved: {full_path}")
        return True

    except ValueError as e:
        print(f"Invalid input: {e}")
        return False
    except RuntimeError as e:
        print(f"Error: {e}")
        return False
    except Exception as e:
        print(f"Error taking screenshot: {e}")
        return False


def scroll(*args, duration=40):
    """
    Universal scroll function for both PyAutoGUI and Selenium.

    Args:
        *args: Variable arguments (see examples below)
        duration: Max seconds when scrolling to 'bottom'/'top' (default: 40)

    PyAutoGUI Examples (scroll any window):
        scroll()                # Scroll down 1 time (default)
        scroll(5)               # Scroll down 5 times
        scroll('down')          # Scroll down 1 time
        scroll('down', 10)      # Scroll down 10 times
        scroll('up', 5)         # Scroll up 5 times
        scroll('bottom')        # Scroll down continuously for 40 seconds
        scroll('bottom', timeout=60)  # Scroll down continuously for 60 seconds
        scroll('top')           # Scroll up continuously for 40 seconds

    Selenium Examples (pass driver object):
        scroll(driver)              # Scroll down 1 time in browser
        scroll(driver, 5)           # Scroll down 5 times in browser
        scroll(driver, 'down')      # Scroll down 1 time in browser
        scroll(driver, 'down', 10)  # Scroll down 10 times in browser
        scroll(driver, 'up', 5)     # Scroll up 5 times in browser
        scroll(driver, 'bottom')    # Scroll to bottom (auto-detect end)
        scroll(driver, 'top')       # Scroll to top (auto-detect end)
        scroll(driver, 'bottom', duration=120)  # Scroll to bottom, max 2 minutes

    Returns:
        True if successful, False otherwise
    """

    wait = 3  # Fixed wait time between scrolls (3 seconds)

    # Parse arguments
    if len(args) == 0:
        # scroll() - default: scroll down 1 time with PyAutoGUI
        use_selenium = False
        driver_obj = None
        direction = 'down'
        count = 1

    elif len(args) > 0 and hasattr(args[0], 'execute_script'):
        # Selenium mode - first arg is WebDriver object
        use_selenium = True
        driver_obj = args[0]

        if len(args) == 1:
            # scroll(driver) - scroll down 1 time
            direction = 'down'
            count = 1
        elif isinstance(args[1], int):
            # scroll(driver, 5) - scroll down 5 times
            direction = 'down'
            count = args[1]
        elif args[1] in ['down', 'up', 'bottom', 'top']:
            # scroll(driver, 'down') or scroll(driver, 'down', 10)
            direction = args[1]
            count = args[2] if len(args) > 2 else 1
        else:
            print(f"Error: Invalid argument '{args[1]}'. Use 'down', 'up', 'bottom', 'top', or a number.")
            return False

    else:
        # PyAutoGUI mode
        use_selenium = False
        driver_obj = None

        if isinstance(args[0], int):
            # scroll(5) - scroll down 5 times
            direction = 'down'
            count = args[0]
        elif args[0] in ['down', 'up', 'bottom', 'top']:
            # scroll('down') or scroll('down', 10)
            direction = args[0]
            count = args[1] if len(args) > 1 else 1
        else:
            print(f"Error: Invalid argument '{args[0]}'. Use 'down', 'up', 'bottom', 'top', or a number.")
            return False

    # Validate direction
    if direction not in ['down', 'up', 'bottom', 'top']:
        print(f"Error: Invalid direction '{direction}'.")
        return False

    try:
        # ============================================================
        # SELENIUM MODE
        # ============================================================
        if use_selenium:
            # Scroll to bottom (with auto-detection)
            if direction == 'bottom':
                print(f"[Selenium] Scrolling to bottom (timeout={duration}s)...")
                start_time = time.time()
                scrolls = 0
                last_height = driver_obj.execute_script("return document.body.scrollHeight")

                while time.time() - start_time < duration:
                    # Scroll to bottom
                    driver_obj.execute_script("window.scrollTo(0, document.body.scrollHeight);")
                    scrolls += 1

                    # Wait for content to load (3 seconds)
                    time.sleep(wait)

                    # Check new height
                    new_height = driver_obj.execute_script("return document.body.scrollHeight")

                    # If height unchanged, wait once more to confirm
                    if new_height == last_height:
                        time.sleep(wait)
                        final_height = driver_obj.execute_script("return document.body.scrollHeight")

                        if final_height == last_height:
                            print(f"Reached bottom after {scrolls} scrolls ({time.time() - start_time:.1f}s)")
                            return True

                    last_height = new_height

                    if scrolls % 10 == 0:
                        print(f"Scrolled {scrolls} times ({time.time() - start_time:.1f}s)")

                print(f"Scrolled {scrolls} times")
                return True

            # Scroll to top
            elif direction == 'top':
                print(f"[Selenium] Scrolling to top...")
                driver_obj.execute_script("window.scrollTo(0, 0);")
                time.sleep(wait)
                print("Scrolled to top")
                return True

            # Scroll down/up N times
            else:
                scroll_pixels = 500 if direction == 'down' else -500
                print(f"[Selenium] Scrolling {direction} {count} time(s)...")

                for i in range(count):
                    driver_obj.execute_script(f"window.scrollBy(0, {scroll_pixels});")
                    time.sleep(wait)

                    if count > 10 and (i + 1) % 10 == 0:
                        print(f"Scrolled {i + 1}/{count} times")

                print(f"Scrolled {direction} {count} time(s)")
                return True

        # ============================================================
        # PYAUTOGUI MODE
        # ============================================================
        else:
            # Scroll to bottom/top (time-based, no detection)
            if direction in ['bottom', 'top']:
                scroll_amount = -500 if direction == 'bottom' else 500
                print(f"[PyAutoGUI] Scrolling {direction} continuously for {duration}s...")

                start_time = time.time()
                scrolls = 0

                while time.time() - start_time < duration:
                    pyautogui.scroll(scroll_amount)
                    scrolls += 1
                    time.sleep(wait)

                    if scrolls % 10 == 0:
                        print(f"Scrolled {scrolls} times ({time.time() - start_time:.1f}s)")

                print(f"Scrolled {direction} for {duration}s ({scrolls} total scrolls)")
                return True

            # Scroll down/up N times
            else:
                scroll_amount = -500 if direction == 'down' else 500
                print(f"[PyAutoGUI] Scrolling {direction} {count} time(s)...")

                for i in range(count):
                    pyautogui.scroll(scroll_amount)
                    time.sleep(wait)

                    if count > 10 and (i + 1) % 10 == 0:
                        print(f"Scrolled {i + 1}/{count} times")

                print(f"Scrolled {direction} {count} time(s)")
                return True

    except Exception as e:
        print(f"Error while scrolling: {e}")
        return False


def second():
    """
    Get current second.

    Returns:
        int: Current second (0-59)

    Example:
        second()  # 45
    """
    return time.localtime().tm_sec


def wait(*args, countdown=True):
    """
    Wait with countdown, wait for element, or wait for color at pixel.

    Args:
        *args: Variable arguments (see examples)
        countdown: If True, shows countdown display (default: True)

    Examples:
        # Countdown wait
        wait(5)                              # Wait 5 seconds with countdown
        wait(10, countdown=False)            # Wait 10 seconds silently
        wait()                               # Wait 3 seconds (default)

        # Wait for element (Selenium) - pass driver object
        wait(driver, 'xpath', '//button')            # Wait max 180s with countdown
        wait(driver, 'id', 'submit-btn', 10)         # Wait max 10s with countdown
        wait(driver, 'class', 'content', 30, countdown=False)  # Wait silently for 30s

        # Wait for color at pixel
        wait(100, 200, 255, 0, 0)            # Wait for red at (100,200) with countdown
        wait(100, 200, 255, 0, 0, 30)        # Wait for red, max 30s with countdown
        wait(500, 300, 0, 255, 0, 60, countdown=False)  # Wait silently

    Returns:
        True if successful, False if error or timeout
    """

    # if no argument is passed
    if len(args) == 0:
        args = (3,)  # Set default to 3 seconds

    # ============================================================
    # MODE 1: COUNTDOWN (1 argument, integer or float)
    # ============================================================
    if len(args) == 1 and isinstance(args[0], (int, float)):
        seconds = args[0]

        if seconds < 0:
            raise ValueError("Seconds must be non-negative")

        if seconds == 0:
            return True  # No wait needed

        # Convert to int for countdown
        seconds_int = int(seconds)

        if countdown:
            print('Waiting:')
            for remaining in range(seconds_int, 0, -1):
                print(f'..{remaining} ', end='\r', flush=True)
                time.sleep(1)
            print('Done   ')
        else:
            time.sleep(seconds)

        return True

    # ============================================================
    # MODE 2: WAIT FOR ELEMENT (driver object + 2 or 3 arguments)
    # ============================================================
    elif len(args) >= 3 and hasattr(args[0], 'find_element'):
        # Selenium mode - driver object detected
        driver_obj = args[0]
        selector_type = args[1]
        selector = args[2]
        timeout = args[3] if len(args) == 4 else 180  # Default 180s (3 minutes)

        # Validate selector type
        if selector_type not in ['id', 'xpath', 'class', 'name', 'css', 'tag', 'text', 'partial']:
            raise ValueError(f"Invalid selector type: {selector_type}")

        if countdown:
            print(f'Waiting for element ({selector_type}: {selector}):')

        # Save original implicit wait and set to 0 for fast checking
        original_implicit_wait = driver_obj.timeouts.implicit_wait
        driver_obj.implicitly_wait(0)

        start_time = time.time()

        # Countdown approach: check element every second, print countdown
        for remaining in range(timeout, 0, -1):
            # Check if element exists using get_web_element
            element = _get_web_element(driver_obj, selector_type, selector)

            # Element found
            if element:
                elapsed = time.time() - start_time
                # Restore original implicit wait
                driver_obj.implicitly_wait(original_implicit_wait)
                print(f"Element found after {elapsed:.1f}s")
                return True

            # Print countdown
            if countdown:
                print(f'..{remaining} ', end='\r', flush=True)

            # Wait 1 second before next check
            time.sleep(1)

        # Restore original implicit wait
        driver_obj.implicitly_wait(original_implicit_wait)

        # Timeout reached
        print(f"Element not found after {timeout}s timeout.")
        return False

    # ============================================================
    # MODE 3: WAIT FOR COLOR AT PIXEL (5 or 6 arguments, all integers)
    # ============================================================
    elif len(args) in [5, 6] and all(isinstance(arg, int) for arg in args):
        x = args[0]
        y = args[1]
        target_r = args[2]
        target_g = args[3]
        target_b = args[4]
        timeout = args[5] if len(args) == 6 else 180  # Default 180s

        if countdown:
            print(f'Waiting at ({x},{y}) for color RGB({target_r},{target_g},{target_b}) :')

        start_time = time.time()

        # Countdown approach: check color every second, print countdown
        for remaining in range(timeout, 0, -1):
            try:
                # Check if color matches at pixel
                if pyautogui.pixelMatchesColor(x, y, (target_r, target_g, target_b)):
                    elapsed = time.time() - start_time
                    print(f"Color found after {elapsed:.1f}s")
                    return True

            except Exception as e:
                raise RuntimeError(f"Error checking pixel color: {e}")

            # Print countdown
            if countdown:
                print(f'..{remaining} ', end='\r', flush=True)

            # Wait 1 second before next check
            time.sleep(1)

        # Timeout reached
        print(f"Color not found after {timeout}s timeout.")
        return False

    # Invalid arguments
    else:
        raise ValueError("Invalid arguments for wait()")


def wait_download(timeout_in_min=20, download_dir=None):
    """
    Waits for a download to complete by monitoring the downloads folder.

    Args:
        timeout_in_min: Maximum minutes to wait for download completion
        download_dir: Custom download directory (optional)
            - If provided: Uses specified path
            - If None: Auto-detects (checks env var, Docker, then default)

    Environment Variables:
        DOWNLOAD_DIR: Custom download directory path

    Returns:
        True if download completed, False if timeout

    Examples:
        wait_download(5)                           # Auto-detect download folder
        wait_download(10, '/downloads')            # Docker with custom path
        wait_download(5, 'D:/MyDownloads')         # Windows custom path
        wait_download(3, '/home/user/Downloads')   # Linux custom path

    Note: If a file was modified within the last 20 seconds before calling this function,
    it will be detected as a recently completed download and the function will return True
    immediately. This handles cases where downloads complete very quickly (before monitoring starts).

    """

    timeout_in_sec = timeout_in_min * 60

    # Determine download directory with smart detection
    if download_dir is not None:
        download_dir = str(download_dir)
    elif os.getenv('DOWNLOAD_DIR'):
        download_dir = os.getenv('DOWNLOAD_DIR')
        print(f'Using DOWNLOAD_DIR from environment: {download_dir}')
    elif os.path.exists('/.dockerenv'):
        download_dir = '/downloads'
        print(f'Docker detected, using: {download_dir}')
    else:
        download_dir = str(Path.home() / "Downloads")

    # Validate directory exists
    if not os.path.exists(download_dir):
        print(f"Warning: Download directory does not exist: {download_dir}")
        print(f"Attempting to create it...")
        try:
            os.makedirs(download_dir, exist_ok=True)
            print(f"Created directory: {download_dir}")
        except Exception as e:
            print(f"Error: Could not create directory: {e}")
            return False

    print(f'Monitoring downloads in: {download_dir}')

    # Temporary file extensions for different browsers
    temp_extensions = ('.crdownload', '.part', '.download', '.tmp', '.temp')

    # Get initial state and timestamp
    try:
        initial_files = set(os.listdir(download_dir))
        initial_temp_files = set([f for f in initial_files if f.endswith(temp_extensions)])
        function_start_time = time.time()

        # Get initial sizes of existing temp files
        initial_temp_sizes = {}
        for f in initial_temp_files:
            try:
                file_path = os.path.join(download_dir, f)
                initial_temp_sizes[f] = os.path.getsize(file_path)
            except Exception as e:
                print(f"Could not get file size of {f}: {e}")

    except Exception as e:
        print(f"Error accessing download directory: {e}")
        return False

    # Wait 10 seconds before checking
    print("Checking for active downloads...")
    time.sleep(10)

    # Check for recently downloaded files (within last 20 seconds)
    try:
        current_files_after_wait = set(os.listdir(download_dir))
        current_time = time.time()

        # Look for files modified in last 20 seconds
        recent_files = []
        for f in current_files_after_wait:
            if not f.endswith(temp_extensions):
                try:
                    file_path = os.path.join(download_dir, f)
                    mtime = os.path.getmtime(file_path)
                    age = current_time - mtime

                    # File modified in last 20 seconds
                    if age <= 20:
                        recent_files.append((f, age))
                except Exception as e:
                    print(f"Could not get modification time of {f}: {e}")

        if recent_files:
            # Sort by age (most recent first)
            recent_files.sort(key=lambda x: x[1])
            most_recent_file, age = recent_files[0]

            print(f"Quick download detected: '{most_recent_file}' (modified {int(age)} seconds ago)")
            print("Download already completed")
            return True

    except Exception as e:
        print(f"Error checking for recent files: {e}")

    download_time = 10  # Already waited 10 seconds
    last_print_time = 0
    download_started = False
    monitoring_files = set()

    while download_time < timeout_in_sec:
        try:
            current_files = set(os.listdir(download_dir))
            current_temp_files = set([f for f in current_files if f.endswith(temp_extensions)])

            # Get current sizes of temp files
            current_temp_sizes = {}
            for f in current_temp_files:
                try:
                    file_path = os.path.join(download_dir, f)
                    current_temp_sizes[f] = os.path.getsize(file_path)
                except Exception as e:
                    print(f"Could not get file size of {f}: {e}")

            # Find NEW temp files (appeared after function started)
            new_temp_files = current_temp_files - initial_temp_files

            # Check if OLD temp files are still active (size changed)
            active_old_temp_files = set()
            if download_time == 10:
                for f in initial_temp_files:
                    if f in current_temp_sizes and f in initial_temp_sizes:
                        if current_temp_sizes[f] != initial_temp_sizes[f]:
                            # File size changed - it's actively downloading
                            active_old_temp_files.add(f)
                            monitoring_files.add(f)

                if active_old_temp_files:
                    download_started = True
                    # Show full filenames (including extension)
                    if len(active_old_temp_files) == 1:
                        print(f"Active download detected: '{list(active_old_temp_files)[0]}'")
                    else:
                        print(
                            f"{len(active_old_temp_files)} active downloads detected: {', '.join(active_old_temp_files)}")

            # Track new temp files
            if new_temp_files:
                # Check if these are truly new (not already in monitoring_files)
                truly_new = new_temp_files - monitoring_files

                if truly_new:
                    if not download_started:
                        download_started = True

                    monitoring_files.update(truly_new)

                    # Show full filenames (including extension)
                    if len(truly_new) == 1:
                        print(f"New download started: '{list(truly_new)[0]}'")
                    else:
                        print(f"{len(truly_new)} new downloads started: {', '.join(truly_new)}")

            # Get all files we're monitoring (old active + new)
            all_monitoring = (current_temp_files & monitoring_files) | new_temp_files

            if all_monitoring:
                # Print progress every 10 seconds
                if download_time - last_print_time >= 10:
                    elapsed = str(datetime.timedelta(seconds=download_time))
                    # Show full filenames
                    if len(all_monitoring) == 1:
                        print(f"Waiting for '{list(all_monitoring)[0]}' to finish downloading... (elapsed: {elapsed})")
                    else:
                        print(
                            f"Waiting for {len(all_monitoring)} files to finish downloading: {', '.join(all_monitoring)} (elapsed: {elapsed})")
                    last_print_time = download_time

            elif download_started:
                # Temp files disappeared - download completed
                # Find new complete files
                new_files = current_files - initial_files
                completed_files = [f for f in new_files if not f.endswith(temp_extensions)]

                if completed_files:
                    elapsed = str(datetime.timedelta(seconds=download_time))
                    if len(completed_files) == 1:
                        print(f"Download completed: '{completed_files[0]}' (took {elapsed})")
                    else:
                        print(f"{len(completed_files)} downloads completed in {elapsed}:")
                        for f in completed_files:
                            print(f"  - {f}")
                    return True
                else:
                    # No new files found, keep waiting
                    if download_time - last_print_time >= 10:
                        elapsed = str(datetime.timedelta(seconds=download_time))
                        print(f"Verifying download... (elapsed: {elapsed})")
                        last_print_time = download_time

            else:
                # No downloads detected yet - check for recently downloaded files every 10 seconds
                if download_time - last_print_time >= 10:
                    # Check if any file was recently downloaded
                    recent_complete_files = []

                    for f in current_files:
                        if not f.endswith(temp_extensions) and f not in initial_files:
                            recent_complete_files.append(f)

                    if recent_complete_files:
                        # Found recently downloaded file
                        most_recent_file = recent_complete_files[0]

                        print(f"Quick download detected: '{most_recent_file}'")
                        print("Download already completed")
                        return True

                    elapsed = str(datetime.timedelta(seconds=download_time))
                    print(f"Waiting for download to start... (elapsed: {elapsed})")
                    last_print_time = download_time

        except Exception as e:
            print(f'Error monitoring downloads: {e}')
            return False

        time.sleep(1)
        download_time += 1

    # Timeout reached
    if monitoring_files or download_started:
        # Check what's still downloading
        current_files = set(os.listdir(download_dir))
        current_temp_files = set([f for f in current_files if f.endswith(temp_extensions)])
        still_downloading = current_temp_files & monitoring_files

        if still_downloading:
            if len(still_downloading) == 1:
                print(
                    f"Timeout after {timeout_in_min} minutes while waiting for '{list(still_downloading)[0]}' to complete")
            else:
                print(
                    f"Timeout after {timeout_in_min} minutes while waiting for {len(still_downloading)} files: {', '.join(still_downloading)}")
        else:
            print(f'Timeout after {timeout_in_min} minutes. Download status unclear.')
    else:
        print(f'Timeout after {timeout_in_min} minutes. No download detected.')

    return False


def wait_retry(x_wait, y_wait, target, timeout=90, x_retry_click=None, y_retry_click=None, tolerance=0,
               wait_interval=3):
    """
    This fucntion was earlier named as wait, but as we have created a new wait by replacing delay,
    we are renaming this one to wait_retry()

    Waits for a specific color or a phrase at a screen location. Optionally clicks another location if not found.

    Args:
        x_wait, y_wait: Coordinates to check for the color/phrase.
        target: RGB color (tuple) or phrase (string) to match.
        timeout: Time in seconds to wait before giving up (default 90 seconds).
        x_retry_click, y_retry_click: Optional coordinates to click if the color/phrase is not found.
        tolerance: The tolerance for color matching (default is 0).
        wait_interval: Time in seconds to wait between checks (default 3 seconds).

    Returns:
        True if the color/phrase is found, False if not found within the timeout.

    Example:
        wait(300, 300, (255, 0, 0)) # wait for a color at a coordinate
        wait(300, 300, "test") # wait for a string at a coordinate
        wait(300, 300, "test", 30) # wait for a string at a coordinate with timeout 30 sec
        wait(300, 300, "test", 30, 1500, 1500) # timeout 30 sec,click at another coordinate if not found as a retry
        wait(300, 300, "test", x_retry_click=1500, y_retry_click=1500) # default timeout, retry coordinate present
        wait(300, 300, "test", 60, wait_interval=10) # wait for string with timeout 60 sec and check every 10 sec
        wait(300, 300, (255, 0, 0), tolerance=20) # wait for colour with tolerance of 20
    """
    end_time = time.time() + timeout
    while True:
        current_time = time.time()
        if current_time > end_time:
            break

        if isinstance(target, tuple):
            rgb_color = cast(Tuple[int, int, int], target)
            if pyautogui.pixelMatchesColor(x_wait, y_wait, rgb_color, tolerance=tolerance):
                print(f"Color {rgb_color} found at ({x_wait}, {y_wait}).")
                return True

        elif isinstance(target, str):
            page_content = copy(x_wait, y_wait)  # Ensure 'copy' function is defined
            if target in page_content:
                print(f"Phrase '{target}' found.")
                return True

        if x_retry_click is not None and y_retry_click is not None:
            pyautogui.click(x_retry_click, y_retry_click)

        remaining_time = int(end_time - current_time)
        print(f"Waiting. Time remaining: {remaining_time} seconds", end='\r')
        time.sleep(wait_interval)

    print(f"\nTimeout reached. Target not found.")
    return False

def window(action=None, target=None, *args):
    """
    Unified window management function for Windows and Linux.

    Platform Support:
         Windows 10/11
         Linux (requires wmctrl and xdotool)
         macOS (not supported)

    Args:
        action: Window operation to perform (default: 'list')
            - None or 'list': Get all window titles
            - 'title': Get active window title (or find full title if target provided)
            - 'focus': Bring window to foreground
            - 'close': Close window
            - 'minimize': Minimize window
            - 'maximize': Maximize window
            - 'resize': Resize window (requires width, height)
            - 'move': Move window (requires x, y)
        target: Window title or pattern (required for most actions)
        *args: Additional parameters (width, height for resize; x, y for move)

    Examples:
        # Get all windows (default)
        window()                                    # ['Chrome', 'Notepad', 'Excel']
        window('list')                              # ['Chrome', 'Notepad', 'Excel']

        # Check if window exists
        if 'Chrome' in window():
            print("Chrome is open!")

        # Get active window title
        window('title')                             # 'Google Chrome - New Tab'

        # Find full title containing text
        window('title', 'Chrome')                   # 'Google Chrome - New Tab'
        window('title', 'Note')                     # 'Untitled - Notepad'

        # Window operations
        window('focus', 'Chrome')                   # Focus window
        window('close', 'Notepad')                  # Close window
        window('minimize', 'Excel')                 # Minimize window
        window('maximize', 'Word')                  # Maximize window

        # Position and size
        window('resize', 'Chrome', 800, 600)        # Resize to 800x600
        window('move', 'Chrome', 100, 200)          # Move to (100, 200)

        # Side-by-side setup (1920x1080 display)
        window('move', 'Chrome', 0, 0)              # Left side
        window('resize', 'Chrome', 960, 1080)       # Half screen width
        window('move', 'Code', 960, 0)              # Right side
        window('resize', 'Code', 960, 1080)         # Half screen width

        # Recording setup
        window('resize', 'Demo', 1280, 720)         # 720p size
        window('move', 'Demo', 320, 180)            # Centered on 1920x1080

    Linux Dependencies:
        sudo apt-get install wmctrl xdotool        # Ubuntu/Debian
        sudo yum install wmctrl xdotool            # RHEL/CentOS/Fedora

    Note:
        On Linux, resize and move automatically remove maximized/minimized
        state before applying changes, ensuring consistent behavior.

    Returns:
        - List of strings: When action is None or 'list'
        - String or None: When action is 'title'
        - True/False: When action is focus/close/minimize/maximize/resize/move

    Raises:
        ValueError: If invalid action or missing required parameters
        NotImplementedError: If called on macOS
    """

    system = platform.system()

    # Check if macOS and reject
    if system == "Darwin":
        raise NotImplementedError(
            "window() is not supported on macOS.\n"
            "Supported platforms: Windows, Linux\n"
            "Note: Selenium functions will still work on macOS."
        )

    # Default action: list all windows
    if action is None:
        action = 'list'

    action = action.lower()

    # Validate action
    valid_actions = ['list', 'title', 'focus', 'close', 'minimize', 'maximize', 'resize', 'move']
    if action not in valid_actions:
        raise ValueError(f"Invalid action '{action}'. Valid actions: {', '.join(valid_actions)}")

    # Check if target is required
    if action in ['focus', 'close', 'minimize', 'maximize', 'resize', 'move'] and target is None:
        raise ValueError(f"Action '{action}' requires a target window")

    try:
        # ============================================================
        # LIST - Get all window titles
        # ============================================================
        if action == 'list':
            if system == "Windows":
                titles = []

                def enum_handler(hwnd, ctx):
                    if win32gui.IsWindowVisible(hwnd):
                        title = win32gui.GetWindowText(hwnd)
                        if title.strip():
                            titles.append(title)

                win32gui.EnumWindows(enum_handler, None)
                return titles

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    titles = []
                    for line in result.stdout.strip().split('\n'):
                        parts = line.split(None, 3)
                        if len(parts) >= 4:
                            titles.append(parts[3])
                    return titles
                else:
                    print("Error: wmctrl not installed.")
                    print("Install with: sudo apt-get install wmctrl")
                    return []

        # ============================================================
        # TITLE - Get active window title OR find window title
        # ============================================================
        elif action == 'title':
            # No target - get active window title
            if target is None:
                if system == "Windows":
                    title = win32gui.GetWindowText(win32gui.GetForegroundWindow()).strip()
                    return title if title else None

                elif system == "Linux":
                    result = subprocess.run(['xdotool', 'getactivewindow', 'getwindowname'],
                                            capture_output=True, text=True)
                    if result.returncode == 0:
                        return result.stdout.strip()
                    else:
                        print("Error: xdotool not installed.")
                        print("Install with: sudo apt-get install xdotool")
                        return None

            # Target provided - find full title containing target text
            else:
                all_windows = window('list')

                # Case-insensitive contains match
                for win_title in all_windows:
                    if target.lower() in win_title.lower():
                        return win_title

                # Not found
                return None

        # ============================================================
        # FOCUS - Bring window to foreground
        # ============================================================
        elif action == 'focus':
            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")

                # Suggest similar windows
                suggestions = get_close_matches(target, all_windows, n=3, cutoff=0.6)
                if suggestions:
                    print(f"Did you mean: {', '.join(suggestions)}")

                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                    win32gui.SetForegroundWindow(hwnd)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-a', window_title],
                                        capture_output=True, text=True)
                return result.returncode == 0

        # ============================================================
        # CLOSE - Close window
        # ============================================================
        elif action == 'close':
            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-c', window_title],
                                        capture_output=True, text=True)
                return result.returncode == 0

        # ============================================================
        # MINIMIZE - Minimize window
        # ============================================================
        elif action == 'minimize':
            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.ShowWindow(hwnd, win32con.SW_MINIMIZE)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    for line in result.stdout.strip().split('\n'):
                        if window_title in line:
                            window_id = line.split()[0]
                            subprocess.run(['xdotool', 'windowminimize', window_id])
                            return True
                return False

        # ============================================================
        # MAXIMIZE - Maximize window
        # ============================================================
        elif action == 'maximize':
            if not window('focus', target):
                return False

            time.sleep(0.5)

            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.ShowWindow(hwnd, win32con.SW_MAXIMIZE)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-r', window_title, '-b',
                                         'add,maximized_vert,maximized_horz'],
                                        capture_output=True, text=True)
                return result.returncode == 0

        # ============================================================
        # RESIZE - Resize window to specific dimensions
        # ============================================================
        elif action == 'resize':
            if len(args) < 2:
                raise ValueError("Action 'resize' requires width and height: window('resize', 'Chrome', 800, 600)")

            try:
                width = int(args[0])
                height = int(args[1])
            except (ValueError, TypeError):
                raise ValueError("Width and height must be integers")

            if width <= 0 or height <= 0:
                raise ValueError("Width and height must be positive")

            if not window('focus', target):
                return False

            time.sleep(0.5)

            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    rect = win32gui.GetWindowRect(hwnd)
                    x, y = rect[0], rect[1]
                    win32gui.MoveWindow(hwnd, x, y, width, height, True)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    for line in result.stdout.strip().split('\n'):
                        if window_title in line:
                            window_id = line.split()[0]
                            # Remove maximized/minimized state before resizing
                            subprocess.run(['wmctrl', '-i', '-r', window_id, '-b', 'remove,maximized_vert,maximized_horz'])
                            subprocess.run(['xdotool', 'windowmap', window_id])  # restore if minimized
                            time.sleep(0.3)
                            subprocess.run(['wmctrl', '-i', '-r', window_id, '-e', f'0,-1,-1,{width},{height}'])
                            return True
                return False

        # ============================================================
        # MOVE - Move window to specific coordinates
        # ============================================================
        elif action == 'move':
            if len(args) < 2:
                raise ValueError("Action 'move' requires x and y coordinates: window('move', 'Chrome', 100, 200)")

            try:
                x = int(args[0])
                y = int(args[1])
            except (ValueError, TypeError):
                raise ValueError("X and Y coordinates must be integers")

            if x < 0 or y < 0:
                raise ValueError("X and Y coordinates must be non-negative")

            if not window('focus', target):
                return False

            time.sleep(0.5)

            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    rect = win32gui.GetWindowRect(hwnd)
                    width = rect[2] - rect[0]
                    height = rect[3] - rect[1]
                    win32gui.MoveWindow(hwnd, x, y, width, height, True)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    for line in result.stdout.strip().split('\n'):
                        if window_title in line:
                            window_id = line.split()[0]
                            # Remove maximized/minimized state before moving
                            subprocess.run(['wmctrl', '-i', '-r', window_id, '-b', 'remove,maximized_vert,maximized_horz'])
                            subprocess.run(['xdotool', 'windowmap', window_id])  # restore if minimized
                            time.sleep(0.3)
                            subprocess.run(['wmctrl', '-i', '-r', window_id, '-e', f'0,{x},{y},-1,-1'])
                            return True
                return False

    except Exception as e:
        print(f"Error in window operation: {e}")
        return False if action in ['focus', 'close', 'minimize', 'maximize', 'resize', 'move'] else None


def write(*keys):
    """
    Write or type text using keyboard (PyAutoGUI or Selenium).

    Args:
        *keys: Variable arguments depending on mode

    Examples:
        # PyAutoGUI mode (types in any active window)
        write("Hello World")                    # Types using PyAutoGUI
        write("user@example.com")               # Types email
        write("12345")                          # Types numbers (as string)
        write("12345")                          # Types numbers (as string)

        # Selenium mode - type on page (driver object)
        write(driver, "Hello World")            # Types on active element in browser
        write(driver, "Search query")           # Types in focused input

        # Selenium mode - type in specific element (driver object)
        write(driver, "id", "username", "john_doe")
        write(driver, "xpath", "//input[@name='email']", "user@example.com")
        write(driver, "class", "search-box", "Python tutorial")

    Returns:
        True if successful, False otherwise
    """

    try:
        # ============================================================
        # PYAUTOGUI MODE (1 argument - text only)
        # ============================================================
        if len(keys) == 1:
            pyautogui.typewrite(keys[0])
            time.sleep(1)
            return True

        # ============================================================
        # SELENIUM MODE (driver object detected)
        # ============================================================
        elif len(keys) >= 2 and hasattr(keys[0], 'find_element'):
            driver_obj = keys[0]

            # ----------------------------------------------------------
            # Type on page (driver, text)
            # ----------------------------------------------------------
            if len(keys) == 2:
                # write(driver, "Hello World")
                text = keys[1]
                action = ActionChains(driver_obj)
                action.send_keys(text).perform()
                return True

            # ----------------------------------------------------------
            # Type in specific element (driver, selector_type, selector, text)
            # ----------------------------------------------------------
            elif len(keys) == 4:
                # write(driver, "id", "username", "john")
                selector_type = keys[1]
                selector = keys[2]
                text = keys[3]

                element = _get_web_element(driver_obj, selector_type, selector)
                if element:
                    element.send_keys(text)
                    return True
                else:
                    print(f"Element not found: {selector_type} - {selector}")
                    return False

            else:
                print("Error: Invalid number of arguments for Selenium write()")
                print("Use: write(driver, text) or write(driver, selector_type, selector, text)")
                return False

        # ============================================================
        # ERROR - Invalid arguments
        # ============================================================
        else:
            print("Error: Invalid arguments for write()")
            print("Use: write(text) or write(driver, text) or write(driver, selector_type, selector, text)")
            return False

    except Exception as e:
        print(f"Error writing text: {e}")
        return False


def year():
    """
    Get current year.

    Returns:
        int: Current year (e.g., 2026)

    Example:
        year()  # 2026
    """
    return time.localtime().tm_year


def zoom(*args):
    """
    Zoom in/out using steps or set zoom percentage.

    Args:
        *args: Variable arguments
            - (value): PyAutoGUI zoom steps/reset
            - (driver, value): Selenium zoom steps/percentage/reset

    Value Logic:
        - -9 to +9: Zoom steps (Ctrl+/Ctrl-)
        - 100: Reset to default/100%
        - Outside range (except 100): Percentage (Selenium only)

    Examples:
        # PyAutoGUI (desktop apps)
        zoom(3)              # Zoom in 3 steps
        zoom(-5)             # Zoom out 5 steps
        zoom(100)            # Reset to default (Ctrl+0)

        # Selenium (browser) - Steps
        zoom(driver, 3)      # Zoom in 3 steps
        zoom(driver, -5)     # Zoom out 5 steps

        # Selenium (browser) - Reset
        zoom(driver, 100)    # Reset to 100%

        # Selenium (browser) - Percentage
        zoom(driver, 150)    # Set zoom to 150%
        zoom(driver, 75)     # Set zoom to 75%
        zoom(driver, 50)     # Set zoom to 50%
        zoom(driver, 200)    # Set zoom to 200%

    Returns:
        True if successful, False otherwise
    """

    # Determine mode
    if len(args) == 1:
        # PyAutoGUI mode
        use_driver = False
        driver_obj = None
        value = args[0]
    elif len(args) == 2 and hasattr(args[0], 'execute_script'):
        # Selenium mode - driver object detected
        use_driver = True
        driver_obj = args[0]
        value = args[1]
    else:
        raise ValueError("Invalid arguments. Use zoom(value) or zoom(driver, value)")

    # Validate value
    if not isinstance(value, int):
        raise ValueError("Zoom value must be an integer")

    # PyAutoGUI restrictions
    if not use_driver and value != 100 and abs(value) > 9:
        raise ValueError("PyAutoGUI mode supports steps (-9 to +9) or reset (100)")

    try:
        if use_driver:
            # ============================================================
            # SELENIUM MODE
            # ============================================================

            # Special case: 100 = Reset
            if value == 100:
                driver_obj.execute_script("document.body.style.zoom='100%'")
                print("Zoom reset to 100%")
                return True

            # Steps: -9 to +9
            elif -9 <= value <= 9:
                if value > 0:
                    # Zoom in
                    for _ in range(value):
                        driver_obj.find_element(By.TAG_NAME, 'body').send_keys(Keys.CONTROL, Keys.ADD)
                        time.sleep(0.3)
                    print(f"Zoomed in {value} step(s)")
                elif value < 0:
                    # Zoom out
                    steps = abs(value)
                    for _ in range(steps):
                        driver_obj.find_element(By.TAG_NAME, 'body').send_keys(Keys.CONTROL, Keys.SUBTRACT)
                        time.sleep(0.3)
                    print(f"Zoomed out {steps} step(s)")
                return True

            # Percentage: Outside -9 to +9 (excluding 100)
            else:
                if value < 1:
                    raise ValueError("Zoom percentage must be at least 1%")

                driver_obj.execute_script(f"document.body.style.zoom='{value}%'")
                print(f"Zoom set to {value}%")
                return True

        else:
            # ============================================================
            # PYAUTOGUI MODE (desktop apps)
            # ============================================================

            # Special case: 100 = Reset
            if value == 100:
                pyautogui.hotkey('ctrl', '0')
                print("Zoom reset to default")
                return True

            # Steps: -9 to +9
            elif value > 0:
                # Zoom in
                for _ in range(value):
                    pyautogui.hotkey('ctrl', '+')
                    time.sleep(0.3)
                print(f"Zoomed in {value} step(s)")
            elif value < 0:
                # Zoom out
                steps = abs(value)
                for _ in range(steps):
                    pyautogui.hotkey('ctrl', '-')
                    time.sleep(0.3)
                print(f"Zoomed out {steps} step(s)")

            return True

    except Exception as e:
        print(f"Error zooming: {e}")
        return False
