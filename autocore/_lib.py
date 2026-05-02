# AutoCore - Automate Core Actions
# Author: Ash
# GitHub: https://github.com/AshAutomates/AutoCore
# Docs: https://autocore.readthedocs.io
# PyPI: https://pypi.org/project/autocore
# Supports: Windows, Linux


#---------------standard library imports---------------

import atexit
import configparser
import csv
import datetime
import email
import io
import json
import logging
import os
import platform
import re
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import time
import traceback
import warnings
import wave
import xml.etree.ElementTree as et
from difflib import get_close_matches
from logging.handlers import RotatingFileHandler
from pathlib import Path

#---------------networking---------------
import requests
from bs4 import BeautifulSoup

#------------------------------------------------------------
# Image Processing
# Safe on all platforms including headless servers
# PIL.Image does not need a display : only PIL.ImageTk does
#------------------------------------------------------------
import numpy as np
from PIL import Image

#---------------Selenium------------------------
import undetected_chromedriver as uc
from selenium.common.exceptions import *
from selenium.webdriver.common.action_chains import ActionChains
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import Select

#--------------- Windows Only---------------
if platform.system() == "Windows":
    import win32con
    import win32gui

#------------------------------------------------------------
# Document Processing
# Safe on all platforms including headless servers
#------------------------------------------------------------
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import ebooklib
from ebooklib import epub
from odf import teletype
from odf.opendocument import load as odf_load
from odf.text import P
import extract_msg
from striprtf.striprtf import rtf_to_text
import yaml

#------------------------------------------------------------
# GUI (requires display)
# Needs: X display server on Linux
# Works on: Linux desktop, Windows, macOS
# Fails on: Linux headless servers, Docker, CI/CD, SSH sessions
# Equivalent of running: export DISPLAY=:0 && xhost +local:
#------------------------------------------------------------
if platform.system() == "Linux":
    if "DISPLAY" not in os.environ:
        os.environ["DISPLAY"] = ":0"
    if os.system("xhost +local: > /dev/null 2>&1") == 0:
        import tkinter as tk
        from PIL import ImageTk
        import pyautogui

        # Enable PyAutoGUI fail-safe: move mouse to any screen corner to abort script
        pyautogui.FAILSAFE = True
        _GUI_AVAILABLE = True
    else:
        tk = None
        ImageTk = None
        pyautogui = None
        _GUI_AVAILABLE = False
else:
    import tkinter as tk
    from PIL import ImageTk
    import pyautogui

    pyautogui.FAILSAFE = True
    _GUI_AVAILABLE = True

#---------------------------------------------------------------------------
# Audio (requires sound system)
# Works on: Linux desktop, Windows
# Fails on: Linux headless servers with no audio drivers
# piper-tts is loaded lazily inside say() : we just check
# here whether the platform has a usable audio system at all
#
# Linux check 1 : aplay command is available (alsa-utils installed)
# Linux check 2 : at least one real audio playback device exists
#                 on Ubuntu Server with no audio hardware,
#                 'aplay -l' returns non-zero with "no soundcards found"
#                 on Ubuntu Desktop / Linux Mint with real audio hardware,
#                 'aplay -l' returns zero : so _AUDIO_AVAILABLE = True
#---------------------------------------------------------------------------
if platform.system() == "Windows":
    _AUDIO_AVAILABLE = True  # winsound is always available on Windows
elif platform.system() == "Linux":
    aplay_exists = os.system("aplay --version > /dev/null 2>&1") == 0  # alsa-utils installed?
    audio_device_exists = os.system("aplay -l > /dev/null 2>&1") == 0  # real audio hardware present?
    _AUDIO_AVAILABLE = aplay_exists and audio_device_exists
else:
    _AUDIO_AVAILABLE = False  # unsupported platform

#------------------------------------------------------------
# Clipboard (requires clipboard manager)
# Works on: Linux desktop, Windows, macOS
# Fails on: Linux headless servers with no clipboard manager
#------------------------------------------------------------
try:
    import pyperclip

    _CLIPBOARD_AVAILABLE = True
except Exception:
    pyperclip = None
    _CLIPBOARD_AVAILABLE = False

# print current working directory so user knows where files will be saved
print(f"Current Working Directory: {os.getcwd()}")

# to avoid 'RuntimeError: maximum recursion depth exceeded'
sys.setrecursionlimit(1500)

# Global variables to track logging state
_log_file_handler: RotatingFileHandler | None = None
_original_stdout = None
_original_stderr = None
_log_folder: Path = Path("logs")
_script_had_error = False  # Track if unhandled exception occurred


def _preprocess_for_ocr(image):
    """
    Preprocesses an image to improve OCR accuracy before passing to EasyOCR.

    Args:
        image: numpy array of the image in RGB format

    Returns:
        numpy array: Cleaned binary image ready for EasyOCR

    Note:
        - Converts to grayscale to remove color noise.
        - Upscales 2x using cubic interpolation for better character recognition.
        - Applies denoising to remove background noise.
        - Applies Otsu thresholding to produce a clean black and white image.
        - Applies morphological closing to fill small gaps in characters.
        - Lazily imports cv2 to avoid slow startup when running 'from autocore import *'.
    """
    import cv2  # pip install opencv-python

    gray_image = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
    gray_image = cv2.resize(gray_image, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
    denoised = cv2.fastNlMeansDenoising(gray_image, h=10)
    threshold_img = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 1))
    cleaned = cv2.morphologyEx(threshold_img, cv2.MORPH_CLOSE, kernel)
    return cleaned


def _click_word_by_ocr(word_to_search, occurrences_to_click, button='left'):
    """
    Click on text found via OCR on screen.

    Args:
        word_to_search: Text to search for on screen
        occurrences_to_click: 0 = click all occurrences, N = click Nth occurrence only
        button: 'left' for left-click (default), 'right' for right-click

    Returns:
        True if any click happened, False otherwise

    Output:
        - Prints how many occurrences were found and which one(s) were clicked.
        - Prints a message if the requested occurrence number does not exist.
        - Prints a message if the text was not found on screen.

    Note:
        - Used internally by click() and click_right().
        - Cross-platform compatible (Windows, Linux).

    """

    import cv2  # pip install opencv-python

    click_count = 0

    try:
        # Get shared OCR reader
        reader = _get_ocr_reader()

        # Capture screen using pyautogui (cross-platform)
        screenshot_img = pyautogui.screenshot()

        # Convert PIL image to numpy array
        image = np.array(screenshot_img)

        # Convert RGB to BGR for OpenCV
        image = cv2.cvtColor(image, cv2.COLOR_RGB2BGR)
        cleaned = _preprocess_for_ocr(image)

        # Extract text with bounding boxes using EasyOCR
        results = reader.readtext(cleaned)

        # Filter matching text and collect positions
        matches = []
        for (bbox, text, confidence) in results:
            if word_to_search.lower() in text.lower():
                # Calculate center point (adjust for 2x upscaling)
                center_x = int((bbox[0][0] + bbox[2][0]) / 2 / 2)
                center_y = int((bbox[0][1] + bbox[2][1]) / 2 / 2)
                matches.append((center_x, center_y, text))

        # Sort matches by position (top-to-bottom, then left-to-right)
        matches.sort(key=lambda m: (m[1], m[0]))

        total_matches = len(matches)

        # Handle case when no matches found
        if total_matches == 0:
            print(f"'{word_to_search}' not found on screen")
            return False

        # Determine click action based on button parameter
        click_action = pyautogui.rightClick if button == 'right' else pyautogui.click
        action_name = "right-clicking" if button == 'right' else "clicking"

        # Click on matches based on occurrences_to_click
        if occurrences_to_click == 0:
            # Click all occurrences
            if total_matches == 1:
                print(f"1 occurrence of '{word_to_search}' found, {action_name} it")
            else:
                print(f"{total_matches} occurrences of '{word_to_search}' found, {action_name} all")
            for (x, y, text) in matches:
                click_action(x, y)
                click_count += 1
        elif 0 < occurrences_to_click <= total_matches:
            # Click specific occurrence (1-indexed)
            if total_matches == 1:
                print(f"1 occurrence of '{word_to_search}' found, {action_name} it")
            else:
                print(
                    f"{total_matches} occurrences of '{word_to_search}' found, {action_name} occurrence #{occurrences_to_click}")
            x, y, text = matches[occurrences_to_click - 1]
            click_action(x, y)
            click_count = 1
        else:
            # Requested occurrence doesn't exist
            if total_matches == 1:
                print(
                    f"1 occurrence of '{word_to_search}' found on screen, but you tried to {action_name[:-3]} occurrence #{occurrences_to_click} which doesn't exist")
            else:
                print(
                    f"{total_matches} occurrences of '{word_to_search}' found on screen, but you tried to {action_name[:-3]} occurrence #{occurrences_to_click} which doesn't exist")
            return False

        # Return True if any click happened
        return click_count > 0

    except Exception as e:
        print(f"Error in OCR click: {e}")
        return False


def _get_ocr_reader():
    """
    Returns a single shared OCR reader instance, creating it only once on first call
    and reusing it across all subsequent calls.

    Output:
        - Prints initialization message on first run.
        - Prints confirmation when OCR engine is ready.

    Note:
        - Lazily imports easyocr inside the function to avoid slow startup time
          when running 'from autocore import *' : easyocr is only loaded when
          OCR is actually needed.
        - Tries GPU first, falls back to CPU if GPU is unavailable.
        - Shared across all OCR functions (click, read, etc.).
        - First run downloads OCR models (~100MB).
        - Subsequent runs load instantly from cache.

    Returns:
        easyocr.Reader: Shared OCR reader instance
    """

    # Third-party imports - OCR
    import easyocr  # pip install "numpy<2" easyocr

    if not hasattr(_get_ocr_reader, 'reader'):
        try:
            print("Initializing OCR engine (first run downloads models, please wait)...")
            # Try GPU first
            _get_ocr_reader.reader = easyocr.Reader(['en'], gpu=True, verbose=False)
        except:
            # Fall back to CPU if GPU fails
            _get_ocr_reader.reader = easyocr.Reader(['en'], gpu=False, verbose=False)
        print("OCR engine ready.")
    return _get_ocr_reader.reader


def _get_web_element(driver_obj, selector_type, selector):
    """
    Locates and returns a single web element using Selenium.

    Args:
        driver_obj: Selenium WebDriver instance
        selector_type: Type of selector ('id', 'xpath', 'class', 'name', 'css', 'tag', 'text', 'partial')
        selector: Selector value/string

    Example:
        ::

            element = _get_web_element(driver, 'id', 'submit-button')
            element = _get_web_element(driver, 'xpath', '//button[@type="submit"]')
            element = _get_web_element(driver, 'class', 'btn-primary')

    Note:
        - Used internally by click(), click_right(), write(), copy(), wait() and other functions.
        - 'text' and 'partial' selector types are case-sensitive.
        - Returns None if element is not found instead of raising an exception.

    Returns:
        WebElement if found, None if not found
    """
    try:
        if selector_type == "id":
            return driver_obj.find_element(By.ID, selector)
        elif selector_type == "xpath":
            return driver_obj.find_element(By.XPATH, selector)
        elif selector_type == "class":
            return driver_obj.find_element(By.CLASS_NAME, selector)
        elif selector_type == "name":
            return driver_obj.find_element(By.NAME, selector)
        elif selector_type == "css":
            return driver_obj.find_element(By.CSS_SELECTOR, selector)
        elif selector_type == "tag":
            return driver_obj.find_element(By.TAG_NAME, selector)
        elif selector_type == "text":  # remember, text search is case-sensitive
            return driver_obj.find_element(By.LINK_TEXT, selector)
        elif selector_type == "partial":  # remember, partial text search is case-sensitive
            return driver_obj.find_element(By.PARTIAL_LINK_TEXT, selector)
        else:
            return None
    except NoSuchElementException:
        return None


class _CustomRotatingFileNameHandler(RotatingFileHandler):
    """
    Custom log rotation handler that extends RotatingFileHandler with two behaviours:

    1. Renames rotated files from log.txt.1 format to log_part_1.txt format
       to keep the .txt extension visible and readable in file explorers.
    2. Automatically deletes the oldest log files when the total logs folder
       size exceeds 100MB, keeping disk usage under control.

    Used internally by log_setup().
    """

    def rotation_filename(self, default_name):
        """
        Converts default rotation filename format to a cleaner format.

        Example:
            log_script.txt.1  -->  log_script_part_1.txt
        """
        match = re.search(r'\.txt\.(\d+)$', default_name)
        if match:
            part_num = match.group(1)
            return default_name.replace(f'.txt.{part_num}', f'_part_{part_num}.txt')
        return default_name

    def doRollover(self):
        """
        Performs log rotation and triggers cleanup of old logs if folder exceeds 100MB.
        """
        # Do the normal rotation first
        super().doRollover()

        # Now cleanup old logs if folder exceeds 100MB
        self._cleanup_old_logs()

    def _cleanup_old_logs(self):
        """
        Deletes the oldest log files in the logs folder when total size exceeds 100MB,
        removing files one by one until the folder is back under 100MB.

        Output:
            - Prints name and remaining folder size after each deleted file.
            - Prints warning if a file could not be deleted.
            - Prints warning if cleanup itself fails.
        """
        try:
            log_folder = Path("logs")
            if not log_folder.exists():
                return

            log_files = []
            total_size = 0

            for log_file in log_folder.glob("log_*.txt*"):
                size = log_file.stat().st_size
                mtime = log_file.stat().st_mtime
                log_files.append((log_file, size, mtime))
                total_size += size

            # Convert to MB
            total_size_mb = total_size / (1024 * 1024)
            max_total_mb = 100

            # Keep deleting oldest files until under 100MB
            if total_size_mb > max_total_mb:
                # Sort by modification time (oldest first)
                log_files.sort(key=lambda x: x[2])

                for log_file, size, _ in log_files:
                    if total_size_mb <= max_total_mb:
                        break

                    try:
                        log_file.unlink()  # deletes the file from the filesystem.
                        total_size_mb -= size / (1024 * 1024)
                        print(f"Deleted old log file: {log_file.name} (folder size: {total_size_mb:.1f}MB)")
                    except Exception as e:
                        print(f"Warning: Could not delete {log_file.name}: {e}")

        except Exception as e:
            print(f"Warning: Log cleanup failed: {e}")


class _LogCapture:
    """
    Intercepts stdout and stderr to mirror all print statements to both
    the terminal and the log file simultaneously.

    Used internally by log_setup() to redirect sys.stdout and sys.stderr.

    Note:
        - Uses a guard flag (_logging) to prevent infinite recursion that would
          occur if the logger itself triggers a write() call.
        - flush() is implemented to satisfy the stream interface expected by
          sys.stdout and sys.stderr.
    """

    def __init__(self, original_stream, logger, level):
        """
        Args:
            original_stream: The original sys.stdout or sys.stderr stream
            logger: The logging.Logger instance to write to
            level: Logging level (logging.INFO for stdout, logging.ERROR for stderr)
        """

        self.original_stream = original_stream
        self.logger = logger
        self.level = level
        self._logging = False  # guard flag to prevent infinite recursion

    def write(self, message):
        """
        Writes message to both terminal and log file.
        Skips empty messages and guards against recursive calls.
        """

        self.original_stream.write(message)
        self.original_stream.flush()

        # Write to log file (strip to avoid double newlines)
        # guard flag prevents recursive calls when logger itself triggers write()
        if message.strip() and not self._logging:
            self._logging = True  # block any further recursive calls
            try:
                self.logger.log(self.level, message.rstrip())
            finally:
                self._logging = False  # always release guard even if exception occurs

    def flush(self):
        """
        Flushes the original stream to satisfy the stream interface.
        """
        self.original_stream.flush()


def browser(url, headless=False, timeout=30, cookie_path=None):
    """
    Initialize and return a browser instance for web automation.

    Args:
        url: Target URL to navigate to
        headless: Run browser in headless mode (default: False)
        timeout: Maximum seconds to wait for elements to appear (default: 30)
        cookie_path: Path to cookies JSON file (optional)
            ::

                - Cookies MUST be in JSON format
                - Export from Chrome using "Cookie-Editor" extension
                - Cookie domain must match the target URL

    Returns:
        WebDriver: Browser instance or None if initialization fails

    Example:
        ::

            # Basic usage
            driver = browser('https://google.com')
            click(driver, 'id', 'search-button')

            # Slow-loading site
            driver = browser('https://slow-site.gov', timeout=90)
            click(driver, 'id', 'submit-btn')  # Waits up to 90s

            # Fast site testing
            driver = browser('https://fast-site.com', timeout=5)
            click(driver, 'id', 'login-btn')  # Fails fast in 5s

            # With cookies
            driver = browser('https://site.com', cookie_path='cookies.json')

            # Headless mode
            driver = browser('https://google.com', headless=True)

            # Trigger a download and wait for it
            driver = browser('https://example.com')
            click(driver, 'id', 'download-button')
            wait_download(download_dir=driver.download_dir)

    Note:
        Uses undetected-chromedriver (uc) to bypass bot detection.
        Requires Google Chrome to be installed.

        Windows:
            winget install Google.Chrome

        Linux (Ubuntu/Debian/Mint):
            wget https://dl.google.com/linux/direct/google-chrome-stable_current_amd64.deb
            sudo dpkg -i google-chrome-stable_current_amd64.deb
            sudo apt-get install -f -y

        Linux (RHEL/CentOS/Fedora):
            wget https://dl.google.com/linux/direct/google-chrome-stable_current_x86_64.rpm
            sudo rpm -i google-chrome-stable_current_x86_64.rpm
    """

    # Auto-add https:// if protocol is missing
    if not url.startswith(('http://', 'https://')):
        url = 'https://' + url
        print(f"Protocol not specified. Using: {url}")

    # Initialize options for Chrome
    options = uc.ChromeOptions()

    # --------------- HEADLESS MODE ---------------
    if headless:
        options.add_argument("--headless=new")
        options.add_argument("--window-size=2560,1440")
        options.add_argument("--disable-gpu")  # disable GPU (prevents SwiftShader WebGL signal)
        options.add_argument("--enable-features=OverlayScrollbar")  # mimics real Chrome scrollbar behavior

    # --------------- CHROME VERSION ---------------
    # Fetch installed Chrome version dynamically to keep user-agent current.
    # A matching Chrome version in user-agent makes the browser look more like a real user.
    # Initialize to None so major_version check below does not raise NameError if detection fails.
    chrome_version = None
    try:
        system = platform.system()
        if system == "Windows":
            result = subprocess.run(
                ['reg', 'query', 'HKEY_CURRENT_USER\\Software\\Google\\Chrome\\BLBeacon', '/v', 'version'],
                capture_output=True, text=True)
            chrome_version = result.stdout.strip().split()[-1]
        elif system == "Linux":
            result = subprocess.run(['google-chrome', '--version'], capture_output=True, text=True)
            chrome_version = result.stdout.strip().split()[-1]
        user_agent = f"Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/{chrome_version} Safari/537.36"
    except Exception:
        # If Chrome version detection fails, use a generic user-agent without version number
        user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome Safari/537.36"

    # Set a user-agent string to make the automated browser look like a regular Chrome browser
    options.add_argument(f"--user-agent={user_agent}")

    # Additional options to enhance realism and disable Selenium detection
    options.add_argument('--start-maximized')  # Start browser maximized
    options.add_argument("--safebrowsing-disable-download-protection")  # allow unverified downloads
    options.add_argument("--disable-features=InsecureDownloadWarnings")  # disable insecure download warnings
    options.add_argument("--disable-features=DownloadBubble")  # disable download bubble UI
    options.add_argument("--no-sandbox")  # disable sandbox restrictions
    options.add_argument("--lang=en-US")  # set language to avoid empty navigator.languages
    options.add_argument("--disable-translate")  # disable Google Translate popup
    options.add_argument("--no-first-run")  # skip first-run welcome popups
    options.add_argument("--no-default-browser-check")  # skip default browser prompt

    # --------------- DOWNLOAD DIRECTORY ---------------
    # Resolve download directory using same priority as wait_download()
    # Pre-create it before Chrome launches so Chrome finds and uses it
    if os.getenv('DOWNLOAD_DIR'):
        download_dir = os.getenv('DOWNLOAD_DIR')
    elif os.path.exists('/.dockerenv'):
        download_dir = '/downloads'
    else:
        download_dir = str(Path.home() / 'Downloads')
    os.makedirs(download_dir, exist_ok=True)

    # --------------- CHROME PREFERENCES ---------------
    # Set preferences to avoid unnecessary pop-ups and block notifications
    prefs = {
        "download.default_directory": download_dir,  # explicitly set download folder
        "profile.default_content_setting_values.notifications": 2,  # block browser notification popups
        'credentials_enable_service': False,  # disable save password popup
        'profile': {'password_manager_enabled': False},  # disable password manager completely
        'profile.password_manager_leak_detection': False,  # breach popup
        "profile.default_content_setting_values.notifications": 2,  # block notification popups
        "translate.enabled": False,  # disable translate bar
        "autofill.profile_enabled": False,  # disable autofill suggestions
        "autofill.credit_card_enabled": False,  # disable credit card autofill
        "safebrowsing.enabled": False,  # disable safe browsing to allow downloads
        "download.prompt_for_download": False,  # no download confirmation prompt
        "download.directory_upgrade": True,  # allow download directory changes
        "plugins.always_open_pdf_externally": True,  # download PDFs instead of opening in browser
        "safebrowsing_for_trusted_sources_enabled": False,  # disable safe browsing for trusted sources
        "safebrowsing.disable_download_protection": True,  # disable download protection completely
    }
    options.add_experimental_option("prefs", prefs)

    # --------------- INITIALIZE DRIVER ---------------
    try:
        # Extract major version from detected Chrome version to prevent ChromeDriver mismatch
        major_version = int(chrome_version.split('.')[0]) if chrome_version else None
        # Pass major version to force matching ChromeDriver download
        driver_instance = uc.Chrome(options=options, version_main=major_version)

    except Exception as e:
        print(f"Error initializing Chrome Driver: {e}")
        return None

    # --------------- BOT DETECTION PATCHES ---------------
    # Inject JS to patch remaining bot detection signals
    driver_instance.execute_cdp_cmd("Page.addScriptToEvaluateOnNewDocument", {
        "source": """
            // Fix navigator.webdriver at prototype level
            Object.defineProperty(navigator, 'webdriver', {
                get: () => false,
                configurable: true,
                enumerable: true
            });
            try { delete navigator.__proto__.webdriver; } catch(e) {}

            // Fix plugins to return proper PluginArray-like object
            const mockPlugins = ['Chrome PDF Plugin', 'Chrome PDF Viewer', 'Native Client'].map(name => {
                const plugin = Object.create(Plugin.prototype);
                Object.defineProperty(plugin, 'name', {get: () => name});
                Object.defineProperty(plugin, 'filename', {get: () => name.toLowerCase().replace(/ /g, '-') + '.dll'});
                Object.defineProperty(plugin, 'description', {get: () => name});
                Object.defineProperty(plugin, 'length', {get: () => 0});
                return plugin;
            });
            Object.defineProperty(navigator, 'plugins', {
                get: () => {
                    const arr = Object.create(PluginArray.prototype);
                    mockPlugins.forEach((p, i) => arr[i] = p);
                    Object.defineProperty(arr, 'length', {get: () => mockPlugins.length});
                    arr.item = i => mockPlugins[i];
                    arr.namedItem = name => mockPlugins.find(p => p.name === name) || null;
                    arr.refresh = () => {};
                    return arr;
                }
            });

            // Fix permissions to return 'prompt' instead of 'denied'
            const originalQuery = window.navigator.permissions.query.bind(window.navigator.permissions);
            window.navigator.permissions.query = (parameters) => {
                const permissionsToSpoof = ['notifications', 'push', 'midi', 'camera', 'microphone', 'geolocation'];
                if (permissionsToSpoof.includes(parameters.name)) {
                    return Promise.resolve({
                        state: 'prompt',
                        permission: 'prompt',
                        status: 'prompt',
                        onchange: null
                    });
                }
                return originalQuery(parameters);
            };

            // Also patch inside iframes
            if (window.top !== window.self) {
                window.navigator.permissions.query = (parameters) => {
                    return Promise.resolve({state: 'prompt', permission: 'prompt', onchange: null});
                };
            }

            // Fix navigator.languages
            Object.defineProperty(navigator, 'languages', {get: () => ['en-US', 'en']});

            // Fix navigator.mimeTypes
            Object.defineProperty(navigator, 'mimeTypes', {get: () => [
                {type: 'application/pdf'},
                {type: 'application/x-nacl'},
            ]});

            // Fix window.chrome
            window.chrome = {
                webstore: {
                    onInstallStageChanged: {},
                    onDownloadProgress: {},
                    install: function() {},
                    constructor: function() {}
                },
                runtime: {
                    onConnect: null,
                    onMessage: null,
                    connect: function() {},
                    sendMessage: function() {},
                    id: undefined
                },
                loadTimes: function() {},
                csi: function() {},
                app: {
                    isInstalled: false,
                    InstallState: {DISABLED: 'disabled', INSTALLED: 'installed', NOT_INSTALLED: 'not_installed'},
                    RunningState: {CANNOT_RUN: 'cannot_run', READY_TO_RUN: 'ready_to_run', RUNNING: 'running'}
                }
            };

            // Fix WebGL renderer
            const getParameter = WebGLRenderingContext.prototype.getParameter;
            WebGLRenderingContext.prototype.getParameter = function(parameter) {
                if (parameter === 37445) return 'Intel Inc.';
                if (parameter === 37446) return 'Intel Iris OpenGL Engine';
                return getParameter.call(this, parameter);
            };
        """
    })

    # Set an implicit wait for elements to be found
    driver_instance.implicitly_wait(timeout)

    # Load a blank page to initialize the driver properly before cookie injection
    driver_instance.get("about:blank")

    # --------------- COOKIES ---------------
    # Load cookies from the specified path, if available.
    if cookie_path and os.path.exists(cookie_path):
        # Load target page to inject cookie
        driver_instance.get(url)

        try:
            with open(cookie_path, "r", encoding="utf-8") as f:
                cookies = json.load(f)

            for cookie in cookies:
                # Remove fields added by browser extensions that Selenium does not support
                cookie.pop("sameSite", None)
                cookie.pop("hostOnly", None)
                cookie.pop("session", None)
                cookie.pop("storeId", None)

                # Selenium rejects cookies with a leading dot in the domain, while browsers accept them, so the dot must be removed for compatibility.
                if cookie.get("domain", "").startswith("."):
                    cookie["domain"] = cookie["domain"][1:]

                driver_instance.add_cookie(cookie)

            # Reload page so cookies show effect
            driver_instance.get(url)

        except Exception as e:
            print(f"Cookie loading failed: {e}")
            return None

    else:
        # No cookies provided so single page load is enough
        driver_instance.get(url)

    # Suppress undetected_chromedriver's __del__ cleanup error on exit
    driver_instance.__class__.__del__ = lambda self: None

    # Store resolved download directory so caller can pass it to wait_download()
    driver_instance.download_dir = download_dir

    # Return the driver instance - user can name it anything!
    return driver_instance


def click(*where):
    """
    Performs left-click based on different input types.

    Modes:
        1. Image matching: Click on visual element
        2. OCR text matching: Click on text found on screen (with nth occurrence support)
        3. Coordinates: Click at specific x, y position
        4. Color matching: Click on specific color in region
        5. Selenium web element: Click on element in browser

    Args:
        *where: Variable arguments depending on click mode

    Returns:
        True if successful, False otherwise

    Example:
        ::

            # Image matching
            click('button.png')

            # OCR text matching
            click('Submit')                    # Click first occurrence
            click('Submit', 2)                 # Click 2nd occurrence
            click('Login', 0)                  # Click all occurrences

            # Coordinates
            click(100, 200)

            # Color matching in region
            click(x1, y1, x2, y2, r, g, b)              # Find and click color
            click(x1, y1, x2, y2, r, g, b, tolerance)   # With tolerance

            # Selenium (pass driver object first)
            click(driver, 'id', 'submit-button')
            click(driver, 'xpath', '//button[@id="submit"]')
            click(driver, 'class', 'btn-primary')
            click(driver, 'name', 'username')
            click(driver, 'css', 'button.submit')
            click(driver, 'tag', 'button')
            click(driver, 'text', 'Click Here')
            click(driver, 'partial', 'Click')
    """
    import cv2

    # --------------- SELENIUM ---------------
    if len(where) > 0 and hasattr(where[0], 'find_element'):
        driver_obj = where[0]

        if len(where) < 3:
            print("Error: Selenium click requires 3 arguments: click(driver, selector_type, selector)")
            return False

        selector_type = where[1]
        selector = where[2]

        try:
            element = _get_web_element(driver_obj, selector_type, selector)
            if element:
                driver_obj.execute_script("arguments[0].scrollIntoView({block: 'center'});", element)
                try:
                    element.click()
                except ElementClickInterceptedException:
                    # Fall back to JavaScript click when element is blocked by an overlay (e.g. ads)
                    driver_obj.execute_script("arguments[0].click();", element)
                return True
            else:
                print(f"Element not found: {selector_type} - {selector}")
                return False
        except NoSuchElementException:
            print(f"Element of {selector_type} - {selector} not found.")
            return False
        except Exception as e:
            print(f"Error clicking element: {e}")
            return False

    # --------------- PYAUTOGUI ---------------
    if not _GUI_AVAILABLE:
        print("Error: click() requires a display.")
        return False

    # --------------- IMAGE MATCHING / OCR TEXT ---------------
    elif len(where) == 1:
        if '.' in where[0]:
            pyautogui.click(where[0])
            return True
        else:
            result = _click_word_by_ocr(where[0], 1, button='left')
            return result

    # --------------- COORDINATES / OCR WITH OCCURRENCE ---------------
    elif len(where) == 2:
        if isinstance(where[0], int) and isinstance(where[1], int):
            pyautogui.click(where[0], where[1])
            return True
        elif isinstance(where[0], str) and isinstance(where[1], int):
            result = _click_word_by_ocr(where[0], where[1], button='left')
            return result
        else:
            print("Error: Invalid arguments for click()")
            return False

    # --------------- COLOR MATCHING IN REGION ---------------
    elif len(where) in [7, 8]:
        x_from, y_from, x_to, y_to, r, g, b = where[:7]
        tolerance = where[7] if len(where) == 8 else 0

        try:
            screenshot_img = pyautogui.screenshot(region=(x_from, y_from, x_to - x_from, y_to - y_from))
            screenshot_img = np.array(screenshot_img)
            screenshot_img = cv2.cvtColor(screenshot_img, cv2.COLOR_RGB2BGR)

            lower = np.array([b - tolerance, g - tolerance, r - tolerance])
            upper = np.array([b + tolerance, g + tolerance, r + tolerance])

            mask = cv2.inRange(screenshot_img, lower, upper)
            points = cv2.findNonZero(mask)

            if points is not None:
                click_x, click_y = points[0][0]
                pyautogui.click(x_from + click_x, y_from + click_y)
                print(f'Pixel found and clicked at ({x_from + click_x}, {y_from + click_y}).')
                return True

            print('Pixel not found.')
            return False

        except Exception as e:
            print(f"Error during color search: {e}")
            return False

    # ---------------Invalid Arguments---------------
    else:
        print("Error: Invalid arguments for click()")
        return False


def click_right(*where):
    """
    Performs right-click (context menu) based on different input types.

    Modes:
        1. Image matching: Right-click on visual element
        2. OCR text matching: Right-click on text found on screen (with nth occurrence support)
        3. Coordinates: Right-click at specific x, y position
        4. Color matching: Right-click on specific color in region
        5. Selenium web element: Right-click on element in browser

    Args:
        *where: Variable arguments depending on click mode

    Returns:
        True if successful, False otherwise

    Example:
        ::

            # Image matching
            click_right('button.png')

            # OCR text matching
            click_right('Submit')              # Right-click first occurrence
            click_right('Submit', 2)           # Right-click 2nd occurrence
            click_right('Login', 0)            # Right-click all occurrences

            # Coordinates
            click_right(100, 200)

            # Color matching in region
            click_right(x1, y1, x2, y2, r, g, b)              # Find and right-click color
            click_right(x1, y1, x2, y2, r, g, b, tolerance)   # With tolerance

            # Selenium (pass driver object first)
            click_right(driver, 'id', 'submit-button')
            click_right(driver, 'xpath', '//button[@id="submit"]')
            click_right(driver, 'class', 'btn-primary')
            click_right(driver, 'name', 'username')
            click_right(driver, 'css', 'button.submit')
            click_right(driver, 'tag', 'button')
            click_right(driver, 'text', 'Click Here')
            click_right(driver, 'partial', 'Click')
    """
    import cv2

    # --------------- SELENIUM ---------------
    if len(where) > 0 and hasattr(where[0], 'find_element'):
        driver_obj = where[0]

        if len(where) < 3:
            print("Error: Selenium click_right requires 3 arguments: click_right(driver, selector_type, selector)")
            return False

        selector_type = where[1]
        selector = where[2]

        try:
            element = _get_web_element(driver_obj, selector_type, selector)
            if element:
                driver_obj.execute_script("arguments[0].scrollIntoView({block: 'center'});", element)
                ActionChains(driver_obj).context_click(element).perform()
                return True
            else:
                print(f"Element not found: {selector_type} - {selector}")
                return False
        except NoSuchElementException:
            print(f"Element of {selector_type} - {selector} not found.")
            return False
        except Exception as e:
            print(f"Error right-clicking element: {e}")
            return False

    # --------------- PYAUTOGUI ---------------
    if not _GUI_AVAILABLE:
        print("Error: click_right() requires a display.")
        return False

    # --------------- IMAGE MATCHING / OCR TEXT ---------------
    elif len(where) == 1:
        if '.' in where[0]:
            pyautogui.rightClick(where[0])
            return True
        else:
            result = _click_word_by_ocr(where[0], 1, button='right')
            return result

    # --------------- COORDINATES / OCR WITH OCCURRENCE ---------------
    elif len(where) == 2:
        if isinstance(where[0], int) and isinstance(where[1], int):
            pyautogui.rightClick(where[0], where[1])
            return True
        elif isinstance(where[0], str) and isinstance(where[1], int):
            result = _click_word_by_ocr(where[0], where[1], button='right')
            return result
        else:
            print("Error: Invalid arguments for click_right()")
            return False

    # --------------- COLOR MATCHING IN REGION ---------------
    elif len(where) in [7, 8]:
        x_from, y_from, x_to, y_to, r, g, b = where[:7]
        tolerance = where[7] if len(where) == 8 else 0

        try:
            screenshot_img = pyautogui.screenshot(region=(x_from, y_from, x_to - x_from, y_to - y_from))
            screenshot_img = np.array(screenshot_img)
            screenshot_img = cv2.cvtColor(screenshot_img, cv2.COLOR_RGB2BGR)

            lower = np.array([b - tolerance, g - tolerance, r - tolerance])
            upper = np.array([b + tolerance, g + tolerance, r + tolerance])

            mask = cv2.inRange(screenshot_img, lower, upper)
            points = cv2.findNonZero(mask)

            if points is not None:
                click_x, click_y = points[0][0]
                pyautogui.rightClick(x_from + click_x, y_from + click_y)
                print(f'Pixel found and right-clicked at ({x_from + click_x}, {y_from + click_y}).')
                return True

            print('Pixel not found.')
            return False

        except Exception as e:
            print(f"Error during color search: {e}")
            return False

    # ---------------Invalid Arguments---------------
    else:
        print("Error: Invalid arguments for click_right()")
        return False


def copy(*where):
    """
    Copies text from various sources: screen, clipboard, Selenium elements or web pages.

    Modes:
        1. Active window: Copy all content from current window
        2. Clipboard: Get current clipboard content
        3. Screen coordinates: Click at position and copy
        4. Selenium webpage: Copy entire page content
        5. Selenium element: Copy element text or attribute value

    Args:
        *where: Variable arguments depending on copy mode

    Returns:
        str: Copied text or '' if nothing was copied

    Example:
        ::

            # Active window - Copy everything from current window
            # Ctrl+A, Ctrl+C from active window
            copy()

            # Clipboard
            # Get current clipboard content
            copy('clipboard')

            # Screen coordinates
            # Click at (500, 300) and copy
            copy(500, 300)

            # Selenium webpage - Copy entire page
            # Copy all webpage content
            copy(driver)

            # Selenium element - Copy text
            copy(driver, 'id', 'username-display')
            copy(driver, 'xpath', '//div[@class="name"]')
            copy(driver, 'class', 'user-info')
            copy(driver, 'name', 'description')
            copy(driver, 'css', 'p.content')
            copy(driver, 'tag', 'h1')
            copy(driver, 'text', 'Welcome')
            copy(driver, 'partial', 'Hello')

            # Selenium element - Copy attribute
            copy(driver, 'id', 'download-link', 'href')         # Get link URL
            copy(driver, 'class', 'product-img', 'src')         # Get image source
            copy(driver, 'id', 'email-field', 'value')          # Get input value
            copy(driver, 'xpath', '//a[@id="link"]', 'title')   # Get title attribute
    """

    result = None

    # --------------- MODE 1 : ACTIVE WINDOW ---------------
    # needs: GUI + clipboard
    if len(where) == 0:
        if not _GUI_AVAILABLE:
            print("Error: copy() requires a display.")
            return ''
        if not _CLIPBOARD_AVAILABLE:
            print("Error: copy() requires a clipboard manager.")
            return ''
        pyperclip.copy('')
        pyautogui.hotkey('ctrl', 'a')
        time.sleep(1)
        pyautogui.hotkey('ctrl', 'c')
        time.sleep(1)
        result = pyperclip.paste().strip()

    # --------------- ONE ARGUMENT ---------------
    elif len(where) == 1:

        # --------------- MODE 2 : CLIPBOARD ---------------
        # needs: clipboard only
        if where[0] == 'clipboard':
            if not _CLIPBOARD_AVAILABLE:
                print("Error: copy('clipboard') requires a clipboard manager.")
                return ''
            result = pyperclip.paste().strip()

        # --------------- MODE 3 : SELENIUM WEBPAGE ---------------
        # needs: clipboard (for paste after Ctrl+C)
        elif hasattr(where[0], 'find_element'):
            if not _CLIPBOARD_AVAILABLE:
                print("Error: copy(driver) requires a clipboard manager.")
                return ''
            driver_obj = where[0]
            pyperclip.copy('')
            try:
                action = ActionChains(driver_obj)
                action.key_down(Keys.CONTROL).send_keys('a').key_up(Keys.CONTROL).perform()
                time.sleep(1)
                action.key_down(Keys.CONTROL).send_keys('c').key_up(Keys.CONTROL).perform()
                time.sleep(1)
                result = pyperclip.paste().strip()
                pyperclip.copy('')
                driver_obj.execute_script("window.getSelection().removeAllRanges();")
            except Exception as e:
                print(f"Error copying webpage content: {e}")

        else:
            print(f"Invalid argument: {where[0]}")

    # --------------- TWO ARGUMENTS ---------------
    elif len(where) == 2:

        # --------------- MODE 4 : SCREEN COORDINATES ---------------
        # needs: GUI + clipboard
        if isinstance(where[0], int) and isinstance(where[1], int):
            if not _GUI_AVAILABLE:
                print("Error: copy(x, y) requires a display.")
                return ''
            if not _CLIPBOARD_AVAILABLE:
                print("Error: copy(x, y) requires a clipboard manager.")
                return ''
            pyperclip.copy('')
            pyautogui.click(where[0], where[1])
            time.sleep(1)
            pyautogui.hotkey('ctrl', 'a')
            time.sleep(1)
            pyautogui.hotkey('ctrl', 'c')
            time.sleep(1)
            result = pyperclip.paste().strip()

        elif hasattr(where[0], 'find_element'):
            print("Error: Selenium copy requires 3 arguments: copy(driver, selector_type, selector)")
            print("Example: copy(driver, 'id', 'username')")

        else:
            print("Invalid arguments for copy()")

    # --------------- MODE 5 : SELENIUM ELEMENT TEXT OR ATTRIBUTE ---------------
    # needs: nothing : reads directly from DOM, no clipboard needed
    elif len(where) in [3, 4]:
        if hasattr(where[0], 'find_element'):
            driver_obj = where[0]
            selector_type = where[1]
            selector = where[2]
            attribute = where[3] if len(where) == 4 else None

            try:
                element = _get_web_element(driver_obj, selector_type, selector)
                if element:
                    if attribute:
                        attr_value = element.get_attribute(attribute)
                        result = attr_value.strip() if attr_value else ''
                    else:
                        result = element.text.strip()
                else:
                    print(f"Element not found: {selector_type} - {selector}")
            except Exception as e:
                print(f"Error copying from element: {e}")
        else:
            print("Error: For Selenium mode, first argument must be driver object")

    # --------------- INVALID ARGUMENTS ---------------
    else:
        print("Invalid arguments provided for copy()")

    # --------------- RETURN RESULT ---------------
    if result == '' or result is None:
        print("No content returned from copy()")
        return ''
    else:
        return result


def csv_to_xlsx(csv_file=None, delete_csv=True):
    """
    Converts CSV file(s) to XLSX format.

    Args:
        csv_file: Path to CSV file or None to auto-detect single CSV in current directory
        delete_csv: If True, deletes original CSV after conversion (default: True)

    Returns:
        str: Path of created XLSX file or None if error

    Output:
            - Prints the detected CSV filename when auto-detected.
            - Prints conversion result showing source and destination filenames.
            - Prints confirmation when original CSV is deleted.

    Example:
        ::

            # Auto-detect single CSV in current directory (deletes CSV by default)
            csv_to_xlsx()                               # Finds, converts and deletes CSV

            # Specific file (deletes CSV by default)
            csv_to_xlsx('data.csv')                     # Converts and deletes data.csv

            # Keep original CSV
            csv_to_xlsx('report.csv', delete_csv=False) # Keeps report.csv
    """

    # Auto-detect CSV if not provided
    if csv_file is None:
        csv_files = list(Path('.').glob('*.csv'))

        if len(csv_files) == 0:
            print("Error: No CSV files found in current directory")
            return None
        elif len(csv_files) == 1:
            csv_file = str(csv_files[0])
            print(f"Auto-detected CSV file: {csv_file}")
        else:
            print(f"Error: Multiple CSV files found ({len(csv_files)}). Please specify which one to convert:")
            for f in csv_files:
                print(f"  - {f.name}")
            return None

    new_wb = Workbook()
    new_ws = new_wb.active

    try:
        with open(csv_file, 'r', encoding='utf-8', newline='') as f:
            reader = csv.reader(f)
            for row in reader:
                new_ws.append(row)

        xlsx_file = Path(csv_file).with_suffix(".xlsx")
        new_wb.save(xlsx_file)
        print(f"Converted: {csv_file} to {xlsx_file.name}")

        if delete_csv:
            os.remove(csv_file)
            print(f"Deleted: {csv_file}")

        return str(xlsx_file)

    except Exception as e:
        print(f"Error converting CSV to XLSX: {e}")
        return None


def date():
    """
    Get current day of month.

    Returns:
        int: Current day (1-31)

    Example:
        date()  # 24
    """
    return time.localtime().tm_mday


def day():
    """
    Get current day of week.

    Returns:
        str: Day name in lowercase (monday, tuesday, wednesday, thursday, friday, saturday, sunday)

    Example:
        ::

            # Weekday check
            if day() == 'monday':
                print("It is Monday today.")
    """
    days = ['monday', 'tuesday', 'wednesday', 'thursday', 'friday', 'saturday', 'sunday']
    return days[time.localtime().tm_wday]


def drag(*args):
    """
    Drag from source to target.

    Modes:
        1. PyAutoGUI screen drag: (x1, y1, x2, y2)
        2. Selenium element drag: (driver, src_type, src_selector, tgt_type, tgt_selector)

    Args:
        PyAutoGUI: (x1, y1, x2, y2)
        Selenium: (driver, src_type, src_selector, tgt_type, tgt_selector)

    Returns:
        True if successful, False otherwise

    Output:
        - Prints drag confirmation showing source and target coordinates (PyAutoGUI).
        - Prints drag confirmation showing source and target selectors (Selenium).

    Example:
        ::

            # Screen drag (PyAutoGUI) - 2 second duration
            drag(100, 200, 500, 600)

            # Web element drag (Selenium)
            drag(driver, 'id', 'card-1', 'class', 'done-column')
            drag(driver, 'xpath', '//li[1]', 'xpath', '//li[5]')

            # Multiple drivers
            driver1 = browser('https://trello.com')
            driver2 = browser('https://jira.com')
            drag(driver1, 'id', 'task-1', 'id', 'done-column')
            drag(driver2, 'class', 'issue', 'class', 'backlog')
    """

    # ------------------------------------------------------------
    # MODE 1 : PYAUTOGUI SCREEN DRAG (4 integer arguments)
    # needs: display
    # ------------------------------------------------------------
    if len(args) == 4 and all(isinstance(arg, int) for arg in args):
        if not _GUI_AVAILABLE:
            print("Error: drag() requires a display.")
            return False
        x1, y1, x2, y2 = args
        try:
            pyautogui.moveTo(x1, y1, 1, pyautogui.easeInOutQuad)
            pyautogui.dragTo(x2, y2, duration=2, button='left')
            print(f"Dragged from ({x1}, {y1}) to ({x2}, {y2})")
            return True
        except Exception as e:
            print(f"Error during drag: {e}")
            return False

    # ------------------------------------------------------------
    # MODE 2 : SELENIUM ELEMENT DRAG (5 arguments, first is WebDriver)
    # needs: nothing : works on all platforms
    # ------------------------------------------------------------
    elif len(args) == 5 and hasattr(args[0], 'find_element'):
        driver_obj = args[0]
        src_type = args[1]
        src_selector = args[2]
        tgt_type = args[3]
        tgt_selector = args[4]

        valid_selectors = ['id', 'xpath', 'class', 'name', 'css', 'tag', 'text', 'partial']
        if src_type not in valid_selectors:
            raise ValueError(f"Invalid source selector type '{src_type}'. Valid: {', '.join(valid_selectors)}")
        if tgt_type not in valid_selectors:
            raise ValueError(f"Invalid target selector type '{tgt_type}'. Valid: {', '.join(valid_selectors)}")

        try:
            source = _get_web_element(driver_obj, src_type, src_selector)
            target = _get_web_element(driver_obj, tgt_type, tgt_selector)

            if not source:
                print(f"Source element not found: {src_type} = '{src_selector}'")
                return False
            if not target:
                print(f"Target element not found: {tgt_type} = '{tgt_selector}'")
                return False

            ActionChains(driver_obj).drag_and_drop(source, target).perform()
            print(f"Dragged element from {src_type}='{src_selector}' to {tgt_type}='{tgt_selector}'")
            return True

        except Exception as e:
            print(f"Error during Selenium drag: {e}")
            return False

    # ---------------Invalid Arguments---------------
    else:
        raise ValueError(
            "Invalid arguments. Use drag(x1, y1, x2, y2) or drag(driver, src_type, src_selector, tgt_type, tgt_selector)")


def dropdown_select(driver_obj, selector_type, selector, selection_criteria):
    """
    Selects an item from a dropdown menu based on the provided criteria.

    Args:
        driver_obj: Selenium WebDriver instance
        selector_type: Type of selector ('id', 'name', 'xpath', 'class', 'css', 'tag', 'text', 'partial')
        selector: The value of the selector
        selection_criteria: Index (int) or visible text (str) for selection

    Returns:
        True if successful, False otherwise

    Output:
        - Prints confirmation showing the selected option index or text.

    Example:
        ::

            # Select by index
            dropdown_select(driver, 'id', 'country-dropdown', 0)        # Select first option
            dropdown_select(driver, 'id', 'country-dropdown', 2)        # Select third option

            # Select by visible text
            dropdown_select(driver, 'id', 'country-dropdown', 'United States')
            dropdown_select(driver, 'name', 'language', 'English')
            dropdown_select(driver, 'xpath', '//select[@name="city"]', 'New York')

            # Different selector types
            dropdown_select(driver, 'class', 'form-select', 'Option 1')
            dropdown_select(driver, 'css', 'select.dropdown', 'Value')
    """

    try:
        # Using the get_web_element function to locate the dropdown element
        dropdown_element = _get_web_element(driver_obj, selector_type, selector)

        if dropdown_element is None:
            print(f"Dropdown element not found: {selector_type} - {selector}")
            return False

        select = Select(dropdown_element)

        # Selecting based on the type of 'selection_criteria'
        if isinstance(selection_criteria, int):
            select.select_by_index(selection_criteria)
            print(f"Selected dropdown option by index: {selection_criteria}")
        elif isinstance(selection_criteria, str):
            select.select_by_visible_text(selection_criteria)
            print(f"Selected dropdown option by text: '{selection_criteria}'")
        else:
            print("Invalid selection criteria. Must be an integer or string.")
            return False

        return True

    except NoSuchElementException:
        print(f"Option not found in the dropdown: {selection_criteria}")
        return False
    except Exception as e:
        print(f"Error during dropdown selection: {e}")
        return False


def erase(*args):
    """
    Erase/clear text from input fields.

    Modes:
        1. PyAutoGUI active window: ()
        2. Selenium specific element: (driver, selector_type, selector)

    Args:
        *args: Variable arguments depending on mode

    Returns:
        True if successful, False otherwise

    Example:
        ::

            # PyAutoGUI mode (erase active window)
            erase()                                  # Select all and delete (Ctrl+A, Delete)

            # Selenium mode (erase specific element)
            erase(driver, 'id', 'username')          # Clear username field
            erase(driver, 'xpath', '//input[@name="email"]')  # Clear email field
            erase(driver, 'class', 'search-box')     # Clear search box
    """

    try:
        # ------------------------------------------------------------
        # MODE 1 : PYAUTOGUI ACTIVE WINDOW (no arguments)
        # needs: display
        # ------------------------------------------------------------
        if len(args) == 0:
            if not _GUI_AVAILABLE:
                print("Error: erase() requires a display.")
                return False
            pyautogui.hotkey('ctrl', 'a', 'delete')
            print("Erased content in active window")
            return True

        # ------------------------------------------------------------
        # MODE 2 : SELENIUM ELEMENT (driver + selector)
        # needs: nothing : works on all platforms
        # ------------------------------------------------------------
        elif len(args) == 3 and hasattr(args[0], 'find_element'):
            driver_obj = args[0]
            selector_type = args[1]
            selector = args[2]

            element = _get_web_element(driver_obj, selector_type, selector)
            if element:
                element.clear()
                print(f"Cleared element: {selector_type} - {selector}")
                return True
            else:
                print(f"Element not found: {selector_type} - {selector}")
                return False

        # ---------------Invalid Arguments---------------
        else:
            print("Error: Invalid arguments for erase()")
            print("Use: erase() or erase(driver, selector_type, selector)")
            return False

    except Exception as e:
        print(f"Error erasing content: {e}")
        return False


def find_browser(*args):
    """
    Find text in browser using Ctrl+F (find function).

    Args:
        *args: Variable arguments depending on mode

    Returns:
        True if successful, False otherwise

    Output:
        - Prints confirmation with the searched text (PyAutoGUI).
        - Prints confirmation if text was found and highlighted (Selenium).
        - Prints a message if text was not found on the page (Selenium).

    Example:
        ::

            # PyAutoGUI mode (any window)
            find_browser('Python')              # Find in active window
            find_browser('error message')       # Find phrase

            # Selenium mode (browser)
            find_browser(driver, 'Python')      # Find in Selenium browser
            find_browser(driver, 'contact us')  # Find phrase in browser

    Note:
        - PyAutoGUI: Opens find dialog (Ctrl+F), types search term, presses Enter, then Esc.
        - Selenium: Uses JavaScript to highlight matching text on the page in yellow
          and scrolls to the first match. Removes any previous highlights before applying new ones.
        - Default wait time between actions is 1 second (PyAutoGUI only).
    """

    wait_time = 1  # Default wait time (1 second)

    try:
        # ------------------------------------------------------------
        # MODE 1 : SELENIUM (driver object passed)
        # needs: nothing : works on all platforms
        # ------------------------------------------------------------
        if len(args) >= 2 and hasattr(args[0], 'execute_script'):
            driver_obj = args[0]
            search_text = args[1]

            # Use JavaScript to find and highlight text
            # Using raw string (r"...") to avoid escape sequence warnings
            script = r"""
            var searchText = arguments[0];
            var body = document.body;
            var innerHTML = body.innerHTML;

            // Remove previous highlights
            innerHTML = innerHTML.replace(/<mark style="background-color: yellow;">([^<]*)<\/mark>/gi, "$1");

            // Add new highlights
            var regex = new RegExp(searchText, "gi");
            innerHTML = innerHTML.replace(regex, '<mark style="background-color: yellow;">$&</mark>');

            body.innerHTML = innerHTML;

            // Scroll to first match
            var firstMatch = document.querySelector('mark');
            if (firstMatch) {
                firstMatch.scrollIntoView({ behavior: 'smooth', block: 'center' });
                return true;
            }
            return false;
            """

            # Pass search_text as argument to avoid string interpolation issues
            result = driver_obj.execute_script(script, search_text)

            if result:
                print(f"[Selenium] Text '{search_text}' found and highlighted")
                return True
            else:
                print(f"[Selenium] Text '{search_text}' not found on page")
                return False

        # ------------------------------------------------------------
        # MODE 2 : PYAUTOGUI (no driver object)
        # needs: display
        # ------------------------------------------------------------
        elif len(args) == 1:
            if not _GUI_AVAILABLE:
                print("Error: find_browser() requires a display.")
                return False
            search_text = args[0]

            # Open browser find dialog
            pyautogui.hotkey('ctrl', 'f')
            time.sleep(wait_time)

            # Type search text
            pyautogui.typewrite(search_text)
            time.sleep(wait_time)

            # Press Enter to find
            pyautogui.press('enter')
            time.sleep(wait_time)

            # Close find dialog
            pyautogui.press('esc')
            time.sleep(wait_time)

            print(f"[PyAutoGUI] Searched for '{search_text}'")
            return True

        # ---------------Invalid Arguments---------------
        else:
            print("Error: Invalid arguments for find_browser()")
            print("Use: find_browser(text) or find_browser(driver, text)")
            return False

    except Exception as e:
        print(f"Error during browser find: {e}")
        return False


def find_key(data, key):
    """
    Recursively finds all values of a specified key in nested data structures
    (dictionaries, lists and tuples). Particularly useful for searching
    deeply nested JSON data from API responses or parsed files.

    Args:
        data: Data structure to search (dict, list or tuple)
        key: Key name to find

    Returns:
        list: All values found for the key (empty list if not found)

    Example:
        ::

            # Single occurrence
            data = {'name': 'John', 'age': 30}
            name = find_key(data, 'name')[0]          # 'John'

            # Multiple occurrences
            data = {
                'user': {'id': 1, 'name': 'Alice'},
                'admin': {'id': 2, 'name': 'Bob'}
            }
            ids = find_key(data, 'id')                # [1, 2]

            # Nested lists/tuples
            data = {'users': [{'age': 25}, {'age': 30}]}
            ages = find_key(data, 'age')              # [25, 30]

            # API response workflow
            response = requests.get('https://api.example.com/users').json()
            ids = find_key(response, 'id')            # finds all 'id' values

            # Parsed file workflow
            data = json.loads(read('data.json'))
            hosts = find_key(data, 'host')            # finds all 'host' values
    """

    results = []

    def extract(obj):
        """Recursively search for key in nested structure"""
        if isinstance(obj, dict):
            for k, v in obj.items():
                if k == key:
                    results.append(v)
                if isinstance(v, (dict, list, tuple)):
                    extract(v)
        elif isinstance(obj, (list, tuple)):
            for item in obj:
                extract(item)

    extract(data)
    return results


def find_str(string, starts_after, ends_before, index=0):
    """
    Extracts substring between two markers.

    Args:
        string: Text to search in
        starts_after: Start extraction after this
        ends_before: End extraction before this
        index: Which match (0=first, -1=last, 1=second, etc.)

    Returns:
        str or None: Extracted string or None if not found

    Example:
        ::

            # Extract version number from string
            text = 'Version: 1.0.5 released'
            version = find_str(text, 'Version: ', ' released')
            # version = '1.0.5'

            # Extract last occurrence using index=-1
            text = 'User: Alice logged in. User: Bob logged in'
            last_user = find_str(text, 'User: ', ' logged', -1)
            # last_user = 'Bob'
    """

    try:
        # Escape special regex characters
        starts_after = re.escape(starts_after)
        ends_before = re.escape(ends_before)

        # Build regex pattern
        pattern = f'{starts_after}(.*?){ends_before}'
        matches = re.findall(pattern, string)

        if not matches:
            return None

        try:
            return matches[index].strip()
        except IndexError:
            return None

    except re.error as e:
        print(f"Regex error in find_str: {e}")
        return None
    except Exception as e:
        print(f"Error in find_str: {e}")
        return None


def hour():
    """
    Get current hour.

    Returns:
        int: Current hour (0-23) in 24-hour format

    Example:
        hour()  # 14 (2 PM)
    """
    return time.localtime().tm_hour


def inspect():
    """
    Opens GUI to inspect pixel position and color with zoomed preview.

    Usage:
        1. Click on the Pixel Inspector window to bring it into focus.
        2. Move the mouse to the desired pixel.
        3. Press 'ESC' to capture.

    Output:
        - Prints position and RGB/HEX values to console.
        - Copies 'x, y, r, g, b' to clipboard.

    """

    # ------------------------------------------------------------
    # GUARD : needs display, tkinter, pyautogui and clipboard
    # ------------------------------------------------------------
    if not _GUI_AVAILABLE:
        print("Error: inspect() requires a display.")
        return
    if not _CLIPBOARD_AVAILABLE:
        print("Error: inspect() requires a clipboard manager.")
        return

    # Track whether the inspector window is still open
    window_open = True

    def capture_pixel(event=None):
        """Called when ESC is pressed. Captures current mouse position and color."""
        nonlocal window_open
        try:
            # Get current mouse position
            x, y = pyautogui.position()

            # Get the RGB color of the pixel at current mouse position
            color = pyautogui.pixel(x, y)

            # Convert RGB to HEX format
            hex_color = f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}".upper()

            # Format the full output string for console
            formatted_output = f"Pixel at ({x}, {y}): RGB {color} | HEX {hex_color}"

            # Format the clipboard string as 'x, y, r, g, b'
            clipboard_text = f"{x}, {y}, {color[0]}, {color[1]}, {color[2]}"

            # Print to console
            print(formatted_output)
            print(f"Copied to clipboard: {clipboard_text}")

            # Copy to clipboard
            pyperclip.copy(clipboard_text)

            # Mark window as closed and destroy the tkinter window
            window_open = False
            root.destroy()

        except Exception as e:
            print(f"Error capturing pixel: {e}")

    # noinspection PyTypeChecker
    def update_color_and_position():
        """Called repeatedly every 100ms to update the live preview."""

        # Stop updating if the window has been closed
        if not window_open:
            return

        try:
            # Get current mouse position
            x, y = pyautogui.position()

            # Get screen dimensions to validate mouse position
            screen_width, screen_height = pyautogui.size()

            # Skip update if mouse is outside screen bounds
            if x < 0 or y < 0 or x >= screen_width or y >= screen_height:
                root.after(30, update_color_and_position)
                return

            # Get RGB color of the pixel at current mouse position
            color = pyautogui.pixel(x, y)

            # Convert RGB to HEX format
            hex_color = f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}".upper()

            # Update the position, RGB and HEX labels in the GUI
            position_label.config(text=f"Position: ({x}, {y})")
            color_label.config(text=f"RGB: {color}")
            hex_label.config(text=f"HEX: {hex_color}")

            # Update the color swatch to show the current pixel color
            color_display.config(bg=hex_color.lower())

            # --- Zoomed Preview ---

            # Define the capture area size around the cursor (30x30 pixels)
            zoom_size = 30

            # Calculate top-left corner of the capture region centered on cursor
            left = max(0, x - zoom_size // 2)
            top = max(0, y - zoom_size // 2)

            # Clamp the region so it does not go beyond screen edges
            left = min(left, screen_width - zoom_size)
            top = min(top, screen_height - zoom_size)

            # Take a screenshot of the small region around the cursor
            screenshot = pyautogui.screenshot(region=(left, top, zoom_size, zoom_size))

            # Resize the region to 240x240 for 8x magnification using nearest
            # neighbor to keep pixels sharp and not blurred
            zoomed = screenshot.resize((240, 240), Image.Resampling.NEAREST)

            # Convert the zoomed image to a format tkinter can display
            photo = ImageTk.PhotoImage(zoomed)

            # Clear the canvas and draw the new zoomed image
            zoom_canvas.delete("all")
            zoom_canvas.create_image(120, 120, image=photo)

            # Keep a reference to prevent garbage collection
            zoom_canvas.image = photo

            # Draw a red crosshair at the center of the zoomed preview
            zoom_canvas.create_line(120, 110, 120, 130, fill='red', width=2)  # vertical
            zoom_canvas.create_line(110, 120, 130, 120, fill='red', width=2)  # horizontal

        except Exception as e:
            print(f"Pixel preview update failed: {e}")
            return

        # Schedule the next update after 100ms
        root.after(30, update_color_and_position)

    # --- Build the tkinter GUI ---

    # Create the main window
    root = tk.Tk()
    root.title("Pixel Inspector")
    root.geometry("400x560")
    root.configure(bg='#4D4D4D')

    # Bind ESC key to capture_pixel so pressing ESC triggers the capture
    root.bind('<Escape>', capture_pixel)

    # Label to show current mouse position
    position_label = tk.Label(root, text="Position:", font=("Helvetica", 12), bg='#4D4D4D', fg='white')
    position_label.pack()

    # Label to show current RGB values
    color_label = tk.Label(root, text="RGB:", font=("Helvetica", 12), bg='#4D4D4D', fg='white')
    color_label.pack()

    # Label to show current HEX value
    hex_label = tk.Label(root, text="HEX:", font=("Helvetica", 12), bg='#4D4D4D', fg='white')
    hex_label.pack()

    # Color swatch frame that fills with the current pixel color
    color_display = tk.Frame(root, height=60, width=200)
    color_display.pack(pady=10)

    # Label above the zoomed preview canvas
    zoom_label = tk.Label(root, text="Zoomed Preview (8x):", font=("Helvetica", 10), bg='#4D4D4D', fg='white')
    zoom_label.pack()

    # Canvas where the zoomed pixel preview is drawn
    zoom_canvas = tk.Canvas(root, width=240, height=240, bg='#2D2D2D', highlightthickness=2,
                            highlightbackground='white')
    zoom_canvas.pack(pady=5)

    # Instruction label at the bottom
    instruction_label = tk.Label(root, text="Press 'ESC' to capture and exit", font=("Helvetica", 10),
                                 bg='#4D4D4D', fg='white')
    instruction_label.pack(pady=10)

    # Start the live update loop
    update_color_and_position()

    # Start the tkinter event loop : blocks here until window is closed
    root.mainloop()


def log_setup(title):
    """
    Sets up logging and terminal styling for the script.

    Combines terminal setup with comprehensive logging and automatic color-coded status indication.
    Creates a logs folder and saves all output with timestamps.
    Shows output in terminal while also saving to file.

    Args:
        title: Name for both terminal title and log file

    Example:
        ::

            log_setup("MyScript")
            print("This gets logged")
            # ... script runs ...
            # Terminal turns GREEN on success or RED on crash

    Log file format:
        ::

            logs/log_MyScript_2026-03-26_14-30-45_IST_session-1.txt        (active - newest logs)
            logs/log_MyScript_2026-03-26_14-30-45_IST_session-1_part_1.txt (2nd newest - rotated)
            logs/log_MyScript_2026-03-26_14-30-45_IST_session-1_part_2.txt (3rd newest)
            ...
            logs/log_MyScript_2026-03-26_14-30-45_IST_session-1_part_9.txt (oldest backup)

    Session numbering:
        ::

            session-1 : First run of this script
            session-2 : Second run of this script
            session-3 : Third run, etc.
            session-N : Automatically increments based on existing log files

    Features:
        - Sets terminal title and colors (blue bg, white text)
        - Automatic color changes: Blue to Green (success) or Blue to Red (crash)
        - Automatic session numbering (increments from previous runs)
        - Captures all print() statements
        - Captures all errors and exceptions
        - Adds timestamp to each entry
        - Shows output in terminal AND saves to file
        - Automatic file rotation (10MB per file, max 10 files = 100MB per session)
        - Automatic cleanup (keeps max 100MB total logs across all sessions)

    Note:
        Terminal colors change automatically based on script outcome::

            Blue background  : Script is running
            Green background : Script completed successfully
            Red background   : Script crashed (unhandled exception)
    """

    global _log_file_handler, _original_stdout, _original_stderr, _log_folder

    # ---------- TERMINAL SETUP ----------
    system = platform.system()

    if system == "Windows":
        # Clear the screen
        os.system('cls')

        # Set the window title (works in both CMD and PowerShell)
        os.system(f'title {title}')

        # Try CMD color command first (Blue background, white text)
        result = os.system('color 1F')

        # If color command failed (PowerShell), use ANSI escape codes
        if result != 0:
            # Enable ANSI color support in PowerShell
            print("\033[97m\033[44m", end="", flush=True)

    elif system == "Linux":
        # Set the window title using ANSI escape codes
        print(f"\33]0;{title}\a", end="", flush=True)
        # Set terminal text color to bright white and background to blue
        # ANSI: 97 = bright white text, 44 = blue background
        print("\033[97m\033[44m", end="")
        # Clear the screen
        os.system('clear')

    elif system == "Darwin":
        print("Terminal colors not supported on macOS (logging will still work)")

    # ---------- LOGGING SETUP ----------

    # Create logs folder if it doesn't exist
    _log_folder.mkdir(exist_ok=True)

    # ---------- DETERMINE SESSION NUMBER ----------

    # Find all existing log files for this script title
    existing_logs = list(_log_folder.glob(f"log_{title}_*_session-*.txt"))

    # Extract session numbers from existing log files
    session_numbers = []
    for log_file in existing_logs:
        # Pattern: log_title_timestamp_timezone_session-N.txt
        filename = log_file.stem  # Get filename without extension

        # Extract session number using regex
        match = re.search(r'_session-(\d+)', filename)
        if match:
            session_numbers.append(int(match.group(1)))

    # Determine next session number
    if session_numbers:
        next_session = max(session_numbers) + 1
    else:
        next_session = 1

    print(f"Starting session {next_session} for '{title}'")

    # ---------- GET TIMESTAMP WITH TIMEZONE ----------

    # Get current time with system's local timezone
    try:
        now = datetime.datetime.now().astimezone()  # auto-detects system timezone
        time_zone_name = now.tzname() or "LOCAL"  # e.g., IST, EST, PDT
    except Exception:
        now = datetime.datetime.now()
        time_zone_name = "LOCAL"

    # Create log filename with session number
    timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
    log_filename = f"log_{title}_{timestamp}_{time_zone_name}_session-{next_session}.txt"
    log_filepath = _log_folder / log_filename

    # ---------- CONFIGURE LOGGER ----------

    # Configure logger
    logger = logging.getLogger(title)
    logger.setLevel(logging.DEBUG)

    # Remove existing handlers
    logger.handlers.clear()

    # Create file handler with rotation (10MB per file)
    _log_file_handler = _CustomRotatingFileNameHandler(
        log_filepath,
        maxBytes=10 * 1024 * 1024,  # 10MB
        backupCount=9,  # 10 files total (1 active + 9 backups) = 100MB
        encoding='utf-8'  # fixes UnicodeEncodeError for special characters
    )

    # Create formatter with timestamp
    formatter = logging.Formatter(
        '%(asctime)s -> %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    _log_file_handler.setFormatter(formatter)
    logger.addHandler(_log_file_handler)

    # Silence the root logger for the entire session so third-party libraries
    # that add StreamHandlers to it cannot produce duplicate "INFO:..." lines in our logs.
    # Setting level above CRITICAL blocks all messages at the source,
    # regardless of how many handlers any library adds to the root logger.
    # Our named logger is unaffected as it has its own independent level.
    logging.getLogger().handlers.clear()
    logging.getLogger().addHandler(logging.NullHandler())
    logging.getLogger().setLevel(logging.CRITICAL + 1)

    # Redirect stdout and stderr
    _original_stdout = sys.stdout
    _original_stderr = sys.stderr

    sys.stdout = _LogCapture(_original_stdout, logger, logging.INFO)
    sys.stderr = _LogCapture(_original_stderr, logger, logging.ERROR)

    # Log start message
    logger.info(f"Logging started for: {title} (Session {next_session})")
    logger.info(f"Log file: {log_filepath}")

    # ---------- SETUP SUCCESS/ERROR COLOR HANDLERS ----------

    global _script_had_error  # Declare at start of this section

    # Store original exception handler
    original_excepthook = sys.excepthook

    def set_success_color():
        """Change terminal to green on successful completion"""
        global _original_stdout, _original_stderr, _script_had_error, _log_file_handler

        # Don't show success if there was an error
        if _script_had_error:
            return

        # LOG SUCCESS MESSAGE directly to file (before restoring streams)
        if _log_file_handler:
            try:
                # Create a log record manually
                record = logging.LogRecord(
                    name="completion",
                    level=logging.INFO,
                    pathname="",
                    lineno=0,
                    msg="Script execution completed successfully",
                    args=(),
                    exc_info=None
                )
                _log_file_handler.emit(record)
                _log_file_handler.flush()  # Ensure it's written
            except Exception as e:
                print(f"Could not write script completion status to log file: {e}")

        # NOW restore original stdout/stderr
        if _original_stdout:
            sys.stdout = _original_stdout
            sys.stderr = _original_stderr

        # Change color (terminal only)
        if platform.system() == "Windows":
            os.system('color 2F')  # Green background, white text
        elif platform.system() == "Linux":
            print("\033[97m\033[42m", end="", flush=True)  # Green bg, white text

    def set_error_color(exc_type, exc_value, exc_traceback):
        """Change terminal to red on unhandled exception"""
        global _original_stdout, _original_stderr, _script_had_error

        # Mark that an error occurred
        _script_had_error = True

        # LOG THE EXCEPTION FIRST (while stderr is still redirected to log)
        error_msg = ''.join(traceback.format_exception(exc_type, exc_value, exc_traceback))
        print(error_msg, file=sys.stderr, end='')

        # NOW restore original stdout/stderr
        if _original_stdout:
            sys.stdout = _original_stdout
            sys.stderr = _original_stderr

        # Change color
        if platform.system() == "Windows":
            os.system('color 4F')  # Red background, white text
        elif platform.system() == "Linux":
            print("\033[97m\033[41m", end="", flush=True)  # Red bg, white text

        # DON'T call original_excepthook - we already printed the error!
        # (Calling it would print the error a second time)

    # Register success handler (runs on normal exit)
    atexit.register(set_success_color)

    # Register error handler (runs on unhandled exception)
    sys.excepthook = set_error_color


def minute():
    """
    Get current minute.

    Returns:
        int: Current minute (0-59)

    Example:
        minute()  # 23
    """
    return time.localtime().tm_min


def month():
    """
    Get current month.

    Returns:
        int: Current month (1-12)

    Example:
        month()  # 2 (February)
    """
    return time.localtime().tm_mon


def press(*keys):
    """
    Press keyboard keys with support for Selenium, PyAutoGUI and key combinations.

    Modes:
        1. PyAutoGUI single key: (key)
        2. PyAutoGUI key N times: (key, count)
        3. PyAutoGUI key combination: (key1, key2, ...)
        4. Selenium driver key: (driver, key)
        5. Selenium driver key N times: (driver, key, count)
        6. Selenium driver key combination: (driver, key1, key2, ...)
        7. Selenium element key: (driver, selector_type, selector, key)

    Args:
        *keys: Variable arguments for key presses

    Returns:
        True if successful, False otherwise

    Example:
        ::

            # PyAutoGUI keys (no driver needed)
            press("tab")                    # Single key
            press("tab", 5)                 # Press 5 times
            press("tab", -5)                # Press 5 times with shift held
            press("ctrl", "a")              # Two-key combo
            press("alt", "ctrl", "z")       # Three-key combo
            press("num5")                   # Numpad 5
            press("volumeup")               # Volume up
            press("mute")                   # Volume mute (short form)
            press("back")                   # Browser back (short form)
            press("forward")                # Browser forward (short form)

            # Selenium driver keys (pass driver object)
            press(driver, "tab")
            press(driver, "tab", 5)           # Press tab 5 times
            press(driver, "tab", -5)          # Press shift+tab 5 times
            press(driver, "ctrl", "c")
            press(driver, "ctrl", "shift", "s")

            # Selenium element + key (driver, selector_type, selector, key)
            press(driver, "xpath", "//input", "enter")
            press(driver, "id", "username", "tab")

    Note:
        - Negative count presses the key with Shift held (e.g. Shift+Tab for reverse navigation).
        - PyAutoGUI-only keys (num0-9, volumeup, volumedown, mute, back, forward, etc.)
          are not supported in Selenium mode.
        - Short forms supported: 'back' for browserback, 'forward' for browserforward,
          'mute' for volumemute.
    """

    # Key mappings for special keys
    key_map = {
        'enter': Keys.RETURN,
        'tab': Keys.TAB,
        'up': Keys.ARROW_UP,
        'down': Keys.ARROW_DOWN,
        'left': Keys.ARROW_LEFT,
        'right': Keys.ARROW_RIGHT,
        'alt': Keys.ALT,
        'backspace': Keys.BACKSPACE,
        'esc': Keys.ESCAPE,
        'escape': Keys.ESCAPE,
        'home': Keys.HOME,
        'end': Keys.END,
        'insert': Keys.INSERT,
        'delete': Keys.DELETE,
        'pageup': Keys.PAGE_UP,
        'pagedown': Keys.PAGE_DOWN,
        'shift': Keys.SHIFT,
        'ctrl': Keys.CONTROL,
        'control': Keys.CONTROL,
        'space': Keys.SPACE,
        'pause': Keys.PAUSE,
        'f1': Keys.F1, 'f2': Keys.F2, 'f3': Keys.F3, 'f4': Keys.F4,
        'f5': Keys.F5, 'f6': Keys.F6, 'f7': Keys.F7, 'f8': Keys.F8,
        'f9': Keys.F9, 'f10': Keys.F10, 'f11': Keys.F11, 'f12': Keys.F12,
        'windows': Keys.COMMAND,
        'command': Keys.COMMAND
    }

    # PyAutoGUI-only keys (not supported in Selenium)
    pyautogui_only_keys = {
        'num0', 'num1', 'num2', 'num3', 'num4', 'num5', 'num6', 'num7', 'num8', 'num9',
        'numlock', 'volumeup', 'volumedown', 'volumemute', 'playpause',
        'browserback', 'browserforward', 'printscreen', 'capslock', 'scrolllock',
        # Short forms
        'back', 'forward', 'mute'
    }

    # Short form mappings for PyAutoGUI keys
    pyautogui_short_forms = {
        'back': 'browserback',
        'forward': 'browserforward',
        'mute': 'volumemute'
    }

    # Valid PyAutoGUI keys (comprehensive list) to check valid keys
    valid_pyautogui_keys = {
        '\t', '\n', '\r', ' ', '!', '"', '#', '$', '%', '&', "'", '(',
        ')', '*', '+', ',', '-', '.', '/', '0', '1', '2', '3', '4', '5', '6', '7',
        '8', '9', ':', ';', '<', '=', '>', '?', '@', '[', '\\', ']', '^', '_', '`',
        'a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n', 'o',
        'p', 'q', 'r', 's', 't', 'u', 'v', 'w', 'x', 'y', 'z', '{', '|', '}', '~',
        'accept', 'add', 'alt', 'altleft', 'altright', 'apps', 'backspace',
        'browserback', 'browserfavorites', 'browserforward', 'browserhome',
        'browserrefresh', 'browsersearch', 'browserstop', 'capslock', 'clear',
        'convert', 'ctrl', 'ctrlleft', 'ctrlright', 'decimal', 'del', 'delete',
        'divide', 'down', 'end', 'enter', 'esc', 'escape', 'execute', 'f1', 'f10',
        'f11', 'f12', 'f13', 'f14', 'f15', 'f16', 'f17', 'f18', 'f19', 'f2', 'f20',
        'f21', 'f22', 'f23', 'f24', 'f3', 'f4', 'f5', 'f6', 'f7', 'f8', 'f9',
        'final', 'fn', 'hanguel', 'hangul', 'hanja', 'help', 'home', 'insert', 'junja',
        'kana', 'kanji', 'launchapp1', 'launchapp2', 'launchmail',
        'launchmediaselect', 'left', 'modechange', 'multiply', 'nexttrack',
        'nonconvert', 'num0', 'num1', 'num2', 'num3', 'num4', 'num5', 'num6',
        'num7', 'num8', 'num9', 'numlock', 'pagedown', 'pageup', 'pause', 'pgdn',
        'pgup', 'playpause', 'prevtrack', 'print', 'printscreen', 'prntscrn',
        'prtsc', 'prtscr', 'return', 'right', 'scrolllock', 'select', 'separator',
        'shift', 'shiftleft', 'shiftright', 'sleep', 'space', 'stop', 'subtract', 'tab',
        'up', 'volumedown', 'volumemute', 'volumeup', 'win', 'winleft', 'winright', 'yen',
        'command', 'option', 'optionleft', 'optionright',
        # Short forms
        'back', 'forward', 'mute'
    }

    try:
        # ------------------------------------------------------------
        # SELENIUM MODE - Check if first arg is WebDriver object
        # needs: nothing : works on all platforms
        # ------------------------------------------------------------
        if len(keys) > 0 and hasattr(keys[0], 'find_element'):
            driver_obj = keys[0]

            # ----------------------------------------------------------
            # SELENIUM ELEMENT + KEY (driver, selector_type, selector, key)
            # ----------------------------------------------------------
            if len(keys) >= 4 and keys[1] in ('id', 'name', 'xpath', 'text', 'partial', 'tag', 'class', 'css'):
                selector_type = keys[1]
                selector = keys[2]
                key_name = keys[3].lower()

                # Check if it's a PyAutoGUI-only key
                if key_name in pyautogui_only_keys:
                    print(f"Warning: '{key_name}' key not supported in Selenium. Use PyAutoGUI instead.")
                    return False

                key = key_map.get(key_name, key_name)  # Use mapping or raw key

                # Find element and send key
                element = _get_web_element(driver_obj, selector_type, selector)
                if element:
                    element.send_keys(key)
                    return True
                else:
                    print(f"Element not found with {selector_type}: {selector}")
                    return False

            # ----------------------------------------------------------
            # SELENIUM DRIVER KEYS (driver + keys)
            # ----------------------------------------------------------
            else:
                action = ActionChains(driver_obj)

                # Single key on driver: press(driver, "tab")
                if len(keys) == 2:
                    key_name = keys[1].lower()

                    # Check if it's a PyAutoGUI-only key
                    if key_name in pyautogui_only_keys:
                        print(
                            f"Warning: '{key_name}' key not supported in Selenium. Use PyAutoGUI press('{key_name}') instead.")
                        return False

                    key = key_map.get(key_name, keys[1])  # Use mapping or raw string
                    action.send_keys(key).perform()
                    return True

                # Key + count: press(driver, "tab", 5)
                elif len(keys) == 3 and isinstance(keys[2], int):
                    count = keys[2]
                    key_name = keys[1].lower()

                    # Check if it's a PyAutoGUI-only key
                    if key_name in pyautogui_only_keys:
                        print(f"Warning: '{key_name}' key not supported in Selenium.")
                        return False

                    key = key_map.get(key_name, keys[1])

                    # Positive count: press normally
                    if count >= 0:
                        for _ in range(count):
                            action.send_keys(key).perform()
                            time.sleep(1)
                    # Negative count: press with shift held
                    else:
                        for _ in range(abs(count)):
                            action.key_down(Keys.SHIFT).send_keys(key).key_up(Keys.SHIFT).perform()
                            time.sleep(1)

                    return True

                # Multiple keys (combination): press(driver, "ctrl", "c")
                elif len(keys) >= 3:
                    # Check if any key is PyAutoGUI-only
                    for i in range(1, len(keys)):
                        if keys[i].lower() in pyautogui_only_keys:
                            print(f"Warning: '{keys[i]}' key not supported in Selenium.")
                            return False

                    # Press all keys down in order
                    for i in range(1, len(keys)):
                        key_name = keys[i].lower()
                        key = key_map.get(key_name, keys[i])

                        # If it's a modifier key, hold it down
                        if key in (Keys.CONTROL, Keys.SHIFT, Keys.ALT, Keys.COMMAND):
                            action.key_down(key)
                        else:
                            # Regular key - just send it
                            action.send_keys(key)

                    # Release all modifier keys in reverse order
                    for i in range(len(keys) - 1, 0, -1):
                        key_name = keys[i].lower()
                        key = key_map.get(key_name, keys[i])

                        if key in (Keys.CONTROL, Keys.SHIFT, Keys.ALT, Keys.COMMAND):
                            action.key_up(key)

                    action.perform()
                    return True

        # ------------------------------------------------------------
        # PYAUTOGUI MODE - No driver object
        # needs: display
        # ------------------------------------------------------------
        else:
            if not _GUI_AVAILABLE:
                print("Error: press() requires a display.")
                return False

            # Single key press
            if len(keys) == 1:
                key_name = keys[0].lower()
                # Translate short form to full form
                actual_key = pyautogui_short_forms.get(key_name, keys[0])

                # Validate key
                if actual_key.lower() not in valid_pyautogui_keys and len(actual_key) > 1:
                    print(f"Error: Invalid key '{keys[0]}'")
                    return False

                pyautogui.press(actual_key)
                return True

            # Key + count (press n times)
            elif len(keys) == 2 and isinstance(keys[1], int):
                count = keys[1]
                key_name = keys[0].lower()
                # Translate short form to full form
                actual_key = pyautogui_short_forms.get(key_name, keys[0])

                # Validate key
                if actual_key.lower() not in valid_pyautogui_keys and len(actual_key) > 1:
                    print(f"Error: Invalid key '{keys[0]}'")
                    return False

                # Positive count: press normally
                if count >= 0:
                    for _ in range(count):
                        pyautogui.press(actual_key)
                        time.sleep(1)
                # Negative count: press with shift held
                else:
                    pyautogui.keyDown('shift')
                    for _ in range(abs(count)):
                        pyautogui.press(actual_key)
                        time.sleep(1)
                    pyautogui.keyUp('shift')
                return True

            # Multiple key combination (hotkey)
            else:
                # Translate short forms to full forms and validate
                translated_keys = []
                for key in keys:
                    key_name = key.lower() if isinstance(key, str) else key
                    translated_key = pyautogui_short_forms.get(key_name, key)

                    # Validate each key
                    if translated_key.lower() not in valid_pyautogui_keys and len(str(translated_key)) > 1:
                        print(f"Error: Invalid key '{key}' in combination")
                        return False

                    translated_keys.append(translated_key)

                pyautogui.hotkey(*translated_keys)
                return True

    except Exception as e:
        print(f"Error pressing keys {keys}: {e}")
        return False


def read(*args):
    """
    Extract text from screen (using OCR), files (by parsing file format) or a Selenium browser window.

    Modes:
        1. No arguments: OCR full screen
        2. 2 integers: OCR from (x, y) to bottom-right corner
        3. 4 integers: OCR specific region (x, y, width, height)
        4. 1 string: Read file by parsing its format
        5. 1 driver object: Take screenshot of Selenium browser and read using OCR

    Supported file formats:
        - Documents: PDF, DOCX, PPTX, ODT, RTF
        - Tabular: CSV, TSV, XLSX, SQLite
        - Structured: JSON, YAML, XML, INI/CFG
        - Text: TXT, LOG, MD
        - Web: HTML
        - Email: EML, MSG
        - eBooks: EPUB
        - Scripts: SH, BAT, PY

    Args:
        *args: Variable arguments depending on mode

    Returns:
        str: Extracted text or None if error

    Example:
        ::

            # OCR - Read entire screen
            text = read()

            # OCR - Read from (100, 200) to bottom-right corner
            text = read(100, 200)

            # OCR - Read specific region: x=100, y=200, width=400, height=300
            text = read(100, 200, 400, 300)

            # Selenium - Read text from browser window using OCR
            d1 = browser('https://example.com')
            text = read(d1)

            # File - Read with extension
            text = read('report.pdf')
            text = read('data.csv')
            text = read('script.py')

            # File - Read without extension (auto-detects)
            text = read('report')      # Finds report.pdf automatically
            text = read('config')      # Finds config.yaml automatically

            # Check if text on screen
            if 'login' in read():
                print("Login visible!")

    Note:
        OCR first run downloads models (~100MB), subsequent runs are fast.
    """

    # ------------------------------------------------------------
    # DETERMINE MODE AND VALIDATE ARGUMENTS
    # ------------------------------------------------------------
    is_ocr = False
    region = None

    # MODE 1: OCR - No arguments (full screen)
    if len(args) == 0:
        is_ocr = True
        region = None  # Full screen

    # MODE 2: OCR - 2 integers (from x,y to bottom-right corner)
    elif len(args) == 2 and all(isinstance(arg, int) for arg in args):
        is_ocr = True
        x, y = args

        # Validate coordinates
        if x < 0 or y < 0:
            print(f"Error: Coordinates must be non-negative (x={x}, y={y})")
            return None

        # Get screen size
        screen_width, screen_height = pyautogui.size()

        # Calculate width and height to bottom-right corner
        width = screen_width - x
        height = screen_height - y

        if width <= 0 or height <= 0:
            print(f"Error: Starting point ({x}, {y}) is outside screen bounds")
            return None

        region = (x, y, width, height)

    # MODE 3: OCR - 4 integers (screen region: x, y, width, height)
    elif len(args) == 4 and all(isinstance(arg, int) for arg in args):
        is_ocr = True
        x, y, width, height = args

        # Validate region
        if width <= 0 or height <= 0:
            print(f"Error: Invalid region dimensions (width={width}, height={height})")
            return None

        if x < 0 or y < 0:
            print(f"Error: Coordinates must be non-negative (x={x}, y={y})")
            return None

        region = (x, y, width, height)

    # ------------------------------------------------------------
    # PERFORM OCR (modes 1, 2, 3 : needs display)
    # ------------------------------------------------------------
    if is_ocr:
        if not _GUI_AVAILABLE:
            print("Error: read() screen/OCR modes require a display.")
            return None

        try:
            # Get shared OCR reader
            reader = _get_ocr_reader()

            # Capture screenshot (region or full screen)
            if region:
                image = pyautogui.screenshot(region=region)
            else:
                image = pyautogui.screenshot()

            # Convert PIL image to numpy array
            image = np.array(image)
            cleaned = _preprocess_for_ocr(image)

            # Extract text using EasyOCR
            results = reader.readtext(cleaned, detail=0)

            # If no text was found
            if not results:
                print(f"OCR: No text found.")
                return ""  # Return blank string (not None) to avoid errors in blank text case

            # Join all detected text with spaces
            text = ' '.join(results)

            return text.strip()

        except ImportError as e:
            print(f"Error: Required library not installed - {e}")
            print("Install with: pip install easyocr opencv-python")
            return None

        except Exception as e:
            print(f"Error during OCR: {e}")
            return None

    # ------------------------------------------------------------
    # MODE 4: Selenium driver - screenshot browser and OCR
    # needs: nothing for screenshot : OCR runs on bytes, no display needed
    # ------------------------------------------------------------
    elif len(args) == 1 and hasattr(args[0], 'find_element'):
        driver_obj = args[0]

        try:
            reader = _get_ocr_reader()

            # Get screenshot as bytes directly without saving to disk
            png_bytes = driver_obj.get_screenshot_as_png()
            image = np.array(Image.open(io.BytesIO(png_bytes)))
            cleaned = _preprocess_for_ocr(image)

            results = reader.readtext(cleaned, detail=0)

            if not results:
                print("OCR: No text found in browser window.")
                return ""

            text = ' '.join(results)
            return text.strip()

        except Exception as e:
            print(f"Could not read text from browser window: {e}")
            return None

    # ------------------------------------------------------------
    # MODE 5: File reading - 1 string (file path)
    # needs: nothing : safe on all platforms including servers
    # ------------------------------------------------------------
    elif len(args) == 1 and isinstance(args[0], str):
        file = args[0]

        # ------------------------------------------------------------
        # AUTO-DETECT FILE EXTENSION IF NOT PROVIDED
        # ------------------------------------------------------------
        if not os.path.splitext(file)[1]:  # No extension provided
            # Get directory and filename
            if os.path.dirname(file):
                directory = os.path.dirname(file)
                filename = os.path.basename(file)
            else:
                directory = '.'
                filename = file

            # Find files matching the name
            try:
                matching_files = [f for f in os.listdir(directory)
                                  if os.path.splitext(f)[0] == filename]

                if len(matching_files) == 1:
                    # Found exactly one file - use it
                    file = os.path.join(directory, matching_files[0])
                    print(f"Auto-detected file: {file}")
                elif len(matching_files) > 1:
                    print(f"Error: Multiple files found with name '{filename}':")
                    for f in matching_files:
                        print(f"  - {f}")
                    print("Please specify file extension.")
                    return None
                else:
                    print(f"Error: No file found with name '{filename}'")
                    return None

            except Exception as e:
                print(f"Error searching for file: {e}")
                return None

        # Get file extension
        _, ext = os.path.splitext(file)
        ext = ext.lower()

        try:
            # ------------------------------------------------------------
            # DOCUMENTS
            # ------------------------------------------------------------

            # PDF files
            if ext == '.pdf':
                with open(file, 'rb') as pdf_file:
                    pdf_reader = PdfReader(pdf_file)
                    return '\r\n'.join([pdf_reader.pages[i].extract_text() for i in range(len(pdf_reader.pages))])

            # Word documents
            elif ext in ['.docx', '.doc']:
                doc = Document(file)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])

            # PowerPoint presentations
            elif ext in ['.pptx', '.ppt']:
                prs = Presentation(file)
                text = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    text.append(f"=== Slide {slide_num} ===")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                    text.append('')
                return '\n'.join(text)

            # OpenDocument Text (LibreOffice)
            elif ext == '.odt':
                textdoc = odf_load(file)
                allparas = textdoc.getElementsByType(P)
                return '\n'.join([teletype.extractText(p) for p in allparas])

            # Rich Text Format
            elif ext == '.rtf':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return rtf_to_text(f.read())

            # ------------------------------------------------------------
            # TABULAR DATA (with >>>Row_X: format)
            # ------------------------------------------------------------

            # CSV files
            elif ext == '.csv':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = list(reader)

                    if not rows:
                        return ''

                    text = []
                    # First row as header
                    text.append('\t'.join(rows[0]))
                    text.append('-' * 60)

                    # Data rows with >>>Row_X: format
                    for row_num, row in enumerate(rows[1:], 1):
                        row_data = '\t'.join(row)
                        text.append(f">>>Row_{row_num}:\t{row_data}")

                    return '\n'.join(text)

            # TSV files (Tab-separated values)
            elif ext == '.tsv':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f, delimiter='\t')
                    rows = list(reader)

                    if not rows:
                        return ''

                    text = []
                    # First row as header
                    text.append('\t'.join(rows[0]))
                    text.append('-' * 60)

                    # Data rows with >>>Row_X: format
                    for row_num, row in enumerate(rows[1:], 1):
                        row_data = '\t'.join(row)
                        text.append(f">>>Row_{row_num}:\t{row_data}")

                    return '\n'.join(text)

            # Excel files
            elif ext in ['.xlsx', '.xlsm']:
                wb = load_workbook(file, data_only=True)
                text = []

                for sheet in wb.worksheets:
                    text.append(f"\n{'=' * 60}")
                    text.append(f"Sheet: {sheet.title}")
                    text.append('=' * 60)

                    rows = list(sheet.iter_rows(values_only=True))

                    if rows:
                        # First row as header
                        text.append('\t'.join([str(cell) if cell is not None else '' for cell in rows[0]]))
                        text.append('-' * 60)

                        # Data rows with >>>Row_X: format
                        for row_num, row in enumerate(rows[1:], 1):
                            row_data = '\t'.join([str(cell) if cell is not None else 'NULL' for cell in row])
                            text.append(f">>>Row_{row_num}:\t{row_data}")

                    text.append('')

                return '\n'.join(text)

            # SQLite databases
            elif ext in ['.db', '.sqlite', '.sqlite3']:
                conn = sqlite3.connect(file)
                cursor = conn.cursor()
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = cursor.fetchall()
                text = []

                for table in tables:
                    table_name = table[0]
                    text.append(f"\n{'=' * 60}")
                    text.append(f"Table: {table_name}")
                    text.append('=' * 60)

                    # Get column names
                    cursor.execute(f"PRAGMA table_info({table_name})")
                    columns = [col[1] for col in cursor.fetchall()]

                    # Add column headers
                    text.append('\t'.join(columns))
                    text.append('-' * 60)

                    # Get data (limit to 100 rows)
                    cursor.execute(f"SELECT * FROM {table_name} LIMIT 100")
                    rows = cursor.fetchall()

                    # Add rows with >>>Row_X: format
                    for row_num, row in enumerate(rows, 1):
                        row_data = '\t'.join([str(cell) if cell is not None else 'NULL' for cell in row])
                        text.append(f">>>Row_{row_num}:\t{row_data}")

                    text.append(f"\nTotal_rows: {len(rows)}\n")

                conn.close()
                return '\n'.join(text)

            # ------------------------------------------------------------
            # STRUCTURED DATA
            # ------------------------------------------------------------

            # JSON files
            elif ext == '.json':
                with open(file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data, indent=2)

            # YAML files
            elif ext in ['.yaml', '.yml']:
                with open(file, 'r', encoding='utf-8') as f:
                    data = yaml.safe_load(f)
                    return yaml.dump(data, default_flow_style=False, sort_keys=False)

            # XML files
            elif ext == '.xml':
                tree = et.parse(file)
                return et.tostring(tree.getroot(), encoding='unicode', method='xml')

            # INI/CFG files
            elif ext in ['.ini', '.cfg']:
                config = configparser.ConfigParser()
                config.read(file)
                text = []
                for section in config.sections():
                    text.append(f"[{section}]")
                    for key, value in config.items(section):
                        text.append(f"{key} = {value}")
                    text.append('')
                return '\n'.join(text)

            # ------------------------------------------------------------
            # TEXT FILES
            # ------------------------------------------------------------

            # Plain text, log files, markdown
            elif ext in ['.txt', '.log', '.md']:
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # ------------------------------------------------------------
            # WEB
            # ------------------------------------------------------------

            # HTML files
            elif ext in ['.html', '.htm']:
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    soup = BeautifulSoup(f.read(), 'html.parser')
                    return soup.get_text('\n').strip()

            # ------------------------------------------------------------
            # EMAIL
            # ------------------------------------------------------------

            # EML files (standard email)
            elif ext == '.eml':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    msg = email.message_from_file(f)
                    text = []
                    text.append(f"From: {msg.get('From')}")
                    text.append(f"To: {msg.get('To')}")
                    text.append(f"Subject: {msg.get('Subject')}")
                    text.append(f"Date: {msg.get('Date')}")
                    text.append("\n--- Body ---")

                    if msg.is_multipart():
                        for part in msg.walk():
                            if part.get_content_type() == "text/plain":
                                payload = part.get_payload(decode=True)
                                if payload:
                                    text.append(payload.decode('utf-8', errors='ignore'))
                    else:
                        payload = msg.get_payload(decode=True)
                        if payload:
                            text.append(payload.decode('utf-8', errors='ignore'))

                    return '\n'.join(text)

            # MSG files (Outlook email)
            elif ext == '.msg':
                msg = extract_msg.Message(file)
                text = []
                text.append(f">>>From: {msg.sender}")
                text.append(f">>>To: {msg.to}")
                text.append(f">>>Subject: {msg.subject}")
                text.append(f">>>Date: {msg.date}")
                text.append("\nBody ---")
                text.append(msg.body)
                msg.close()
                return '\n'.join(text)

            # ------------------------------------------------------------
            # EBOOKS
            # ------------------------------------------------------------

            # EPUB files
            elif ext == '.epub':
                book = epub.read_epub(file)
                text = []
                for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text.append(soup.get_text())
                return '\n\n'.join(text)

            # ------------------------------------------------------------
            # SCRIPT FILES
            # ------------------------------------------------------------

            # Shell scripts (Linux)
            elif ext == '.sh':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # Batch files (Windows)
            elif ext == '.bat':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # Python scripts
            elif ext == '.py':
                with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # ------------------------------------------------------------
            # UNSUPPORTED FORMAT
            # ------------------------------------------------------------

            else:
                print(f"Unsupported file format: {ext}")
                print(f"Supported formats: PDF, DOCX, PPTX, ODT, RTF, CSV, TSV, XLSX, SQLite,")
                print(f"                   JSON, YAML, XML, INI/CFG, TXT, LOG, MD, HTML,")
                print(f"                   EML, MSG, EPUB, SH, BAT, PY")
                return None

        except Exception as e:
            print(f"Error reading file {file}: {e}")
            return None

    # ---------------Invalid Arguments---------------
    else:
        print("Error: Invalid arguments for read()")
        print("Use: read() - OCR full screen")
        print("     read(x, y) - OCR from (x,y) to bottom-right")
        print("     read(x, y, width, height) - OCR specific region")
        print("     read(driver) - OCR browser window")
        print("     read('file.pdf') - Read file")
        return None


def run(target):
    """
    Runs a file or application on Windows and Linux.

    Args:
        target: File path or application name to execute

            - If target is a file path: Opens with default application
            - If target is an application name: Launches the application
            - For applications, the command must be available in system PATH

    Raises:
        NotImplementedError: If called on macOS

    Output:
        - Prints error message if file or application was not found.
        - Prints error message if permission was denied.

    Example:
        ::

            # Open files with default application
            run("sample.txt")           # Opens in default text editor
            run("document.pdf")         # Opens in default PDF viewer
            run("C:\\Users\\file.xlsx") # Opens Excel file

            # Launch applications (Windows)
            run("calc")                 # Calculator
            run("notepad")              # Notepad
            run("mspaint")              # Paint

            # Launch applications (Linux)
            run("gedit")                # Text editor
            run("firefox")              # Browser
            run("gnome-calculator")     # Calculator

    Note:
        - Windows: Uses os.startfile for files, subprocess for applications.
        - Linux: Uses xdg-open for files, direct execution for applications.
          xdg-utils is required (included in Linux dependencies).
    """

    # Check if macOS and reject
    if sys.platform.startswith('darwin') or platform.system() == "Darwin":
        raise NotImplementedError(
            "run() is not supported on macOS.\n"
            "Supported platforms: Windows, Linux\n"
            "Use subprocess.Popen() directly for macOS-specific needs."
        )

    try:
        # ------------------------------------------------------------
        # WINDOWS
        # ------------------------------------------------------------
        if sys.platform.startswith('win'):
            if os.path.isfile(target):
                # File exists - open with default application
                # Using os.startfile is more reliable than shell=True
                os.startfile(target)
            else:
                # Assume it's an application name
                # Use shell=True to allow PATH resolution
                subprocess.Popen(target, shell=True)

        # ------------------------------------------------------------
        # LINUX
        # ------------------------------------------------------------
        elif sys.platform.startswith('linux'):
            if os.path.isfile(target):
                # File exists - open with default application using xdg-open
                item_path = os.path.abspath(target)
                subprocess.Popen(['xdg-open', item_path])
            else:
                # Assume it's an application name - try to execute directly
                subprocess.Popen([target])

        else:
            # Unknown platform
            raise NotImplementedError(
                f"run() is not supported on platform: {sys.platform}\n"
                "Supported platforms: Windows (win32), Linux (linux)"
            )

    except FileNotFoundError:
        print(f"Error: Application or file not found: '{target}'")
        print("For applications, ensure the command is in your system PATH.")
    except PermissionError:
        print(f"Error: Permission denied to run: '{target}'")
        print("Check file permissions or try running with appropriate privileges.")
    except Exception as e:
        print(f"Error running '{target}': {e}")


def say(text, volume=1.0):
    """
    Speak text using offline Text-to-Speech via Piper TTS.

    Args:
        text:   Text to speak
        volume: Volume level 0.0 to 1.0 (default: 1.0)

    Returns:
        None

    Example:
        ::

            say("Hello, how are you?")
            say("Download complete.")
            say("Error occurred, please try again.", volume=0.7)
            say("Warning: Low battery.", volume=0.5)

    Note:
        - Automatically logs spoken text when log_setup() is active.
        - Requires: pip install piper-tts huggingface_hub
        - Linux requires: sudo apt install espeak-ng alsa-utils
        - Model files are saved in: Windows → %LOCALAPPDATA%\\autocore\\piper_models\\
                                    Linux   → ~/.local/share/autocore/piper_models/
        - Model size is approximately 60MB, downloaded once and reused.
        - Browse all voices at: https://huggingface.co/rhasspy/piper-voices
    """

    # ------------------------------------------------------------
    # GUARD — needs audio system
    # _AUDIO_AVAILABLE is set at file level during library import.
    # if audio init failed (headless server, no sound drivers etc.)
    # we bail out early instead of crashing deep inside piper
    # ------------------------------------------------------------
    if not _AUDIO_AVAILABLE:
        print("Error: say() requires an audio system.")
        return

    # ensure caller passed a string — piper will crash on anything else
    if not isinstance(text, str):
        raise TypeError("Text to speak must be a string")

    # piper's SynthesisConfig will silently clamp or error outside this range
    if not (0.0 <= volume <= 1.0):
        raise ValueError("Volume must be between 0.0 and 1.0")

    try:
        # ------------------------------------------------------------
        # LAZY IMPORTS
        # kept here intentionally and NOT moved to the top of the file
        # because:
        #   - piper and huggingface_hub are optional heavy dependencies
        #     if not installed, the rest of the library still works fine
        #   - winsound is Windows-only, importing at file level would
        #     crash on Linux
        # ------------------------------------------------------------
        from piper.voice import PiperVoice, SynthesisConfig

        # ------------------------------------------------------------
        # MODEL PATHS
        # models are stored in a user-writable location on each platform
        # so we never hit permission errors writing into Program Files
        # or other system-protected directories
        #
        # Windows → C:\Users\YourUsername\AppData\Local\autocore\piper_models\
        # Linux   → /home/yourusername/.local/share/autocore/piper_models/
        #
        # os.environ["LOCALAPPDATA"] on Windows always points to the
        # current user's AppData\Local folder which is always writable
        # without admin rights, unlike C:\Program Files\
        # ------------------------------------------------------------
        model_name = "en_US-libritts_r-medium.onnx"
        model_json = "en_US-libritts_r-medium.onnx.json"

        if platform.system() == "Windows":
            model_dir = Path(os.environ["LOCALAPPDATA"]) / "autocore" / "piper_models"
        else:
            model_dir = Path.home() / ".local" / "share" / "autocore" / "piper_models"

        model_path = model_dir / model_name  # full path to the .onnx model
        model_json_path = model_dir / model_json  # full path to the .json config

        # ------------------------------------------------------------
        # VALIDATION HELPERS
        # defined as inner functions so they can access model_dir
        # from the enclosing scope without passing it as an argument
        # ------------------------------------------------------------
        def _is_valid_onnx(path):
            # valid onnx must exist AND be at least 50MB
            # we know the model is ~60MB so anything under 50MB
            # means the download was interrupted or corrupted
            p = Path(path)
            return p.exists() and p.stat().st_size > 50 * 1024 * 1024  # 50MB in bytes

        def _is_valid_json(path):
            # valid json config must exist, be non-empty, and parse
            # without errors — a corrupt json would crash piper at
            # load time so we catch it here and re-download instead
            p = Path(path)
            if not p.exists() or p.stat().st_size == 0:
                return False
            try:
                with open(p, 'r', encoding='utf-8') as f:
                    json.load(f)  # json already imported at file level
                return True
            except Exception:
                return False

        def _download_file(filename, is_onnx=False):
            # downloads one file from HuggingFace into model_dir
            # hf_hub_download saves into HuggingFace's own cache first
            # (~/.cache/huggingface/) then we move it to piper_models/
            # retries up to 3 times for network hiccups or incomplete downloads
            # returns the final Path if successful, None if all retries failed
            try:
                from huggingface_hub import hf_hub_download  # lazy — optional dependency
            except ImportError:
                print("Error: huggingface_hub not installed. Run: pip install huggingface_hub")
                return None

            max_retries = 3
            for attempt in range(1, max_retries + 1):
                try:
                    # suppress HuggingFace's own progress/warning output
                    # so our print messages stay clean and readable
                    with warnings.catch_warnings():
                        warnings.simplefilter("ignore")
                        cached_path = hf_hub_download(
                            repo_id="rhasspy/piper-voices",
                            filename=f"en/en_US/libritts_r/medium/{filename}",  # path inside the HuggingFace repo
                            local_dir=str(model_dir),
                        )

                    # hf_hub_download may save into a subdirectory inside model_dir
                    # move it to model_dir root so our fixed model_path always finds it
                    final_path = model_dir / filename
                    if os.path.abspath(cached_path) != os.path.abspath(final_path):
                        shutil.move(cached_path, final_path)  # move from hf cache subdir to model_dir root

                    # validate the downloaded file before declaring success
                    # size check for onnx, json parse check for config file
                    valid = _is_valid_onnx(final_path) if is_onnx else _is_valid_json(final_path)
                    if valid:
                        return final_path  # download succeeded and file is healthy

                    # file exists but failed validation — delete it so
                    # the next attempt starts a completely fresh download
                    print(f"Attempt {attempt}/{max_retries}: {filename} appears incomplete, retrying...")
                    if final_path.exists():
                        final_path.unlink()  # delete corrupt/incomplete file
                    if Path(cached_path).exists():
                        Path(cached_path).unlink()  # also clean up the hf cache copy

                except Exception as e:
                    print(f"Attempt {attempt}/{max_retries} failed: {e}")
                    final_path = model_dir / filename
                    if final_path.exists():
                        final_path.unlink()  # clean up any partial download

                if attempt < max_retries:
                    print("Retrying in 3 seconds...")
                    time.sleep(3)  # brief pause before retry to let network recover

            print(f"Error: Failed to download {filename} after {max_retries} attempts.")
            return None

        # ------------------------------------------------------------
        # DOWNLOAD MODEL IF MISSING OR INCOMPLETE
        # both files are checked independently — if only the json is
        # missing or corrupt we skip re-downloading the 60MB onnx and
        # just fetch the small config file again
        # ------------------------------------------------------------
        onnx_valid = _is_valid_onnx(model_path)
        json_valid = _is_valid_json(model_json_path)

        if not onnx_valid or not json_valid:
            print("Piper voice model not found or incomplete. Downloading...")
            model_dir.mkdir(parents=True, exist_ok=True)  # create piper_models/ if it doesn't exist yet

            if not onnx_valid:
                if model_path.exists():
                    print(f"Removing incomplete {model_name}...")
                    model_path.unlink()  # delete the corrupt file before re-downloading
                print(f"Downloading {model_name} (~60MB)...")
                if _download_file(model_name, is_onnx=True) is None:
                    return  # download failed after all retries, abort
                print(f"{model_name} ready.")

            if not json_valid:
                if model_json_path.exists():
                    print(f"Removing incomplete {model_json}...")
                    model_json_path.unlink()  # delete the corrupt config before re-downloading
                print(f"Downloading {model_json}...")
                if _download_file(model_json, is_onnx=False) is None:
                    return  # download failed after all retries, abort
                print(f"{model_json} ready.")

            print("Piper model ready.")

        # ------------------------------------------------------------
        # SYNTHESIZE SPEECH
        # PiperVoice.load() reads the .onnx model into memory
        # SynthesisConfig carries volume into piper so it scales
        # the audio during synthesis — no manual audio work needed
        # synthesize_wav() runs the neural network and writes raw
        # PCM audio into the wav file handle
        # ------------------------------------------------------------
        voice = PiperVoice.load(str(model_path))  # load model from piper_models/
        syn_config = SynthesisConfig(volume=volume)  # package volume setting for piper

        # delete=False because the file must persist after the 'with'
        # block closes so aplay/winsound can open and play it —
        # we manually delete it ourselves after playback is done
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            tmp_path = tmp.name  # grab the path before the handle closes

        with wave.open(tmp_path, "wb") as wav_file:
            voice.synthesize_wav(text, wav_file, syn_config=syn_config)  # neural network runs here

        # ------------------------------------------------------------
        # PLAY AUDIO — platform specific
        # each OS has its own built-in audio playback tool:
        #   Linux   — aplay   (part of alsa-utils, -q means quiet/no output)
        #   Windows — winsound (Python standard library, Windows only)
        # ------------------------------------------------------------
        if platform.system() == "Linux":
            os.system(f"aplay -q {tmp_path}")
        elif platform.system() == "Windows":
            import winsound  # lazy — Windows only
            winsound.PlaySound(tmp_path, winsound.SND_FILENAME)

        # temp WAV is deleted immediately after playback —
        # only the model files in piper_models/ are kept permanently
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

    except ImportError:
        print("Error: piper-tts not installed. Run: pip install piper-tts")
        return
    except Exception as e:
        print(f"Error using piper TTS: {e}")
        return

    # integrates with log_setup() if active —
    # provides an audit trail of everything spoken during a session
    print(f"Spoken: {text}")


def screenshot(*args):
    """
    Takes a screenshot and saves it to the current working directory.

    Modes:
        1. Full screen, auto-named: ()
        2. Full screen, custom name: (filename)
        3. From (x,y) to screen edge, auto-named: (x, y)
        4. From (x,y) to screen edge, custom name: (x, y, filename)
        5. Specific region, auto-named: (x, y, width, height)
        6. Specific region, custom name: (x, y, width, height, filename)
        7. Selenium variants of all above: (driver, ...)

    Args:
        *args: Variable arguments depending on usage
            ::

                - ()                              : Full screen, auto-named
                - (filename)                      : Full screen, custom filename
                - (x, y)                          : From (x,y) to screen edge, auto-named
                - (x, y, filename)                : From (x,y) to screen edge, custom filename
                - (x, y, width, height)           : Specific region, auto-named
                - (x, y, width, height, filename) : Specific region, custom filename
                - (driver, ...)                   : Same as above but captures from Selenium browser

        Where:
            - driver: Selenium WebDriver instance
            - x, y: Top-left corner coordinates of the screenshot region
            - width, height: Dimensions of the screenshot region
            - filename: Custom name to save the screenshot
                - .png extension is added automatically if not provided
                - If not provided, auto-generates: screenshot_YYYY-MM-DD_HH-MM-SS_<unix>.png
                  Example: screenshot_2025-02-18_14-30-45_1708268445.png

    Returns:
        True if successful, False otherwise

    Output:
        - Prints the full path of the saved screenshot on success.
        - Prints error message if invalid arguments or coordinates are provided.

    Example:
        ::

            # Full screen (PyAutoGUI)
            screenshot()                                    # Full screen, auto-named
            screenshot('desktop.png')                       # Full screen, custom name

            # Selenium full page
            screenshot(driver)                              # Selenium full page, auto-named
            screenshot(driver, 'webpage.png')               # Selenium full page, custom name

            # From top-left point to edge (PyAutoGUI)
            screenshot(100, 200)                            # From (100,200) to edge, auto-named
            screenshot(100, 200, 'crop.png')                # From (100,200) to edge, custom name

            # From top-left point to edge (Selenium)
            screenshot(driver, 100, 200)                    # Selenium from (100,200) to edge
            screenshot(driver, 100, 200, 'page.png')        # Selenium, custom name

            # Specific region (PyAutoGUI)
            screenshot(0, 0, 500, 300)                      # Region: top-left (0,0), 500x300, auto-named
            screenshot(0, 0, 500, 300, 'region.png')        # Region: top-left (0,0), 500x300, custom name

            # Specific region (Selenium)
            screenshot(driver, 0, 0, 800, 600)              # Selenium region, auto-named
            screenshot(driver, 0, 0, 800, 600, 'sel.png')   # Selenium region, custom name
    """

    try:
        # ------------------------------------------------------------
        # DETERMINE MODE : Selenium or PyAutoGUI
        # ------------------------------------------------------------
        if len(args) > 0 and hasattr(args[0], 'save_screenshot'):
            # First argument is a WebDriver object
            use_driver = True
            driver_obj = args[0]
            remaining_args = args[1:]
        else:
            # PyAutoGUI mode : guard display check here
            if not _GUI_AVAILABLE:
                print("Error: screenshot() requires a display.")
                return False
            use_driver = False
            driver_obj = None
            remaining_args = args

        # Default values
        x, y, width, height, filename = None, None, None, None, None

        # ------------------------------------------------------------
        # PARSE ARGUMENTS
        # ------------------------------------------------------------
        if len(remaining_args) == 0:
            # screenshot() or screenshot(driver)
            # Full screen, auto-named
            x, y, width, height, filename = 0, 0, None, None, None

        elif len(remaining_args) == 1:
            # screenshot('file.png') or screenshot(driver, 'file.png')
            # Full screen, custom name
            x, y, width, height = 0, 0, None, None
            filename = remaining_args[0]

        elif len(remaining_args) == 2:
            # screenshot(100, 200) or screenshot(driver, 100, 200)
            # From (x,y) to screen edge, auto-named
            x, y = remaining_args[0], remaining_args[1]
            width, height, filename = None, None, None

        elif len(remaining_args) == 3:
            # screenshot(100, 200, 'file.png') or screenshot(driver, 100, 200, 'file.png')
            # From (x,y) to screen edge, custom name
            x, y, filename = remaining_args[0], remaining_args[1], remaining_args[2]
            width, height = None, None

        elif len(remaining_args) == 4:
            # screenshot(0, 0, 500, 300) or screenshot(driver, 0, 0, 500, 300)
            # Specific region, auto-named
            x, y, width, height = remaining_args[0], remaining_args[1], remaining_args[2], remaining_args[3]
            filename = None

        elif len(remaining_args) == 5:
            # screenshot(0, 0, 500, 300, 'file.png') or screenshot(driver, 0, 0, 500, 300, 'file.png')
            # Specific region, custom name
            x, y, width, height, filename = remaining_args

        else:
            raise ValueError("Too many arguments")

        # Convert coordinates to integers
        x = int(x) if x is not None else 0
        y = int(y) if y is not None else 0

        # ------------------------------------------------------------
        # GET SCREEN DIMENSIONS
        # ------------------------------------------------------------
        if use_driver:
            screen_width = driver_obj.execute_script("return window.innerWidth")
            screen_height = driver_obj.execute_script("return window.innerHeight")
        else:
            screen_width, screen_height = pyautogui.size()

        # Calculate width and height if not provided
        if width is None:
            width = screen_width - x  # From x to right edge
        if height is None:
            height = screen_height - y  # From y to bottom edge

        # Convert to integers
        width, height = int(width), int(height)

        # Validate
        if width <= 0 or height <= 0:
            raise ValueError("Width and height must be positive")
        if x < 0 or y < 0:
            raise ValueError("Coordinates must be non-negative")

        # ------------------------------------------------------------
        # GENERATE FILENAME
        # ------------------------------------------------------------
        if filename is None:
            now = datetime.datetime.now()
            timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
            unix_time = int(now.timestamp())
            filename = f"screenshot_{timestamp}_{unix_time}.png"
        else:
            # Add .png extension if missing
            if not os.path.splitext(filename)[1]:
                filename += '.png'

        # Get full path
        full_path = os.path.join(os.getcwd(), filename)

        # ------------------------------------------------------------
        # TAKE SCREENSHOT
        # ------------------------------------------------------------
        if use_driver:
            if x == 0 and y == 0 and width == screen_width and height == screen_height:
                # Full browser window, no cropping needed
                driver_obj.save_screenshot(full_path)
            else:
                # Region crop requested
                png_bytes = driver_obj.get_screenshot_as_png()
                image = Image.open(io.BytesIO(png_bytes))
                img_width, img_height = image.size
                crop_x2 = min(x + width, img_width)
                crop_y2 = min(y + height, img_height)
                cropped = image.crop((x, y, crop_x2, crop_y2))
                cropped.save(full_path)
        else:
            # PyAutoGUI: capture region directly
            img = pyautogui.screenshot(region=(x, y, width, height))
            img.save(full_path)

        print(f"Screenshot saved: {full_path}")
        return True

    except ValueError as e:
        print(f"Invalid input: {e}")
        return False
    except RuntimeError as e:
        print(f"Error: {e}")
        return False
    except Exception as e:
        print(f"Error taking screenshot: {e}")
        return False


def scroll(*args, timeout=30):
    """
    Universal scroll function for both PyAutoGUI and Selenium.

    Args:
        *args: Variable arguments (see examples below)
        timeout: Max seconds when scrolling to 'bottom'/'top' (default: 30)

    Returns:
        True if successful, False otherwise

    Output:
        - Prints scroll direction and count on completion.
        - Prints progress every 10 scrolls for large scroll counts.
        - Prints confirmation when bottom or top is reached (Selenium).

    Example:
        ::

            # PyAutoGUI Examples (scroll any window):
            scroll()                      # Scroll down 1 time (default)
            scroll(5)                     # Scroll down 5 times
            scroll('down')                # Scroll down 1 time
            scroll('down', 10)            # Scroll down 10 times
            scroll('up', 5)               # Scroll up 5 times
            scroll('bottom')              # Scroll down continuously for 30 seconds
            scroll('bottom', timeout=60)  # Scroll down continuously for 60 seconds
            scroll('top')                 # Scroll up continuously for 30 seconds

            # Selenium Examples (pass driver object):
            scroll(driver)                        # Scroll down 1 time in browser
            scroll(driver, 5)                     # Scroll down 5 times in browser
            scroll(driver, 'down')                # Scroll down 1 time in browser
            scroll(driver, 'down', 10)            # Scroll down 10 times in browser
            scroll(driver, 'up', 5)               # Scroll up 5 times in browser
            scroll(driver, 'bottom')              # Scroll to bottom (auto-detect end)
            scroll(driver, 'top')                 # Scroll to top (auto-detect end)
            scroll(driver, 'bottom', timeout=120) # Scroll to bottom, max 2 minutes
            scroll(driver, 'Login')               # Scroll to 1st instance of 'Login'
            scroll(driver, 'Login', 2)            # Scroll to 2nd instance of 'Login'
            scroll(driver, 'Login', -1)           # Scroll to last instance of 'Login'
            scroll(driver, 'Login', -2)           # Scroll to 2nd last instance of 'Login'
    """

    wait = 3  # Fixed wait time between scrolls (3 seconds)

    # --------------- PARSE ARGUMENTS ---------------
    if len(args) == 0:
        # scroll() - default: scroll down 1 time with PyAutoGUI
        use_selenium = False
        driver_obj = None
        direction = 'down'
        count = 1
        text_mode = False

    elif len(args) > 0 and hasattr(args[0], 'execute_script'):
        # Selenium mode - first arg is WebDriver object
        use_selenium = True
        driver_obj = args[0]
        text_mode = False

        if len(args) == 1:
            # scroll(driver) - scroll down 1 time
            direction = 'down'
            count = 1
        elif isinstance(args[1], int):
            # scroll(driver, 5) - scroll down 5 times
            direction = 'down'
            count = args[1]
        elif args[1] in ['down', 'up', 'bottom', 'top']:
            # scroll(driver, 'down') or scroll(driver, 'down', 10)
            direction = args[1]
            count = args[2] if len(args) > 2 else 1
        elif isinstance(args[1], str):
            # scroll(driver, 'Login') or scroll(driver, 'Login', 2) or scroll(driver, 'Login', -1)
            text_mode = True
            search_text = args[1]
            index = args[2] if len(args) > 2 else 1
        else:
            print(f"Error: Invalid argument '{args[1]}'. Use 'down', 'up', 'bottom', 'top' or a number.")
            return False

    else:
        # PyAutoGUI mode - guard display check here
        if not _GUI_AVAILABLE:
            print("Error: scroll() requires a display.")
            return False

        use_selenium = False
        driver_obj = None
        text_mode = False

        if isinstance(args[0], int):
            # scroll(5) - scroll down 5 times
            direction = 'down'
            count = args[0]
        elif args[0] in ['down', 'up', 'bottom', 'top']:
            # scroll('down') or scroll('down', 10)
            direction = args[0]
            count = args[1] if len(args) > 1 else 1
        else:
            print(f"Error: Invalid argument '{args[0]}'. Use 'down', 'up', 'bottom', 'top' or a number.")
            return False

    # Also guard the default scroll() case (no args = PyAutoGUI)
    if not use_selenium and not _GUI_AVAILABLE:
        print("Error: scroll() requires a display.")
        return False

    # Validate direction (only when not text mode)
    if not text_mode and direction not in ['down', 'up', 'bottom', 'top']:
        print(f"Error: Invalid direction '{direction}'.")
        return False

    try:
        # --------------- SELENIUM MODE ---------------
        # needs: nothing : works on all platforms
        if use_selenium:

            # --------------- SCROLL TO TEXT INSTANCE ---------------
            if text_mode:

                # Reset selection and scroll to top so find always starts from beginning
                driver_obj.execute_script("window.scrollTo(0, 0);")
                driver_obj.execute_script("window.getSelection().removeAllRanges();")

                # Count total instances using TreeWalker (text nodes only, no DOM corruption)
                total = driver_obj.execute_script("""
                    var text = arguments[0].toLowerCase();
                    var count = 0;
                    var walker = document.createTreeWalker(document.body, NodeFilter.SHOW_TEXT, null, false);
                    while (walker.nextNode()) {
                        var nodeText = walker.currentNode.nodeValue.toLowerCase();
                        var pos = 0;
                        while ((pos = nodeText.indexOf(text, pos)) !== -1) {
                            count++;
                            pos += text.length;
                        }
                    }
                    return count;
                """, search_text)

                if total == 0:
                    print(f"Error: No instances of '{search_text}' found.")
                    return False

                # Resolve index (1-based positive, negative from end)
                if index == 0:
                    print(f"Error: index=0 is invalid. Use 1 to {total} or -1 to -{total}.")
                    return False
                elif index < 0:
                    if abs(index) > total:
                        print(f"Error: index={index} out of range. Found {total} instance(s) of '{search_text}'.")
                        return False
                    resolved_index = total + index + 1
                else:
                    if index > total:
                        print(f"Error: index={index} out of range. Found {total} instance(s) of '{search_text}'.")
                        return False
                    resolved_index = index

                # Use window.find() resolved_index times to land on correct instance
                # window.find() works like Ctrl+F - highlights just the word and scrolls to it
                driver_obj.execute_script("""
                    var text = arguments[0];
                    var count = arguments[1];
                    for (var i = 0; i < count; i++) {
                        window.find(text, false, false, true, false, false, false);
                    }
                """, search_text, resolved_index)

                print(f"[Selenium] Scrolled to '{search_text}' (instance {resolved_index} of {total})")
                return True

            # --------------- SCROLL TO BOTTOM ---------------
            elif direction == 'bottom':
                print(f"[Selenium] Scrolling to bottom (timeout={timeout}s)...")
                start_time = time.time()
                scrolls = 0

                while time.time() - start_time < timeout:
                    driver_obj.execute_script("window.scrollBy(0, 1000);")
                    scrolls += 1
                    time.sleep(1)

                    new_height = driver_obj.execute_script("return document.documentElement.scrollHeight")
                    scroll_position = driver_obj.execute_script("return window.pageYOffset + window.innerHeight")

                    # Check if scroll position reached the bottom
                    if scroll_position >= new_height:
                        print(f"Reached bottom after {scrolls} scrolls ({time.time() - start_time:.1f}s)")
                        return True

                    if scrolls % 10 == 0:
                        print(f"Scrolled {scrolls} times ({time.time() - start_time:.1f}s)")

                print(f"Scrolled {scrolls} times")
                return True

            # --------------- SCROLL TO TOP ---------------
            elif direction == 'top':
                print(f"[Selenium] Scrolling to top...")
                driver_obj.execute_script("window.scrollTo(0, 0);")
                time.sleep(wait)
                print("Scrolled to top")
                return True

            # --------------- SCROLL DOWN / UP N TIMES ---------------
            else:
                scroll_pixels = 500 if direction == 'down' else -500
                print(f"[Selenium] Scrolling {direction} {count} time(s)...")

                for i in range(count):
                    driver_obj.execute_script(f"window.scrollBy(0, {scroll_pixels});")
                    time.sleep(wait)

                    if count > 10 and (i + 1) % 10 == 0:
                        print(f"Scrolled {i + 1}/{count} times")

                print(f"Scrolled {direction} {count} time(s)")
                return True

        # --------------- PYAUTOGUI MODE ---------------
        # needs: display (already guarded above)
        else:

            # --------------- SCROLL TO BOTTOM / TOP ---------------
            if direction in ['bottom', 'top']:
                scroll_amount = -500 if direction == 'bottom' else 500
                print(f"[PyAutoGUI] Scrolling {direction} continuously for {timeout}s...")

                start_time = time.time()
                scrolls = 0

                while time.time() - start_time < timeout:
                    pyautogui.scroll(scroll_amount)
                    scrolls += 1
                    time.sleep(wait)

                    if scrolls % 10 == 0:
                        print(f"Scrolled {scrolls} times ({time.time() - start_time:.1f}s)")

                print(f"Scrolled {direction} for {timeout}s ({scrolls} total scrolls)")
                return True

            # --------------- SCROLL DOWN / UP N TIMES ---------------
            else:
                scroll_amount = -500 if direction == 'down' else 500
                print(f"[PyAutoGUI] Scrolling {direction} {count} time(s)...")

                for i in range(count):
                    pyautogui.scroll(scroll_amount)
                    time.sleep(wait)

                    if count > 10 and (i + 1) % 10 == 0:
                        print(f"Scrolled {i + 1}/{count} times")

                print(f"Scrolled {direction} {count} time(s)")
                return True

    except Exception as e:
        print(f"Error while scrolling: {e}")
        return False


def second():
    """
    Get current second.

    Returns:
        int: Current second (0-59)

    Example:
        second()  # 45
    """
    return time.localtime().tm_sec


def wait(*args, countdown=True):
    """
    Wait with countdown, wait for element or wait for color at pixel.

    Args:
        *args: Variable arguments (see examples)
        countdown: If True, shows countdown display (default: True)

    Returns:
        True if successful, False if error or timeout

    Example:
        ::

            # Countdown wait
            wait(5)                              # Wait 5 seconds with countdown
            wait(10, countdown=False)            # Wait 10 seconds silently
            wait()                               # Wait 3 seconds (default)

            # Wait for element (Selenium) - pass driver object
            wait(driver, 'xpath', '//button')            # Wait max 180s with countdown
            wait(driver, 'id', 'submit-btn', 10)         # Wait max 10s with countdown
            wait(driver, 'class', 'content', 30, countdown=False)  # Wait silently for 30s

            # Wait for color at pixel
            wait(100, 200, 255, 0, 0)            # Wait for red (255,0,0) at (100,200) with countdown
            wait(100, 200, 255, 0, 0, 30)        # Wait for red, max 30s with countdown
            wait(500, 300, 0, 255, 0, 60, countdown=False)  # Wait silently
    """

    # if no argument is passed
    if len(args) == 0:
        args = (3,)  # Set default to 3 seconds

    # ------------------------------------------------------------
    # MODE 1 : COUNTDOWN (1 argument, integer or float)
    # needs: nothing : safe on all platforms
    # ------------------------------------------------------------
    if len(args) == 1 and isinstance(args[0], (int, float)):
        seconds = args[0]

        if seconds < 0:
            raise ValueError("Seconds must be non-negative")

        if seconds == 0:
            return True  # No wait needed

        # Convert to int for countdown
        seconds_int = int(seconds)

        if countdown:
            print('Waiting:')
            for remaining in range(seconds_int, 0, -1):
                print(f'..{remaining} ', end='\r', flush=True)
                time.sleep(1)
            print('Done   ')
        else:
            time.sleep(seconds)

        return True

    # ------------------------------------------------------------
    # MODE 2 : WAIT FOR ELEMENT (driver object + 2 or 3 arguments)
    # needs: nothing : works on all platforms
    # ------------------------------------------------------------
    elif len(args) >= 3 and hasattr(args[0], 'find_element'):
        # Selenium mode - driver object detected
        driver_obj = args[0]
        selector_type = args[1]
        selector = args[2]
        timeout = args[3] if len(args) == 4 else 180  # Default 180s (3 minutes)

        # Validate selector type
        if selector_type not in ['id', 'xpath', 'class', 'name', 'css', 'tag', 'text', 'partial']:
            raise ValueError(f"Invalid selector type: {selector_type}")

        if countdown:
            print(f'Waiting for element ({selector_type}: {selector}):')

        # Save original implicit wait and set to 0 for fast checking
        original_implicit_wait = driver_obj.timeouts.implicit_wait
        driver_obj.implicitly_wait(0)

        start_time = time.time()

        # Countdown approach: check element every second, print countdown
        for remaining in range(timeout, 0, -1):
            # Check if element exists using get_web_element
            element = _get_web_element(driver_obj, selector_type, selector)

            # Element found
            if element:
                elapsed = time.time() - start_time
                # Restore original implicit wait
                driver_obj.implicitly_wait(original_implicit_wait)
                print(f"Element found after {elapsed:.1f}s")
                return True

            # Print countdown
            if countdown:
                print(f'..{remaining} ', end='\r', flush=True)

            # Wait 1 second before next check
            time.sleep(1)

        # Restore original implicit wait
        driver_obj.implicitly_wait(original_implicit_wait)

        # Timeout reached
        print(f"Element not found after {timeout}s timeout.")
        return False

    # ------------------------------------------------------------
    # MODE 3 : WAIT FOR COLOR AT PIXEL (5 or 6 arguments, all integers)
    # needs: display
    # ------------------------------------------------------------
    elif len(args) in [5, 6] and all(isinstance(arg, int) for arg in args):
        if not _GUI_AVAILABLE:
            print("Error: wait() color mode requires a display.")
            return False

        x = args[0]
        y = args[1]
        target_r = args[2]
        target_g = args[3]
        target_b = args[4]
        timeout = args[5] if len(args) == 6 else 180  # Default 180s

        if countdown:
            print(f'Waiting at ({x},{y}) for color RGB({target_r},{target_g},{target_b}) :')

        start_time = time.time()

        # Countdown approach: check color every second, print countdown
        for remaining in range(timeout, 0, -1):
            try:
                # Check if color matches at pixel
                if pyautogui.pixelMatchesColor(x, y, (target_r, target_g, target_b)):
                    elapsed = time.time() - start_time
                    print(f"Color found after {elapsed:.1f}s")
                    return True

            except Exception as e:
                raise RuntimeError(f"Error checking pixel color: {e}")

            # Print countdown
            if countdown:
                print(f'..{remaining} ', end='\r', flush=True)

            # Wait 1 second before next check
            time.sleep(1)

        # Timeout reached
        print(f"Color not found after {timeout}s timeout.")
        return False

    # ---------------Invalid Arguments---------------
    else:
        raise ValueError("Invalid arguments for wait()")


def wait_download(timeout=1200, url=None, filename=None, download_dir=None):
    """
    Wait for a browser-initiated download to complete or download a file directly via URL.

    Modes:
        1. URL mode (url provided): Downloads file directly using requests. Useful for
           file types blocked by browsers (.exe, .msix, .msi, etc.). File is saved to
           Python's current working directory.
        2. Monitor mode (url not provided): Monitors the downloads folder for a
           browser-initiated download to complete.

    Args:
        timeout: Maximum seconds to wait for download completion (default: 1200)
        url: Direct download URL (optional)::

            - If provided: Downloads file directly via requests
            - If None: Monitors downloads folder for browser-initiated download
        filename: Custom filename to save/rename the downloaded file (optional)::

            - With extension (e.g. "myapp.exe")  : used as-is
            - Without extension (e.g. "myapp")   : extension borrowed from original file
            - If None: Original filename is kept
            - If multiple files are downloaded, only the first completed file is renamed
        download_dir: Custom download directory to monitor (monitor mode only, optional)::

            - If provided: Uses specified path and skips all auto-detection
            - If None: Auto-detects using the priority order described in Note below

    Returns::

        str: Full path of the downloaded file (always includes extension) on success
        False: On failure (download error, timeout, directory access issue, etc.)

    Output:
        - Prints download progress every 10 seconds showing elapsed time and file size.
        - Prints confirmation with final filename and saved path on completion.
        - Prints timeout message if download does not complete in time.

    Example:
        ::

            wait_download()                                                   # Monitor downloads folder
            wait_download(url='https://abc.com/file.msix')                    # Direct download via URL
            wait_download(url='https://abc.com/file.msix', filename='myapp')  # Custom name, borrows extension
            wait_download(300, filename='our_log.txt')                        # Monitor and rename on completion
            wait_download(600, download_dir='/downloads')                     # Docker with custom path
            wait_download(300, download_dir='D:/MyDownloads')                 # Windows custom path

            # Use with browser() : pass driver.download_dir to guarantee alignment
            driver = browser('https://example.com')
            click(driver, 'id', 'download-button')
            wait_download(download_dir=driver.download_dir)

    Note:
        - When download_dir is not provided, the folder is auto-detected in this order:
            1. DOWNLOAD_DIR environment variable (if set at OS level)
            2. /downloads folder (if running inside Docker)
            3. ~/Downloads (default fallback)
        - If a file was modified within the last 20 seconds before calling this function,
          it will be detected as a recently completed download and returned immediately.
          This handles cases where downloads complete very quickly before monitoring starts.
    """

    def _resolve_final_filename(custom_name, original_name):
        """Resolve the final filename, borrowing extension from original if custom name has none."""
        if custom_name is None:
            return original_name
        original_ext = os.path.splitext(original_name)[1]
        custom_base, custom_ext = os.path.splitext(custom_name)
        return custom_name if custom_ext else custom_base + original_ext

    def _rename_if_needed(directory, original_name, final_name):
        """Rename the downloaded file if a custom filename was provided."""
        if final_name != original_name:
            os.rename(
                os.path.join(directory, original_name),
                os.path.join(directory, final_name)
            )
            print(f'"{original_name}" renamed to "{final_name}"')

    def _format_size(bytes_count):
        """Convert bytes to a human-readable MB string."""
        return f"{bytes_count / (1024 * 1024):.1f} MB"

    # ------------------------------------------------------------
    # MODE 1 : Direct download via requests
    # ------------------------------------------------------------
    if url is not None:
        try:
            original_filename = url.split('/')[-1].split('?')[0]
            final_filename = _resolve_final_filename(filename, original_filename)
            print(f"Downloading: {final_filename}")
            response = requests.get(url, stream=True, timeout=30)  # 30s for connection/chunk timeout
            response.raise_for_status()
            deadline = time.time() + (timeout)  # enforce total download time limit
            download_start = time.time()
            last_print_time = 0
            total_bytes = 0
            with open(final_filename, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if time.time() > deadline:
                        print(f"Timeout after {timeout} seconds while downloading '{final_filename}'")
                        return False
                    f.write(chunk)
                    total_bytes += len(chunk)
                    elapsed_sec = int(time.time() - download_start)
                    # Print progress every 10 seconds
                    if elapsed_sec - last_print_time >= 10:
                        elapsed = str(datetime.timedelta(seconds=elapsed_sec))
                        print(f"Downloading... {elapsed} elapsed ({_format_size(total_bytes)})")
                        last_print_time = elapsed_sec
            print(f"Downloaded file saved at: {os.path.abspath(final_filename)}")
            return os.path.abspath(final_filename)
        except Exception as e:
            print(f"Could not download file: {e}")
            return False

    # ------------------------------------------------------------
    # MODE 2 : Monitor browser downloads folder
    # ------------------------------------------------------------

    # Determine download directory using the following priority order:
    # 1. User provided path via download_dir argument
    # 2. DOWNLOAD_DIR environment variable (if set)
    # 3. /downloads folder (if running inside Docker)
    # 4. ~/Downloads folder (default, if none of the above apply)
    if download_dir is not None:
        download_dir = str(download_dir)
    elif os.getenv('DOWNLOAD_DIR'):
        download_dir = os.getenv('DOWNLOAD_DIR')
        print(f'Using DOWNLOAD_DIR from environment: {download_dir}')
    elif os.path.exists('/.dockerenv'):
        download_dir = '/downloads'
        print(f'Docker detected, using: {download_dir}')
    else:
        download_dir = str(Path.home() / "Downloads")

    # Ensure the download directory exists
    if not os.path.exists(download_dir):
        print(f"Warning: Download directory does not exist: {download_dir}")
        print(f"Attempting to create it...")
        try:
            os.makedirs(download_dir, exist_ok=True)
            print(f"Created directory: {download_dir}")
        except Exception as e:
            print(f"Error: Could not create directory: {e}")
            return False

    print(f'Monitoring downloads in: {download_dir}')

    # Temporary file extensions used by different browsers while downloading
    temp_extensions = ('.crdownload', '.part', '.download', '.tmp', '.temp')

    # Snapshot the initial state of the downloads folder before monitoring begins
    try:
        initial_files = set(os.listdir(download_dir))
        initial_temp_files = set([f for f in initial_files if f.endswith(temp_extensions)])

        # Record sizes of any temp files that already exist, to detect if they are actively downloading
        initial_temp_sizes = {}
        for f in initial_temp_files:
            try:
                initial_temp_sizes[f] = os.path.getsize(os.path.join(download_dir, f))
            except Exception as e:
                print(f"Could not get file size of {f}: {e}")

    except Exception as e:
        print(f"Error accessing download directory: {e}")
        return False

    # Wait 10 seconds before first check to allow download to start
    print("Checking for active downloads...")
    time.sleep(10)

    # Check if a file was already downloaded quickly (modified within last 20 seconds)
    try:
        current_files_after_wait = set(os.listdir(download_dir))
        current_time = time.time()

        recent_files = []
        for f in current_files_after_wait:
            if not f.endswith(temp_extensions):
                try:
                    mtime = os.path.getmtime(os.path.join(download_dir, f))
                    age = current_time - mtime
                    if age <= 20:
                        recent_files.append((f, age))
                except Exception as e:
                    print(f"Could not get modification time of {f}: {e}")

        if recent_files:
            recent_files.sort(key=lambda x: x[1])  # sort by age, most recent first
            most_recent_file, age = recent_files[0]
            final_filename = _resolve_final_filename(filename, most_recent_file)
            _rename_if_needed(download_dir, most_recent_file, final_filename)
            print(f"Quick download detected: '{final_filename}' (modified {int(age)} seconds ago)")
            print(f"Downloaded file saved at: {os.path.join(download_dir, final_filename)}")
            return os.path.join(download_dir, final_filename)

    except Exception as e:
        print(f"Error checking for recent files: {e}")

    download_time = 10  # already waited 10 seconds above
    last_print_time = 0
    download_started = False
    monitoring_files = set()
    completed_files_so_far = []  # tracks files already reported as completed mid-session
    filename_used = False  # ensures filename rename only applies to first completed file

    while download_time < timeout:
        try:
            current_files = set(os.listdir(download_dir))
            current_temp_files = set([f for f in current_files if f.endswith(temp_extensions)])

            # Get current sizes of all active temp files
            current_temp_sizes = {}
            for f in current_temp_files:
                try:
                    current_temp_sizes[f] = os.path.getsize(os.path.join(download_dir, f))
                except Exception as e:
                    print(f"Could not get file size of {f}: {e}")

            new_temp_files = current_temp_files - initial_temp_files

            # On the first pass (10s mark), check if any pre-existing temp files are actively growing
            if download_time == 10:
                active_old_temp_files = set()
                for f in initial_temp_files:
                    if f in current_temp_sizes and f in initial_temp_sizes:
                        if current_temp_sizes[f] != initial_temp_sizes[f]:
                            active_old_temp_files.add(f)
                            monitoring_files.add(f)

                if active_old_temp_files:
                    download_started = True
                    if len(active_old_temp_files) == 1:
                        print(f"Active download detected: '{list(active_old_temp_files)[0]}'")
                    else:
                        print(
                            f"{len(active_old_temp_files)} active downloads detected: {', '.join(active_old_temp_files)}")

            # Track any brand new temp files that appeared after monitoring started
            if new_temp_files:
                truly_new = new_temp_files - monitoring_files
                if truly_new:
                    if not download_started:
                        download_started = True
                    monitoring_files.update(truly_new)
                    if len(truly_new) == 1:
                        print(f"New download started: '{list(truly_new)[0]}'")
                    else:
                        print(f"{len(truly_new)} new downloads started: {', '.join(truly_new)}")

            # Check for any newly completed files : catches mid-session completions including
            # small files that complete too quickly for temp file tracking to catch
            newly_completed = [
                f for f in (current_files - initial_files)
                if not f.endswith(temp_extensions)
                   and f not in completed_files_so_far
            ]
            for f in newly_completed:
                elapsed = str(datetime.timedelta(seconds=download_time))
                if not filename_used:
                    final_name = _resolve_final_filename(filename, f)
                    _rename_if_needed(download_dir, f, final_name)
                    filename_used = True
                else:
                    final_name = f
                    if filename is not None:
                        print(f'"{f}" kept original name (filename only applies to first file)')
                print(f"Download completed: '{final_name}' (took {elapsed})")
                completed_files_so_far.append(final_name)

            all_monitoring = (current_temp_files & monitoring_files) | new_temp_files

            # Return immediately if a file just completed and no temp files are still active
            if newly_completed and not (current_temp_files & monitoring_files):
                primary = completed_files_so_far[0]
                print(f"Downloaded file saved at: {os.path.join(download_dir, primary)}")
                return os.path.join(download_dir, primary)

            if all_monitoring:
                # Still downloading : print progress every 10 seconds with size of each file
                if download_time - last_print_time >= 10:
                    elapsed = str(datetime.timedelta(seconds=download_time))
                    if len(all_monitoring) == 1:
                        f = list(all_monitoring)[0]
                        size_str = _format_size(current_temp_sizes.get(f, 0))
                        print(f"Downloading... {elapsed} elapsed ({size_str})")
                    else:
                        size_parts = ', '.join(
                            f"{f}: {_format_size(current_temp_sizes.get(f, 0))}"
                            for f in all_monitoring
                        )
                        print(f"Downloading... {elapsed} elapsed ({size_parts})")
                    last_print_time = download_time

            elif download_started:
                # All temp files are gone : find any remaining unreported completed files
                new_files = current_files - initial_files
                remaining_completed = [
                    f for f in new_files
                    if not f.endswith(temp_extensions)
                       and f not in completed_files_so_far
                ]

                if remaining_completed:
                    elapsed = str(datetime.timedelta(seconds=download_time))
                    for f in remaining_completed:
                        if not filename_used:
                            final_name = _resolve_final_filename(filename, f)
                            _rename_if_needed(download_dir, f, final_name)
                            filename_used = True
                        else:
                            final_name = f
                            if filename is not None:
                                print(f'"{f}" kept original name (filename only applies to first file)')
                        print(f"Download completed: '{final_name}' (took {elapsed})")
                        completed_files_so_far.append(final_name)

                    # Return the first completed file as the primary result
                    primary = completed_files_so_far[0]
                    print(f"Downloaded file saved at: {os.path.join(download_dir, primary)}")
                    return os.path.join(download_dir, primary)
                else:
                    # Temp files gone but no completed file found yet : keep waiting briefly
                    if download_time - last_print_time >= 10:
                        elapsed = str(datetime.timedelta(seconds=download_time))
                        print(f"Verifying download... (elapsed: {elapsed})")
                        last_print_time = download_time

            else:
                # No download detected yet : no size to report, just print elapsed time
                if download_time - last_print_time >= 10:
                    recent_complete_files = [
                        f for f in current_files
                        if not f.endswith(temp_extensions) and f not in initial_files
                    ]

                    if recent_complete_files:
                        original_name = recent_complete_files[0]
                        final_filename = _resolve_final_filename(filename, original_name)
                        _rename_if_needed(download_dir, original_name, final_filename)
                        print(f"Quick download detected: '{final_filename}'")
                        print(f"Downloaded file saved at: {os.path.join(download_dir, final_filename)}")
                        return os.path.join(download_dir, final_filename)

                    elapsed = str(datetime.timedelta(seconds=download_time))
                    print(f"Waiting for download to start... {elapsed} elapsed")
                    last_print_time = download_time

        except Exception as e:
            print(f'Error monitoring downloads: {e}')
            return False

        time.sleep(1)
        download_time += 1

    # Report any files that completed mid-session before the timeout
    if completed_files_so_far:
        if len(completed_files_so_far) == 1:
            print(f"Note: 1 file completed during session: '{list(completed_files_so_far)[0]}'")
        else:
            print(f"Note: {len(completed_files_so_far)} files completed during session:")
            for f in completed_files_so_far:
                print(f"  - {f}")

    # Timeout reached : report what was still in progress if anything
    if monitoring_files or download_started:
        current_files = set(os.listdir(download_dir))
        current_temp_files = set([f for f in current_files if f.endswith(temp_extensions)])
        still_downloading = current_temp_files & monitoring_files

        if still_downloading:
            if len(still_downloading) == 1:
                print(f"Timeout after {timeout} seconds while waiting for '{list(still_downloading)[0]}' to complete.")
            else:
                print(
                    f"Timeout after {timeout} seconds while waiting for {len(still_downloading)} files: {', '.join(still_downloading)}")
        else:
            print(f'Timeout after {timeout} seconds. Download status unclear.')
    else:
        print(f'Timeout after {timeout} seconds. No download detected.')

    return False


def window(action=None, target=None, *args):
    """
    Unified window management function for Windows and Linux.

    Args:
        action: Window operation to perform (default: 'list')::

            'list'     : Get all window titles
            'title'    : Get active window title (or find full title if target provided)
            'focus'    : Bring window to foreground
            'close'    : Close window
            'minimize' : Minimize window
            'maximize' : Maximize window
            'resize'   : Resize window (requires width, height)
            'move'     : Move window (requires x, y)

        target: Window title or pattern (required for most actions)
        *args: Additional parameters (width, height for resize; x, y for move)

    Returns:
        Return type depends on action::

            list/None : List of strings when action is None or 'list'
            title     : String or None
            others    : True if successful, False otherwise

    Raises:
        ValueError: If invalid action, missing required parameters or invalid
                    dimensions/coordinates
        NotImplementedError: If called on macOS

    Output:
        - Prints error if window not found, with suggestions for similar window titles
          (focus action only).
        - Prints error if wmctrl or xdotool is not installed (Linux only).

    Example:
        ::

            # Get all windows (default)
            window()                                    # ['Chrome', 'Notepad', 'Excel']
            window('list')                              # ['Chrome', 'Notepad', 'Excel']

            # Check if window exists
            if 'Chrome' in window():
                print("Chrome is open!")

            # Get active window title
            window('title')                             # 'Google Chrome - New Tab'

            # Find full title containing text
            window('title', 'Chrome')                   # 'Google Chrome - New Tab'
            window('title', 'Note')                     # 'Untitled - Notepad'

            # Window operations
            window('focus', 'Chrome')                   # Focus window
            window('close', 'Notepad')                  # Close window
            window('minimize', 'Excel')                 # Minimize window
            window('maximize', 'Word')                  # Maximize window

            # Position and size
            window('resize', 'Chrome', 800, 600)        # Resize to 800x600
            window('move', 'Chrome', 100, 200)          # Move to (100, 200)

            # Side-by-side setup (1920x1080 display)
            window('move', 'Chrome', 0, 0)              # Left side
            window('resize', 'Chrome', 960, 1080)       # Half screen width
            window('move', 'Code', 960, 0)              # Right side
            window('resize', 'Code', 960, 1080)         # Half screen width

            # Recording setup
            window('resize', 'Demo', 1280, 720)         # 720p size
            window('move', 'Demo', 320, 180)            # Centered on 1920x1080

    Note:
        - On Linux, resize and move automatically remove maximized/minimized
          state before applying changes, ensuring consistent behavior.
        - Target matching is case-insensitive and partial (e.g. 'Chrome' matches
          'Google Chrome - New Tab').
    """

    system = platform.system()

    # Check if macOS and reject
    if system == "Darwin":
        raise NotImplementedError(
            "window() is not supported on macOS.\n"
            "Supported platforms: Windows, Linux\n"
            "Note: Selenium functions will still work on macOS."
        )

    # Default action: list all windows
    if action is None:
        action = 'list'

    action = action.lower()

    # Validate action
    valid_actions = ['list', 'title', 'focus', 'close', 'minimize', 'maximize', 'resize', 'move']
    if action not in valid_actions:
        raise ValueError(f"Invalid action '{action}'. Valid actions: {', '.join(valid_actions)}")

    # Check if target is required
    if action in ['focus', 'close', 'minimize', 'maximize', 'resize', 'move'] and target is None:
        raise ValueError(f"Action '{action}' requires a target window")

    try:
        # ------------------------------------------------------------
        # LIST - Get all window titles
        # ------------------------------------------------------------
        if action == 'list':
            if system == "Windows":
                titles = []

                def enum_handler(hwnd, ctx):
                    if win32gui.IsWindowVisible(hwnd):
                        title = win32gui.GetWindowText(hwnd)
                        if title.strip():
                            titles.append(title)

                win32gui.EnumWindows(enum_handler, None)
                return titles

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    titles = []
                    for line in result.stdout.strip().split('\n'):
                        parts = line.split(None, 3)
                        if len(parts) >= 4:
                            titles.append(parts[3])
                    return titles
                else:
                    print("Error: wmctrl not installed.")
                    print("Install with: sudo apt-get install wmctrl")
                    return []

        # ------------------------------------------------------------
        # TITLE - Get active window title OR find window title
        # ------------------------------------------------------------
        elif action == 'title':
            # No target - get active window title
            if target is None:
                if system == "Windows":
                    title = win32gui.GetWindowText(win32gui.GetForegroundWindow()).strip()
                    return title if title else None

                elif system == "Linux":
                    result = subprocess.run(['xdotool', 'getactivewindow', 'getwindowname'],
                                            capture_output=True, text=True)
                    if result.returncode == 0:
                        return result.stdout.strip()
                    else:
                        print("Error: xdotool not installed.")
                        print("Install with: sudo apt-get install xdotool")
                        return None

            # Target provided - find full title containing target text
            else:
                all_windows = window('list')

                # Case-insensitive contains match
                for win_title in all_windows:
                    if target.lower() in win_title.lower():
                        return win_title

                # Not found
                return None

        # ------------------------------------------------------------
        # FOCUS - Bring window to foreground
        # ------------------------------------------------------------
        elif action == 'focus':
            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")

                # Suggest similar windows
                suggestions = get_close_matches(target, all_windows, n=3, cutoff=0.6)
                if suggestions:
                    print(f"Did you mean: {', '.join(suggestions)}")

                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                    win32gui.SetForegroundWindow(hwnd)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-a', window_title],
                                        capture_output=True, text=True)
                return result.returncode == 0

        # ------------------------------------------------------------
        # CLOSE - Close window
        # ------------------------------------------------------------
        elif action == 'close':
            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-c', window_title],
                                        capture_output=True, text=True)
                return result.returncode == 0

        # ------------------------------------------------------------
        # MINIMIZE - Minimize window
        # ------------------------------------------------------------
        elif action == 'minimize':
            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.ShowWindow(hwnd, win32con.SW_MINIMIZE)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    for line in result.stdout.strip().split('\n'):
                        if window_title in line:
                            window_id = line.split()[0]
                            subprocess.run(['xdotool', 'windowminimize', window_id])
                            return True
                return False

        # ------------------------------------------------------------
        # MAXIMIZE - Maximize window
        # ------------------------------------------------------------
        elif action == 'maximize':
            if not window('focus', target):
                return False

            time.sleep(0.5)

            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    win32gui.ShowWindow(hwnd, win32con.SW_MAXIMIZE)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-r', window_title, '-b',
                                         'add,maximized_vert,maximized_horz'],
                                        capture_output=True, text=True)
                return result.returncode == 0

        # ------------------------------------------------------------
        # RESIZE - Resize window to specific dimensions
        # ------------------------------------------------------------
        elif action == 'resize':
            if len(args) < 2:
                raise ValueError("Action 'resize' requires width and height: window('resize', 'Chrome', 800, 600)")

            try:
                width = int(args[0])
                height = int(args[1])
            except (ValueError, TypeError):
                raise ValueError("Width and height must be integers")

            if width <= 0 or height <= 0:
                raise ValueError("Width and height must be positive")

            if not window('focus', target):
                return False

            time.sleep(0.5)

            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    rect = win32gui.GetWindowRect(hwnd)
                    x, y = rect[0], rect[1]
                    win32gui.MoveWindow(hwnd, x, y, width, height, True)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    for line in result.stdout.strip().split('\n'):
                        if window_title in line:
                            window_id = line.split()[0]
                            # Remove maximized/minimized state before resizing
                            subprocess.run(
                                ['wmctrl', '-i', '-r', window_id, '-b', 'remove,maximized_vert,maximized_horz'])
                            subprocess.run(['xdotool', 'windowmap', window_id])  # restore if minimized
                            time.sleep(0.3)
                            subprocess.run(['wmctrl', '-i', '-r', window_id, '-e', f'0,-1,-1,{width},{height}'])
                            return True
                return False

        # ------------------------------------------------------------
        # MOVE - Move window to specific coordinates
        # ------------------------------------------------------------
        elif action == 'move':
            if len(args) < 2:
                raise ValueError("Action 'move' requires x and y coordinates: window('move', 'Chrome', 100, 200)")

            try:
                x = int(args[0])
                y = int(args[1])
            except (ValueError, TypeError):
                raise ValueError("X and Y coordinates must be integers")

            if x < 0 or y < 0:
                raise ValueError("X and Y coordinates must be non-negative")

            if not window('focus', target):
                return False

            time.sleep(0.5)

            all_windows = window('list')
            matches = [w for w in all_windows if target.lower() in w.lower()]

            if not matches:
                print(f"Window not found: {target}")
                return False

            window_title = matches[0]

            if system == "Windows":

                def find_window_hwnd(title):
                    hwnd_list = []

                    def enum_handler(hwnd, ctx):
                        if win32gui.IsWindowVisible(hwnd):
                            win_text = win32gui.GetWindowText(hwnd)
                            if win_text == title:
                                hwnd_list.append(hwnd)

                    win32gui.EnumWindows(enum_handler, None)
                    return hwnd_list[0] if hwnd_list else None

                hwnd = find_window_hwnd(window_title)
                if hwnd:
                    rect = win32gui.GetWindowRect(hwnd)
                    width = rect[2] - rect[0]
                    height = rect[3] - rect[1]
                    win32gui.MoveWindow(hwnd, x, y, width, height, True)
                    return True
                return False

            elif system == "Linux":
                result = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True)
                if result.returncode == 0:
                    for line in result.stdout.strip().split('\n'):
                        if window_title in line:
                            window_id = line.split()[0]
                            # Remove maximized/minimized state before moving
                            subprocess.run(
                                ['wmctrl', '-i', '-r', window_id, '-b', 'remove,maximized_vert,maximized_horz'])
                            subprocess.run(['xdotool', 'windowmap', window_id])  # restore if minimized
                            time.sleep(0.3)
                            subprocess.run(['wmctrl', '-i', '-r', window_id, '-e', f'0,{x},{y},-1,-1'])
                            return True
                return False

    except Exception as e:
        print(f"Error in window operation: {e}")
        return False if action in ['focus', 'close', 'minimize', 'maximize', 'resize', 'move'] else None


def write(*keys):
    """
    Write or type text using keyboard (PyAutoGUI or Selenium).

    Args:
        *keys: Variable arguments depending on mode

    Returns:
        True if successful, False otherwise

    Output:
        - Prints error message if element was not found (Selenium).
        - Prints error message if invalid arguments are provided.

    Example:
        ::

            # PyAutoGUI mode (types in any active window)
            write("Hello World")                                        # Types in active window
            write("user@example.com")                                   # Types email
            write("12345")                                              # Types numbers as string

            # Selenium mode - type on active element in browser
            write(driver, "Hello World")                                # Types on active element
            write(driver, "Search query")                               # Types in focused input

            # Selenium mode - type in specific element
            write(driver, "id", "username", "john_doe")
            write(driver, "xpath", "//input[@name='email']", "user@example.com")
            write(driver, "class", "search-box", "Python tutorial")

    Note:
        - PyAutoGUI uses typewrite() which types one character at a time.
        - Selenium uses send_keys() which types the entire string at once.
    """

    try:
        # ------------------------------------------------------------
        # MODE 1 : PYAUTOGUI (1 argument - text only)
        # needs: display
        # ------------------------------------------------------------
        if len(keys) == 1:
            if not _GUI_AVAILABLE:
                print("Error: write() requires a display.")
                return False
            pyautogui.typewrite(keys[0])
            time.sleep(1)
            return True

        # ------------------------------------------------------------
        # SELENIUM MODE (driver object detected)
        # needs: nothing : works on all platforms
        # ------------------------------------------------------------
        elif len(keys) >= 2 and hasattr(keys[0], 'find_element'):
            driver_obj = keys[0]

            # ----------------------------------------------------------
            # MODE 2 : Type on page (driver, text)
            # ----------------------------------------------------------
            if len(keys) == 2:
                # write(driver, "Hello World")
                text = keys[1]
                action = ActionChains(driver_obj)
                action.send_keys(text).perform()
                return True

            # ----------------------------------------------------------
            # MODE 3 : Type in specific element (driver, selector_type, selector, text)
            # ----------------------------------------------------------
            elif len(keys) == 4:
                # write(driver, "id", "username", "john")
                selector_type = keys[1]
                selector = keys[2]
                text = keys[3]

                element = _get_web_element(driver_obj, selector_type, selector)
                if element:
                    element.send_keys(text)
                    return True
                else:
                    print(f"Element not found: {selector_type} - {selector}")
                    return False

            else:
                print("Error: Invalid number of arguments for Selenium write()")
                print("Use: write(driver, text) or write(driver, selector_type, selector, text)")
                return False

        # ---------------Invalid Arguments---------------
        else:
            print("Error: Invalid arguments for write()")
            print("Use: write(text) or write(driver, text) or write(driver, selector_type, selector, text)")
            return False

    except Exception as e:
        print(f"Error writing text: {e}")
        return False


def year():
    """
    Get current year.

    Returns:
        int: Current year (e.g., 2026)

    Example:
        year()  # 2026
    """
    return time.localtime().tm_year


def zoom(*args):
    """
    Zoom in/out using steps or set zoom percentage.

    Modes:
        1. PyAutoGUI steps/reset: (value)
        2. Selenium steps: (driver, value) where value is -9 to +9
        3. Selenium percentage: (driver, value) where value is outside -9 to +9
        4. Selenium reset: (driver, 100) or (driver, 0)

    Args:
        *args: Variable arguments::

            - (value): PyAutoGUI zoom steps/reset
            - (driver, value): Selenium zoom steps/percentage/reset

    Returns:
        True if successful, False otherwise

    Raises:
        ValueError: If arguments are invalid, zoom value is not an integer,
                    PyAutoGUI value is outside -9 to +9 (except 100),
                    or Selenium percentage is less than 1%.

    Output:
        - Prints zoom direction and step count (PyAutoGUI).
        - Prints new zoom percentage after change (Selenium).
        - Prints confirmation when zoom is reset.

    Value Logic:
        - -9 to +9: Zoom steps (Ctrl+/Ctrl-)
        - 100 or 0: Reset to default/100%
        - Outside range (except 100): Percentage (Selenium only)

    Example:
        ::

            # PyAutoGUI (desktop apps)
            zoom(3)              # Zoom in 3 steps
            zoom(-5)             # Zoom out 5 steps
            zoom(100) or zoom(0) # Reset to default (Ctrl+0)
            When zoom in/out is performed using UI (Ctrl and +/-) in Chrome,
            the min %, zoom states % and max % follow this order:
                (25, 33, 50, 67, 75, 80, 90, 100, 110, 125, 150, 175, 200, 250, 300, 400, 500)

            # Selenium (browser) - Steps
            zoom(driver, 3)      # Zoom in current + 3 * 10%
            zoom(driver, -5)     # Zoom out current - 5 * 10%

            # Selenium (browser) - Reset
            zoom(driver, 100)    # Reset to 100%
            zoom(driver, 0)      # Reset to 100%

            # Selenium (browser) - Percentage
            zoom(driver, 150)    # Set zoom to 150%
            zoom(driver, 75)     # Set zoom to 75%
            zoom(driver, 50)     # Set zoom to 50%
            zoom(driver, 200)    # Set zoom to 200%

    Note:
        - Selenium zoom is applied via JavaScript and is not reflected in
          the Chrome URL bar or the kebab menu zoom indicator.
        - PyAutoGUI reset (0 or 100) uses Ctrl+0 which resets to the
          application's default zoom, which may not always be 100%
          (e.g. a PDF viewer may default to 'fit to page' instead).
        - Selenium reset explicitly sets zoom to exactly 100%.
    """

    # ------------------------------------------------------------
    # DETERMINE MODE
    # ------------------------------------------------------------
    if len(args) == 1:
        # PyAutoGUI mode
        use_driver = False
        driver_obj = None
        value = args[0]
    elif len(args) == 2 and hasattr(args[0], 'execute_script'):
        # Selenium mode - driver object detected
        use_driver = True
        driver_obj = args[0]
        value = args[1]
    else:
        raise ValueError("Invalid arguments. Use zoom(value) or zoom(driver, value)")

    # Validate value
    if not isinstance(value, int):
        raise ValueError("Zoom value must be an integer")

    # PyAutoGUI restrictions
    if not use_driver and value != 100 and abs(value) > 9:
        raise ValueError("PyAutoGUI mode supports steps (-9 to +9) or reset (100) or (0)")

    try:
        # ------------------------------------------------------------
        # SELENIUM MODE
        # needs: nothing : works on all platforms
        # ------------------------------------------------------------
        if use_driver:

            # Special case: 100 or 0 = Reset
            if value == 100 or value == 0:
                driver_obj.execute_script("document.body.style.zoom='100%'")
                # "reset to 100%" here because Selenium explicitly sets zoom to 100%
                # unlike PyAutoGUI's Ctrl+0 which resets to the default zoom level
                # of whatever application is currently active
                print("Zoom reset to 100%")
                return True

            # Steps: -9 to +9
            elif -9 <= value <= 9:
                # Get current zoom level
                current_zoom = driver_obj.execute_script("return document.body.style.zoom || '100%'")
                current_zoom = int(current_zoom.replace('%', '')) if current_zoom else 100

                if value > 0:
                    new_zoom = current_zoom + (value * 10)  # each step = 10%
                    driver_obj.execute_script(f"document.body.style.zoom='{new_zoom}%'")
                    print(f"Zoomed in {value} step(s) - zoom: {new_zoom}%")
                elif value < 0:
                    steps = abs(value)
                    new_zoom = current_zoom - (steps * 10)  # each step = 10%
                    new_zoom = max(10, new_zoom)  # minimum 10% zoom
                    driver_obj.execute_script(f"document.body.style.zoom='{new_zoom}%'")
                    print(f"Zoomed out {steps} step(s) - zoom: {new_zoom}%")
                return True

            # Percentage: Outside -9 to +9 (excluding 100)
            else:
                if value < 1:
                    raise ValueError("Zoom percentage must be at least 1%")

                driver_obj.execute_script(f"document.body.style.zoom='{value}%'")
                print(f"Zoom set to {value}%")
                return True

        # ------------------------------------------------------------
        # PYAUTOGUI MODE (desktop apps)
        # needs: display
        # ------------------------------------------------------------
        else:
            if not _GUI_AVAILABLE:
                print("Error: zoom() requires a display.")
                return False

            # Special case: 100 or 0 = Reset
            if value == 100 or value == 0:
                pyautogui.hotkey('ctrl', '0')
                # "reset to default" not "reset to 100%" because Ctrl+0 resets to
                # the application's default zoom, which may not always be 100%
                # e.g. a PDF viewer may default to "fit to page" instead of 100%
                print("Zoom reset to default.")
                return True

            # Steps: -9 to +9
            elif value > 0:
                # Zoom in
                for _ in range(value):
                    pyautogui.hotkey('ctrl', '+')
                    time.sleep(0.3)
                print(f"Zoomed in {value} step(s)")
            elif value < 0:
                # Zoom out
                steps = abs(value)
                for _ in range(steps):
                    pyautogui.hotkey('ctrl', '-')
                    time.sleep(0.3)
                print(f"Zoomed out {steps} step(s)")

            return True

    except Exception as e:
        print(f"Error zooming: {e}")
        return False
